import os
os.environ["GRADIO_MCP_SERVER"] = "false"
import base64
import re
import tempfile
from typing import List, Tuple, Annotated, TypedDict, Optional, Literal
import operator
import gradio as gr
import markdown as md
import numpy as np
import pandas as pd

from openai import OpenAI
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage
from pydantic import BaseModel, Field
from langgraph.graph import END, StateGraph

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.colors import LinearSegmentedColormap


def _register_cjk_font():
    """Register a CJK-capable font so matplotlib can render Traditional Chinese
    (the base Docker image has no CJK font installed unless one is bundled)."""
    try:
        # Priority matters: some CJK "fallback" fonts (e.g. Droid Sans Fallback) map every
        # codepoint to a glyph -- including a blank/tofu box for Latin digits/punctuation --
        # so matplotlib never triggers real fallback to a Latin font for them. Noto Sans CJK
        # and WenQuanYi Zen Hei render both CJK and Latin correctly, so they must come first.
        priority = ("NotoSansCJK", "NotoSerifCJK", "wqy-zenhei", "DroidSansFallback")
        candidates = []
        search_roots = ("/usr/share/fonts", "/usr/local/share/fonts", os.path.expanduser("~/.fonts"))
        for root in search_roots:
            if not os.path.isdir(root):
                continue
            for dirpath, _, files in os.walk(root):
                for fn in files:
                    for rank, key in enumerate(priority):
                        if key in fn:
                            candidates.append((rank, os.path.join(dirpath, fn)))
                            break
        candidates.sort(key=lambda t: t[0])
        family_names = []
        for _, path in candidates:
            try:
                fm.fontManager.addfont(path)
                family_names.append(fm.FontProperties(fname=path).get_name())
            except Exception:
                continue
        families = family_names + [
            "Noto Sans CJK TC", "Noto Sans CJK SC", "WenQuanYi Zen Hei", "DejaVu Sans",
        ]
        plt.rcParams["font.sans-serif"] = families
        plt.rcParams["font.family"] = "sans-serif"
        plt.rcParams["axes.unicode_minus"] = False
    except Exception:
        pass


_register_cjk_font()

# ── Clients ────────────────────────────────────────────────────────────────────
_api_key = os.environ.get("OPENAI_API_KEY")
openai_client = OpenAI(api_key=_api_key)
llm_4o = ChatOpenAI(model="gpt-4o", temperature=0.7, api_key=_api_key)
llm_mini = ChatOpenAI(model="gpt-4o-mini", temperature=0, api_key=_api_key)

# ── File extraction helpers ────────────────────────────────────────────────────
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".svg"}
MIME_MAP = {"jpg": "jpeg", "jpeg": "jpeg", "png": "png", "gif": "gif",
            "webp": "webp", "bmp": "bmp", "tiff": "tiff", "svg": "svg+xml"}

def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)

def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        slides.append(f"[投影片 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)

def _extract_excel(path: str) -> str:
    import pandas as pd
    xl = pd.ExcelFile(path)
    parts = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        parts.append(f"[工作表: {sheet}]\n{df.to_string(index=False)}")
    return "\n\n".join(parts)

def _extract_csv(path: str) -> str:
    import pandas as pd
    df = pd.read_csv(path)
    return df.to_string(index=False)

def _extract_image(path: str) -> str:
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    mime = f"image/{MIME_MAP.get(ext, 'png')}"
    with open(path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    resp = openai_client.chat.completions.create(
        model="gpt-4o",
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": "請詳細描述這張圖表或圖片的所有內容，包含數據、標題、標籤、趨勢與重要訊息，以便用於課程教材設計。"},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
            ],
        }],
        max_tokens=1500,
    )
    return f"[圖表內容分析]\n{resp.choices[0].message.content}"

def extract_file_text(file_obj) -> str:
    if file_obj is None:
        return ""
    path = file_obj.name if hasattr(file_obj, "name") else str(file_obj)
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext == ".docx":
            return _extract_docx(path)
        elif ext in {".pptx", ".ppt"}:
            return _extract_pptx(path)
        elif ext in {".xlsx", ".xls"}:
            return _extract_excel(path)
        elif ext == ".csv":
            return _extract_csv(path)
        elif ext in IMAGE_EXTS:
            return _extract_image(path)
        else:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as e:
        return f"[文件讀取失敗：{e}]"

def extract_files_text(file_list) -> str:
    if not file_list:
        return ""
    if not isinstance(file_list, list):
        file_list = [file_list]
    parts = []
    for f in file_list:
        name = os.path.basename(f.name if hasattr(f, "name") else str(f))
        text = extract_file_text(f)
        parts.append(f"=== {name} ===\n{text}")
    return "\n\n".join(parts)

# ── HTML export helper ─────────────────────────────────────────────────────────
HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="zh-TW">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>AI課程教材轉譯結果</title>
  <style>
    body {{ font-family: "Noto Sans TC", "Microsoft JhengHei", sans-serif;
            max-width: 960px; margin: 40px auto; padding: 0 24px;
            color: #1a1a1a; line-height: 1.8; background: #fafafa; }}
    h1, h2, h3 {{ color: #2563eb; margin-top: 2em; }}
    h2 {{ border-bottom: 2px solid #e5e7eb; padding-bottom: 6px; }}
    table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
    th {{ background: #2563eb; color: white; padding: 10px 14px; text-align: left; }}
    td {{ padding: 8px 14px; border: 1px solid #d1d5db; }}
    tr:nth-child(even) {{ background: #f3f4f6; }}
    code {{ background: #f1f5f9; padding: 2px 6px; border-radius: 4px; font-size: 0.9em; }}
    pre {{ background: #1e293b; color: #e2e8f0; padding: 16px; border-radius: 8px; overflow-x: auto; }}
    hr {{ border: none; border-top: 1px solid #e5e7eb; margin: 2em 0; }}
    blockquote {{ border-left: 4px solid #2563eb; margin: 0; padding: 8px 16px;
                  background: #eff6ff; color: #1e40af; }}
  </style>
</head>
<body>
{body}
</body>
</html>"""

def generate_html_file(markdown_content: str):
    if not markdown_content or markdown_content.startswith("*送出資料後"):
        return None
    body = md.markdown(markdown_content, extensions=["tables", "fenced_code", "nl2br"])
    html = HTML_TEMPLATE.format(body=body)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html",
                                      mode="w", encoding="utf-8")
    tmp.write(html)
    tmp.close()
    return tmp.name

def generate_word_file(markdown_content: str):
    if not markdown_content or markdown_content.startswith("*送出資料後"):
        return None
    import re
    from docx import Document
    from docx.shared import RGBColor

    doc = Document()
    lines = markdown_content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("|"):
            # collect contiguous table lines
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            data_rows = [r for r in table_lines
                         if not re.match(r"^\|[-| :]+\|$", r.strip())]
            if data_rows:
                parsed = [[c.strip() for c in r.strip("|").split("|")]
                          for r in data_rows]
                ncols = max(len(r) for r in parsed)
                tbl = doc.add_table(rows=len(parsed), cols=ncols)
                tbl.style = "Table Grid"
                for ri, row in enumerate(parsed):
                    for ci, cell in enumerate(row):
                        if ci < ncols:
                            tbl.cell(ri, ci).text = cell
            continue
        elif line.startswith("---") or line == "":
            i += 1
            continue
        else:
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
            text = re.sub(r"\*(.+?)\*", r"\1", text)
            if text.strip():
                doc.add_paragraph(text)
        i += 1

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    tmp.close()
    doc.save(tmp.name)
    return tmp.name

# ══════════════════════════════════════════════════════════════════════════════
# TAB 1 — AI課程教材轉譯與受眾案例設計Agent (Plan-Execute架構)
# ══════════════════════════════════════════════════════════════════════════════

COURSE_AGENT_SYSTEM = """
你是我的專屬課程教材設計助理，名稱為「AI課程教材轉譯與受眾案例設計Agent」。
你熟悉我的教學風格：先有趣，再理論，再實作；重視受眾痛點、實務案例、工作流程、AI工具操作、簡報呈現與課堂活動。
你的任務不是整理資料，而是根據我提供的資料，轉化成適合「企業與跨職能AI初學者」的課程案例、教材內容、實作活動、提示詞與簡報頁面。

## 確認受眾輪廓
這批學員多數已經聽過或使用過AI，但自評能力仍以新手與初階為主。
他們不是完全沒有動機，而是缺乏清楚的方法、範本、工作流程與實際操作經驗。
課程不能只介紹AI工具，也不能一開始就講太深的技術原理。
最有效的教學方式是從他們每天遇到的工作痛點切入，讓他們看到AI如何協助完成資料分析、文件整理、簡報製作、流程自動化與行銷內容產出。

**學員關鍵特徵：**
- 多數已有AI接觸經驗，但缺乏系統化方法
- 多數自評為新手或初階
- 幾乎全體對建立AI Agent有興趣
- 主要需求集中在資料分析、文書整理、流程自動化、提案簡報、會議紀錄、行銷文案
- 產業以製造業為最大宗，但也包含服務業、教育、資訊科技與其他產業
- 職能橫跨數據分析、業務、行銷、社群、主管、行政與教育訓練
- 常用文件包含PPT、Excel、Word、企劃書、會議紀錄、社群貼文與報價單

**教學涵義：**
這批學員最需要的不是AI概念介紹，而是「把AI放進工作流程」。
課程應該用簡報、Excel、報告、會議紀錄、業務資料、行銷文案等共同文件作為主要案例。
每個單元都要讓學員帶走一個可複製的模板、提示詞或小型SOP。

## 我的教學風格
- 先用痛點或有趣案例開場
- 再拆解背後概念
- 接著示範AI工具如何處理
- 最後讓學員自己操作
- 不只教工具功能，要教「什麼情境該怎麼用」
- 內容要好講、好理解、好練習
- 案例要貼近企業現場與跨職能工作
- 簡報內容不能只有圖多字少，也不能像密密麻麻的講義
- 每個概念都要能轉成案例、活動、提示詞或工作流程

## 課程設計定位
這門課要讓學員從「零散使用AI」進階到「用AI改造一個工作流程」。
教學重點不是炫技，而是讓學員知道：
- 我現在最耗時的工作，可以怎麼拆解？
- 哪些步驟適合交給AI？
- 怎麼下提示詞？
- 怎麼檢查AI產出？
- 怎麼把結果整理成簡報、報告或可交付成果？

## 五大優先主題
1. **AI資料整理與報告自動化** — 最大公約數需求，適合結合Excel、資料摘要、週報、主管報告與簡報初稿
   - 案例：Excel銷售資料整理成分析報告、客戶名單轉成分群建議、生產或業務數據轉成主管簡報、問卷資料轉成洞察與課程建議
2. **AI簡報與提案製作** — PPT是最高頻文件之一，適合所有職能，最容易讓學員快速有感
   - 案例：把企劃書轉成10頁簡報、把會議資料整理成主管簡報、把產品資料轉成業務提案、把課程內容轉成招生簡報
3. **個人AI Agent與工作流程設計** — 幾乎所有學員都想建立AI Agent，需用低門檻方式教
   - 案例：Email重點整理助理、客戶追蹤提醒助理、會議紀錄整理助理、競品資料蒐集與摘要助理、報價合約重點整理助理
4. **AI行銷文案與社群內容工廠** — 適合行銷、業務、創業與副業族群，可作為高需求應用模組
   - 案例：社群貼文批次生成、廣告文案改寫、產品賣點萃取、競品文案分析、SEO與AEO內容優化
5. **AI會議紀錄與文書整理** — 入門門檻低，立即見效，適合作為課程暖身或第一個實作
   - 案例：會議逐字稿轉成決議與待辦、客戶訪談紀錄轉成需求表、工作日誌轉成週報、合約或報價單轉成重點摘要

## 當使用者提供資料時
不要直接摘要。請先分析：
- 這份資料適合用來教哪個概念？
- 對這批學員而言，哪個工作場景最有感？
- 可以變成什麼案例？
- 可以安排什麼實作？
- 可以產出什麼簡報頁？

必須轉化成以下所有項目：受眾聽得懂的白話說法、企業現場案例、教師上課講解稿、學員實作任務、可直接使用的提示詞、簡報頁面內容、課堂提問、作業或延伸任務、檢查AI產出的標準。

## 案例設計規則
每個案例必須包含：案例名稱、適合職能、對應痛點、原始資料類型、AI協助方式、操作步驟、可使用提示詞、預期產出、教師講解重點、學員容易犯的錯、延伸應用。

## 輸出前品質檢查
產出前請確認：
- 是否符合企業與跨職能初學者
- 是否對應資料分析、簡報、文書、流程自動化或行銷需求
- 是否能讓學員當場操作
- 是否有可帶走的模板或SOP
- 是否避免只介紹AI工具
- 是否避免過度技術化
- 是否符合先有趣、再理論、再實作
- 是否適合拿來做簡報與上課

## 限制
- 不要只摘要資料
- 不要寫成泛用AI課程
- 不要把學員當成完全不懂AI的人
- 不要把課程設計得太技術導向
- 不要使用離企業工作太遠的案例
- 不要只給概念，要給案例、提示詞、活動與產出
- 不要一直反問，資訊不足時先做合理假設並標示

## 預設受眾（如未指定）
已接觸AI但仍屬新手到初階的企業學員，工作內容橫跨資料分析、業務、行銷、主管、行政與教育訓練。
產業以製造業為主要案例核心，再補充服務業、教育、資訊科技與創業副業情境。
"""

SECTION_FORMATS = {
    "一、這份資料適合怎麼教": """\
請用以下表格格式輸出：
| 判斷項目 | 內容 |
|---|---|
| 適合的課程單元 |  |
| 對應受眾痛點 |  |
| 最適合的職能族群 |  |
| 可轉成的案例 |  |
| 不適合直接教的地方 |  |""",

    "二、教材轉譯表": """\
請用以下表格格式輸出：
| 原始概念 | 學員聽得懂的說法 | 對應工作場景 | 教學重點 |
|---|---|---|---|""",

    "三、企業現場案例設計": """\
請用以下表格格式輸出，每個案例包含所有必填欄位：
| 案例名稱 | 適合職能 | 情境 | AI任務 | 預期產出 |
|---|---|---|---|---|
之後再逐一展開每個案例的完整設計（適合職能、對應痛點、原始資料類型、AI協助方式、操作步驟、可使用提示詞、預期產出、教師講解重點、學員容易犯的錯、延伸應用）。""",

    "四、課堂實作活動": """\
請用以下表格格式輸出：
| 階段 | 時間 | 教師做什麼 | 學員做什麼 | 產出 |
|---|---:|---|---|---|""",

    "五、學員可直接使用的提示詞": """\
請提供完整提示詞，包含角色、背景、任務、限制、輸出格式與檢查標準。每個提示詞單獨呈現，清楚標示使用場景。""",

    "六、教師講解稿": """\
請提供一段自然、可直接上課使用的講解內容，語氣口語、好講、有節奏，包含開場、概念說明、案例示範與小結。""",

    "七、簡報頁面建議": """\
請用以下表格格式輸出：
| 頁面 | 標題 | 核心訊息 | 畫面建議 | 講解重點 |
|---|---|---|---|---|""",

    "八、作業或延伸任務": """\
請用以下表格格式輸出：
| 類型 | 任務 | 評量重點 |
|---|---|---|""",
}

SECTIONS = [
    "一、這份資料適合怎麼教",
    "二、教材轉譯表",
    "三、企業現場案例設計",
    "四、課堂實作活動",
    "五、學員可直接使用的提示詞",
    "六、教師講解稿",
    "七、簡報頁面建議",
    "八、作業或延伸任務",
]

TOPIC_CHOICES = [
    "自動判斷",
    "AI資料整理與報告自動化",
    "AI簡報與提案製作",
    "個人AI Agent與工作流程設計",
    "AI行銷文案與社群內容工廠",
    "AI會議紀錄與文書整理",
]

# ── Course Design: State & Models ──────────────────────────────────────────────
class CourseState(TypedDict):
    input: str
    topic_focus: str
    plan: List[str]
    past_steps: Annotated[List[Tuple[str, str]], operator.add]
    response: str

class CoursePlan(BaseModel):
    """教材轉譯章節計畫"""
    steps: List[str] = Field(description="依序需要產出的章節名稱列表")

# ── Course Design: Nodes ───────────────────────────────────────────────────────
def course_planner(state: CourseState):
    topic_hint = f"\n聚焦主題：{state['topic_focus']}" if state.get("topic_focus") else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", COURSE_AGENT_SYSTEM),
        ("human",
         "請分析以下教材內容，決定最適合產出的章節順序。\n"
         "可選章節：" + "／".join(SECTIONS) + "\n"
         "請列出所有八個章節，依最佳教學邏輯排序。{topic_hint}\n\n"
         "教材內容：\n{input}")
    ])
    planner = prompt | llm_4o.with_structured_output(CoursePlan)
    result = planner.invoke({"input": state["input"], "topic_hint": topic_hint})
    return {"plan": result.steps}

def course_executor(state: CourseState):
    section = state["plan"][0]
    completed = "\n\n".join(
        f"### {s}\n{r}" for s, r in state["past_steps"]
    ) if state["past_steps"] else "（尚無已完成章節）"
    section_key = next((k for k in SECTION_FORMATS if k in section), None)
    format_hint = f"\n\n**格式要求：**\n{SECTION_FORMATS[section_key]}" if section_key else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", COURSE_AGENT_SYSTEM),
        ("human",
         "原始教材：\n{input}\n\n"
         "已完成章節（供參考，保持一致性）：\n{completed}\n\n"
         "請現在產出「{section}」的完整內容。內容要具體、可直接使用。{format_hint}")
    ])
    chain = prompt | llm_4o
    result = chain.invoke({
        "input": state["input"],
        "completed": completed,
        "section": section,
        "format_hint": format_hint,
    })
    return {"past_steps": [(section, result.content)]}

def course_replanner(state: CourseState):
    remaining = state["plan"][1:]
    if not remaining:
        full_output = "\n\n---\n\n".join(
            f"## {section}\n\n{content}"
            for section, content in state["past_steps"]
        )
        return {"response": full_output}
    return {"plan": remaining}

def course_should_end(state: CourseState):
    return "end" if state.get("response") else "continue"

def build_course_graph():
    wf = StateGraph(CourseState)
    wf.add_node("planner", course_planner)
    wf.add_node("executor", course_executor)
    wf.add_node("re_planner", course_replanner)
    wf.set_entry_point("planner")
    wf.add_edge("planner", "executor")
    wf.add_edge("executor", "re_planner")
    wf.add_conditional_edges("re_planner", course_should_end, {"continue": "executor", "end": END})
    return wf.compile()

course_graph = build_course_graph()

def run_course_design(materials: str, topic_focus: str, user_prompt: str = ""):
    if user_prompt and user_prompt.strip():
        materials = materials + f"\n\n---\n使用者額外指示：\n{user_prompt.strip()}"
    if not materials.strip():
        yield "⚠️ 請輸入教材內容後再送出。"
        return

    inputs = {
        "input": materials,
        "topic_focus": topic_focus if topic_focus != "自動判斷" else "",
    }

    output = ""
    plan_shown = False

    for event in course_graph.stream(inputs, {"recursion_limit": 30}):
        for node_name, value in event.items():
            if node_name == "planner" and "plan" in value and not plan_shown:
                plan_shown = True
                output = "### 📋 章節規劃\n" + "\n".join(
                    f"{i+1}. {s}" for i, s in enumerate(value["plan"])
                ) + "\n\n---\n\n*逐章節生成中……*\n\n"
                yield output
            elif node_name == "executor" and "past_steps" in value:
                for section, content in value["past_steps"]:
                    # Replace the "generating" notice and append completed section
                    output = output.replace("*逐章節生成中……*\n\n", "")
                    output += f"## {section}\n\n{content}\n\n---\n\n*逐章節生成中……*\n\n"
                    yield output
            elif node_name == "re_planner" and value.get("response"):
                # All sections done — show clean final output
                yield value["response"]


# ══════════════════════════════════════════════════════════════════════════════
# TAB 2 — 客戶提案總監Agent (Plan-Execute架構)
# ══════════════════════════════════════════════════════════════════════════════

PROPOSAL_AGENT_SYSTEM = """
你是一位 B2B 製造業客戶提案總監Agent，名稱為「客戶提案總監Agent」。
你擅長整合商品型錄圖片、Excel 報價表、客戶需求、業務筆記與初步大綱，
產出圖文並茂、有深度、有邏輯、可交付給客戶的正式提案計畫書。

## 主要目標
協助使用者把零散資料整理成正式商務提案，
讓客戶快速理解問題、看見方案價值、降低採購疑慮，
並願意安排下一步會議、PoC 試點或下訂。

## 目標受眾
- 製造業客戶主管
- 採購決策者
- 廠務與品保主管
- 經營者
- 業務與提案團隊

## 輸入資料處理方式
- **商品型錄圖片**：整理產品名稱、功能重點、規格資訊與建議放置章節
- **Excel 報價表**：整理報價項目、金額摘要、付款條件，轉成客戶看得懂的投資說明
- **客戶需求**：分析客戶現況、痛點、決策因素與可能疑慮
- **業務筆記**：將零散想法轉成提案主軸、銷售論點與客戶利益
- **提案大綱**：補強原始大綱，建立正式商務提案架構

## 思考規則
- 先理解客戶痛點，再介紹產品
- 每一個產品功能都要對應一個客戶問題
- 每一個報價項目都要說明它帶來的價值
- 提案不能只是產品型錄重排，也不能只是報價單整理
- 必須有現況診斷、解決方案、導入流程、效益、風險配套與下一步行動
- 不得編造不存在的數字、成果、保證或客戶案例
- 資料不足時，請標示「需補充資料」，並列出建議補充問題
- 語氣要正式、清楚、具商務說服力，適合提供給企業客戶閱讀

## 銷售邏輯
- 先說明客戶目前可能承擔的風險或損失
- 再說明方案如何降低風險、提升效率或改善管理
- 將技術規格轉成客戶能理解的營運價值
- 將報價金額轉成投資項目與合作價值
- 建議採用「先 PoC 試點，再多線擴展」降低客戶疑慮
- 結尾必須給出明確下一步行動

## 寫作風格
語氣：正式、專業、清楚、有商務說服力
避免：空泛口號、過度誇張、制式模板語氣、只列規格不說價值、只列價格不說投資理由
偏好：客戶痛點導向、問題診斷導向、解決方案導向、低風險試點導向、成效驗證導向、下一步行動導向

## 輸出前品質檢查
- 是否清楚說明客戶痛點？
- 是否把產品規格轉成客戶價值？
- 是否把報價表轉成投資說明？
- 是否有圖片配置建議？
- 是否有必要表格？
- 是否有導入流程？
- 是否有風險配套？
- 是否有下一步行動？
- 是否避免誇大與亂編？
- 是否像一份可交付給客戶的正式商務提案？

若資料不足，請先列出不足處，但仍要依現有資料產出可用初稿。
"""

PROPOSAL_SECTIONS = [
    "提案封面資訊",
    "執行摘要",
    "客戶現況與問題診斷",
    "本次提案核心主張",
    "解決方案總覽",
    "商品型錄圖片整合說明",
    "產品與規格摘要",
    "客戶痛點與方案對應表",
    "導入流程與時程規劃",
    "報價與投資說明",
    "預期效益與驗證方式",
    "風險與配套措施",
    "為什麼選擇我們",
    "建議下一步",
    "附錄",
]

PROPOSAL_SECTION_FORMATS = {
    "提案封面資訊": "請列出：提案名稱、提案對象（公司名稱與聯絡人）、提案日期、提案單位、版本號。",
    "執行摘要": "用200字以內濃縮整份提案的核心價值主張、方案亮點與建議下一步行動。",
    "客戶現況與問題診斷": "說明客戶目前面臨的問題、風險與損失，以及不解決問題的代價。",
    "本次提案核心主張": "用一段話說明本次提案的核心主張與差異化價值。",
    "解決方案總覽": "概述提案的解決方案架構與模組，搭配圖表說明整體方案邏輯。",
    "商品型錄圖片整合說明": "針對每張圖片提供：建議放置章節、圖片標題、圖片說明文字、在提案中的說服作用。若無圖片資料，請標示需補充並說明建議拍攝或準備的圖片類型。",
    "產品與規格摘要": """\
請用以下表格格式輸出：
| 產品名稱 | 型號／規格 | 核心功能 | 適用場景 | 備註 |
|---|---|---|---|---|""",
    "客戶痛點與方案對應表": """\
請用以下表格格式輸出：
| 客戶痛點 | 目前影響 | 本方案如何解決 | 預期改善 |
|---|---|---|---|""",
    "導入流程與時程規劃": """\
請用以下表格格式輸出：
| 階段 | 工作項目 | 負責方 | 時程 | 交付物 |
|---|---|---|---|---|""",
    "報價與投資說明": """\
請用以下表格格式輸出：
| 項目 | 說明 | 數量 | 單價 | 小計 | 客戶獲得的價值 |
|---|---|---|---|---|---|
最後加上付款條件與合作期限說明。若無報價資料，請標示需補充。""",
    "預期效益與驗證方式": """\
請用以下表格格式輸出：
| 效益項目 | 目前狀況 | 導入後預期 | 驗證方式 | 衡量指標 |
|---|---|---|---|---|""",
    "風險與配套措施": """\
請用以下表格格式輸出：
| 潛在風險 | 風險等級 | 配套措施 | 責任方 |
|---|---|---|---|""",
    "為什麼選擇我們": "說明核心差異化優勢、過往經驗、服務保障與合作誠意，語氣正式有說服力。",
    "建議下一步": """\
請用以下表格格式輸出：
| 步驟 | 行動項目 | 建議時間 | 負責方 |
|---|---|---|---|
最後加上一段促成下一步會議或 PoC 試點的邀請語。""",
    "附錄": "放置補充資料、參考規格、詞彙說明，或列出「需補充資料」清單與建議補充問題。",
}

class ProposalState(TypedDict):
    input: str
    plan: List[str]
    past_steps: Annotated[List[Tuple[str, str]], operator.add]
    response: str

class ProposalPlan(BaseModel):
    """提案章節計畫"""
    steps: List[str] = Field(description="依序需要產出的提案章節名稱列表")

def proposal_planner(state: ProposalState):
    prompt = ChatPromptTemplate.from_messages([
        ("system", PROPOSAL_AGENT_SYSTEM),
        ("human",
         "請分析以下資料，決定最適合的提案章節順序。\n"
         "可選章節：" + "／".join(PROPOSAL_SECTIONS) + "\n"
         "請列出所有章節，依最佳商務提案邏輯排序。\n\n"
         "資料內容：\n{input}")
    ])
    planner = prompt | llm_4o.with_structured_output(ProposalPlan)
    result = planner.invoke({"input": state["input"]})
    return {"plan": result.steps}

def proposal_executor(state: ProposalState):
    section = state["plan"][0]
    completed = "\n\n".join(
        f"### {s}\n{r}" for s, r in state["past_steps"]
    ) if state["past_steps"] else "（尚無已完成章節）"
    section_key = next((k for k in PROPOSAL_SECTION_FORMATS if k in section), None)
    format_hint = f"\n\n**格式要求：**\n{PROPOSAL_SECTION_FORMATS[section_key]}" if section_key else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", PROPOSAL_AGENT_SYSTEM),
        ("human",
         "原始資料：\n{input}\n\n"
         "已完成章節（供參考，保持一致性）：\n{completed}\n\n"
         "請現在產出「{section}」的完整內容。內容要具體、正式、可直接交付給客戶。{format_hint}")
    ])
    chain = prompt | llm_4o
    result = chain.invoke({
        "input": state["input"],
        "completed": completed,
        "section": section,
        "format_hint": format_hint,
    })
    return {"past_steps": [(section, result.content)]}

def proposal_replanner(state: ProposalState):
    remaining = state["plan"][1:]
    if not remaining:
        full_output = "\n\n---\n\n".join(
            f"## {section}\n\n{content}"
            for section, content in state["past_steps"]
        )
        return {"response": full_output}
    return {"plan": remaining}

def proposal_should_end(state: ProposalState):
    return "end" if state.get("response") else "continue"

def build_proposal_graph():
    wf = StateGraph(ProposalState)
    wf.add_node("planner", proposal_planner)
    wf.add_node("executor", proposal_executor)
    wf.add_node("re_planner", proposal_replanner)
    wf.set_entry_point("planner")
    wf.add_edge("planner", "executor")
    wf.add_edge("executor", "re_planner")
    wf.add_conditional_edges("re_planner", proposal_should_end, {"continue": "executor", "end": END})
    return wf.compile()

proposal_graph = build_proposal_graph()

def run_proposal_design(materials: str, user_prompt: str = ""):
    if user_prompt and user_prompt.strip():
        materials = materials + f"\n\n---\n使用者額外指示：\n{user_prompt.strip()}"
    if not materials.strip():
        yield "⚠️ 請輸入資料後再送出。"
        return

    output = ""
    plan_shown = False

    for event in proposal_graph.stream({"input": materials}, {"recursion_limit": 40}):
        for node_name, value in event.items():
            if node_name == "planner" and "plan" in value and not plan_shown:
                plan_shown = True
                output = "### 📋 提案章節規劃\n" + "\n".join(
                    f"{i+1}. {s}" for i, s in enumerate(value["plan"])
                ) + "\n\n---\n\n*逐章節生成中……*\n\n"
                yield output
            elif node_name == "executor" and "past_steps" in value:
                for section, content in value["past_steps"]:
                    output = output.replace("*逐章節生成中……*\n\n", "")
                    output += f"## {section}\n\n{content}\n\n---\n\n*逐章節生成中……*\n\n"
                    yield output
            elif node_name == "re_planner" and value.get("response"):
                yield value["response"]


# ══════════════════════════════════════════════════════════════════════════════
# TAB 3 — Image to Commerce Decision Growth Agent (Plan-Execute架構)
# ══════════════════════════════════════════════════════════════════════════════

COMMERCE_AGENT_SYSTEM = """
你是一位品牌策略、商品定位、電商轉換、社群內容企劃、競品分析與高擴散內容設計專家，名稱為「Image to Commerce Decision Growth Agent」。
你擅長根據商品圖片，判斷商品類型、設計語言、情緒價值、使用情境、潛在客群、購買猶豫點、競品差異與市場切入位置。
你的工作是把商品圖片轉成一套可用於社群、電商頁、短影音與品牌企劃的內容策略。
每一次輸出都要讓品牌更容易被看見、被記住、被比較後選中。

## 核心原則
好文案的核心，是幫消費者更快做出選擇。每一段內容都要回答三個問題：
- 為什麼這個商品值得停下來看？
- 為什麼它值得被擁有？
- 為什麼它比其他類似選項更有記憶點？

請把商品特徵轉成具體情境、情緒價值與購買理由。
若文案偏向形容詞堆疊，請改成更具畫面感的生活場景。
若文案偏向一般商品介紹，請補上差異化觀點與消費者選擇理由。
若不同平台的內容太接近，請依平台任務重新調整語氣與結構。

## 任務清單
- 根據商品圖片判斷商品類型、設計風格、受眾與使用情境
- 將外觀特徵轉成情緒價值、購買理由、送禮理由與收藏理由
- 找出商品相對於替代品的差異化位置
- 預判消費者購買前的猶豫，並用文案自然處理
- 產出高吸引力標題、社群文案、商品頁文案與短影音腳本
- 設計具備停留、收藏、留言、分享潛力的高擴散內容
- 在沒有外部搜尋能力時，根據使用者提供的資料或圖片做保守判斷
- 維持品牌質感，讓內容有銷售力，也保有自然與可信度
- 在資訊不足時，清楚標示可見事實、合理推測與需要補充的資訊

## 圖片資訊分類原則
先整理圖片能支持的資訊，分成三類：
- **可見事實**：商品外觀、造型、色彩、比例、表面質感、包裝、情境元素、圖片中明確可見的文字或符號
- **合理推測**：可合理推測的商品風格、使用情境、情緒價值、目標客群、送禮或收藏可能性（請標示「從圖片推測」）
- **需要補充確認**：材質、尺寸、產地、耐用度、安全認證、品牌故事、銷售成績、使用者評價、實際功效

## 視覺分析原則
根據圖片進行視覺判讀，把外觀轉成消費者能感受到的語言。
重點放在商品如何進入生活，而非單純描述它長什麼樣子。
每一個視覺判斷都要連回消費者感受，推測內容請自然標示「從圖片推測」。
資訊不足處請留給使用者補充，以具體情境取代籠統形容。

## 消費者需求判斷
- 消費者為什麼會被這類商品吸引？
- 它滿足的是實用需求、情緒需求、送禮需求、收藏需求，還是空間風格需求？
- 消費者可能會拿它和哪些商品比較？
- 這個商品最適合佔據哪一個心智位置？
- 消費者滑到這張商品圖時，哪個理由會讓他停下來？

## 競品比較原則
目前無外部搜尋能力，請標示「未進行外部查找」。
可根據商品類型推測常見替代品，並標示為推測。
若要做精準比較，請提醒使用者提供競品連結、截圖或品牌名稱。
比較重點放在消費者選擇理由，以客觀比較取代攻擊語氣。

## 高擴散內容原則
高擴散內容要讓人有理由收藏、分享、標記朋友、留言或轉發，同時保有品牌質感。
以具體情境取代空泛情緒，以身份認同、生活場景、送禮時刻或實用收藏帶動擴散。

擴散觸發點類型：
- **身份認同觸發**：讓消費者覺得這很像自己、朋友、伴侶或某種生活狀態
- **情緒觸發**：放大療癒、陪伴、驚喜、可愛、幽默、儀式感、懷舊感或安全感
- **社交觸發**：讓人想標記朋友、分享給特定對象或引發互動
- **好奇觸發**：用反差、細節、造型、使用情境或意外感讓人停下來看
- **實用觸發**：讓內容具有收藏價值，例如送禮清單、佈置靈感、挑選指南

## 各平台文案任務
- **Instagram**：停留、收藏、分享；有情緒、有生活感；短句有節奏；含 hook、情境描寫、CTA、5-8個 hashtag
- **Threads**：引發共鳴與討論；像朋友聊天，真實有觀點；含 hook、延伸感受、可互動的問題或觀點
- **Facebook**：建立信任與互動；敘事完整帶溫度；含 hook、情境故事、商品價值、問句互動、CTA
- **小紅書**：分享推薦、收藏、搜尋；真實分享有體驗感；含吸睛標題、心得式描述、推薦理由、收藏導向語句
- **Pinkoi 商品頁**：設計價值與送禮理由；設計品牌感細膩；含商品亮點、使用情境、設計理念、送禮說明
- **電商商品頁**：降低購買決策成本；明確好懂購買導向；含賣點、情境、購買理由、行動引導
- **短影音腳本**：前三秒留住注意力；快節奏立刻抓住；含三秒開場畫面+台詞、商品亮點帶出、使用情境、CTA

## 語言風格規則
- 把「必買」改成具體擁有理由
- 把「很夯」改成可觀察的生活情境
- 把「超值」改成具體使用或送禮價值
- 把「很可愛」改成具體表情、細節或陪伴感
- 把「很有質感」改成材質感、色彩、線條或空間氛圍
- 文字自然、有生活感、有品牌感，也有銷售力
- 有行動引導，但不急迫

## 標題設計原則
產出十個高吸引力標題，涵蓋：好奇型、情境型、情緒型、問題型、利益型、送禮型、收藏型、設計感型、反差型、社交互動型。
標題要有生活感與記憶點，寫出想擁有的理由，貼合商品，不只追求聲量。
最後選出最推薦的標題並說明選擇理由與最適合的平台。

## 輸出品質標準
- 消費者能在三秒內知道商品值得看
- 消費者能理解商品和其他選項不同在哪裡
- 文案能把商品外觀轉成情緒、場景與購買理由
- 平台文案有明顯任務差異
- 高擴散內容有明確分享、收藏或互動理由
- 競品比較能反推我方切入位置
- 最終內容可直接交給品牌、社群或電商團隊使用
"""

COMMERCE_SECTIONS = [
    "商品圖片洞察",
    "消費者會想買的理由",
    "差異化定位",
    "競品與替代品比較",
    "高擴散內容角度",
    "十個標題候選與最佳推薦",
    "各平台文案",
    "關鍵字與 hashtag",
    "品質檢查與可補充資訊",
]

COMMERCE_SECTION_FORMATS = {
    "商品圖片洞察": """\
請分三部分輸出：
**1. 可見事實**：列出圖片中明確可見的商品外觀、造型、色彩、比例、表面質感、包裝與情境元素。
**2. 合理推測**：列出可合理推測的商品風格、使用情境、情緒價值、目標客群、送禮或收藏可能性（請標示「從圖片推測」）。
**3. 需要補充確認**：列出材質、尺寸、產地等需要使用者提供才能確認的資訊。""",

    "消費者會想買的理由": """\
請輸出：
- **情緒價值**：這個商品帶給消費者的情感滿足
- **使用情境**：最能觸發購買的具體生活場景（2-3個）
- **送禮或收藏理由**：為什麼有人會把它當禮物或收藏品
- **最強購買觸發點**：一句話說出讓人下單的核心理由
- **購買猶豫與處理**：預判主要猶豫點（2-3個）與對應的文案策略""",

    "差異化定位": """\
請輸出：
- **一句話市場定位**（說清楚這個商品在市場上的位置）
- **核心賣點**（1-2個，要比「好看」更具體）
- **和一般同類商品最大的不同之處**
- **文案應該放大的主軸**
- **文案應該避開的普通化角度**""",

    "競品與替代品比較": """\
請標示「本次未進行外部搜尋」，並：
- 推測消費者可能拿這個商品比較的替代品類型（標示為推測）
- 分析競品常見說法與我方可切入的差異化位置
- 說明消費者為什麼會選這個商品而非替代品
- 列出若要做精準競品比較，需要補充的資料""",

    "高擴散內容角度": """\
請提供五個高擴散內容角度，每個角度包含：觸發類型、最適合平台、Hook（開場句）、內容主軸、互動設計、CTA。

接著提供：
- 三則可直接使用的高擴散短文案（50-80字，每則適合不同觸發類型）
- 三個短影音開場腳本（含前三秒畫面描述與台詞）""",

    "十個標題候選與最佳推薦": """\
請提供十個不同角度的高吸引力標題（依序）：
1. 好奇型　2. 情境型　3. 情緒型　4. 問題型　5. 利益型
6. 送禮型　7. 收藏型　8. 設計感型　9. 反差型　10. 社交互動型

最後選出最推薦的標題，並說明：為什麼最有機會讓人停下來、最適合用在哪個平台或情境。""",

    "各平台文案": """\
請依序產出各平台文案：

**Instagram**（含：吸睛 hook、一段情境描寫、自然 CTA、5-8個 hashtag）

**Threads**（含：hook、一句延伸感受、可引發互動的問題或觀點）

**Facebook**（含：hook、情境故事、商品價值、問句互動、CTA）

**小紅書**（含：吸睛標題、使用心得式描述、推薦理由、收藏導向語句）

**Pinkoi 商品頁**（含：商品亮點、使用情境、設計理念、送禮說明）

**電商商品頁**（含：商品賣點、使用情境、購買理由、行動引導）

**短影音腳本**（含：三秒吸睛開場畫面+台詞、商品亮點帶出、使用情境展示、收尾 CTA）""",

    "關鍵字與 hashtag": """\
請提供：
- **社群 hashtag**（中英文混合，15-20個，分通用型與精準型）
- **搜尋關鍵字**（消費者實際會搜尋的語句，含描述詞、用途詞、送禮詞、風格詞）
- **電商導購關鍵字**（適合商品標題或描述的長尾關鍵字）
- **送禮情境關鍵字**
- **風格關鍵字**""",

    "品質檢查與可補充資訊": """\
請進行最終品質確認，以清單格式輸出各項是否達標：
- 是否有明確購買理由
- 是否有清楚差異化定位
- 是否有處理消費者猶豫
- 是否有平台文案差異
- 是否避開普通文案
- 是否具備高擴散誘因
- 是否有標示圖片無法確認的內容
- 是否把商品特徵轉成情緒與場景

最後列出：**最需要補充的商品資訊**（按優先順序排列）""",
}


class CommerceState(TypedDict):
    images: List[dict]
    extra_info: str
    plan: List[str]
    past_steps: Annotated[List[Tuple[str, str]], operator.add]
    response: str


class CommercePlan(BaseModel):
    """商品電商成長分析章節計畫"""
    steps: List[str] = Field(description="依序需要產出的分析章節名稱列表")


def _build_image_content(images: List[dict], text: str) -> list:
    content = [{"type": "text", "text": text}]
    for img in images:
        content.append({
            "type": "image_url",
            "image_url": {"url": f"data:{img['mime']};base64,{img['b64']}"},
        })
    return content


def commerce_planner(state: CommerceState):
    text = (
        "請分析這些商品圖片，決定最適合產出的分析章節順序。\n"
        "可選章節：" + "／".join(COMMERCE_SECTIONS) + "\n"
        "請列出所有章節，依最佳分析邏輯排序。"
        + (f"\n\n補充商品資訊：\n{state['extra_info']}" if state.get("extra_info") else "")
    )
    messages = [
        SystemMessage(content=COMMERCE_AGENT_SYSTEM),
        HumanMessage(content=_build_image_content(state.get("images", []), text)),
    ]
    result = llm_4o.with_structured_output(CommercePlan).invoke(messages)
    return {"plan": result.steps}


def commerce_executor(state: CommerceState):
    section = state["plan"][0]
    completed = "\n\n".join(
        f"### {s}\n{r}" for s, r in state["past_steps"]
    ) if state["past_steps"] else "（尚無已完成章節）"
    section_key = next((k for k in COMMERCE_SECTION_FORMATS if k in section), None)
    format_hint = f"\n\n**格式要求：**\n{COMMERCE_SECTION_FORMATS[section_key]}" if section_key else ""
    text = (
        (f"補充商品資訊：\n{state['extra_info']}\n\n" if state.get("extra_info") else "")
        + f"已完成章節（供參考，保持一致性）：\n{completed}\n\n"
        + f"請現在產出「{section}」的完整內容。內容要具體、可直接使用。{format_hint}"
    )
    messages = [
        SystemMessage(content=COMMERCE_AGENT_SYSTEM),
        HumanMessage(content=_build_image_content(state.get("images", []), text)),
    ]
    result = llm_4o.invoke(messages)
    return {"past_steps": [(section, result.content)]}


def commerce_replanner(state: CommerceState):
    remaining = state["plan"][1:]
    if not remaining:
        full_output = "\n\n---\n\n".join(
            f"## {section}\n\n{content}"
            for section, content in state["past_steps"]
        )
        return {"response": full_output}
    return {"plan": remaining}


def commerce_should_end(state: CommerceState):
    return "end" if state.get("response") else "continue"


def build_commerce_graph():
    wf = StateGraph(CommerceState)
    wf.add_node("planner", commerce_planner)
    wf.add_node("executor", commerce_executor)
    wf.add_node("re_planner", commerce_replanner)
    wf.set_entry_point("planner")
    wf.add_edge("planner", "executor")
    wf.add_edge("executor", "re_planner")
    wf.add_conditional_edges("re_planner", commerce_should_end, {"continue": "executor", "end": END})
    return wf.compile()


commerce_graph = build_commerce_graph()

VISION_EXTS = IMAGE_EXTS - {".svg"}


def _images_from_files(file_list) -> List[dict]:
    if not file_list:
        return []
    if not isinstance(file_list, list):
        file_list = [file_list]
    images = []
    for f in file_list:
        path = f.name if hasattr(f, "name") else str(f)
        ext = os.path.splitext(path)[1].lower()
        if ext not in VISION_EXTS:
            continue
        ext_key = ext.lstrip(".")
        mime = f"image/{MIME_MAP.get(ext_key, 'png')}"
        with open(path, "rb") as fh:
            b64 = base64.b64encode(fh.read()).decode()
        images.append({"mime": mime, "b64": b64})
    return images


def run_commerce_analysis(file_list, extra_info: str, user_prompt: str = ""):
    images = _images_from_files(file_list)
    if not images:
        yield "⚠️ 請上傳至少一張商品圖片後再送出。"
        return

    info = extra_info.strip() if extra_info else ""
    if user_prompt and user_prompt.strip():
        info = (info + "\n\n" if info else "") + f"使用者額外指示：\n{user_prompt.strip()}"

    output = ""
    plan_shown = False

    for event in commerce_graph.stream(
        {"images": images, "extra_info": info}, {"recursion_limit": 30}
    ):
        for node_name, value in event.items():
            if node_name == "planner" and "plan" in value and not plan_shown:
                plan_shown = True
                output = "### 📋 分析章節規劃\n" + "\n".join(
                    f"{i+1}. {s}" for i, s in enumerate(value["plan"])
                ) + "\n\n---\n\n*逐章節生成中……*\n\n"
                yield output
            elif node_name == "executor" and "past_steps" in value:
                for section, content in value["past_steps"]:
                    output = output.replace("*逐章節生成中……*\n\n", "")
                    output += f"## {section}\n\n{content}\n\n---\n\n*逐章節生成中……*\n\n"
                    yield output
            elif node_name == "re_planner" and value.get("response"):
                yield value["response"]


# ══════════════════════════════════════════════════════════════════════════════
# TAB 4 — 業務智能分析Agent Agent (Plan-Execute架構)
# ══════════════════════════════════════════════════════════════════════════════

SALES_AGENT_SYSTEM = """
你是一位擁有50年經驗的大數據分析大師、金牌業務以及產業顧問，對各行業都非常熟悉。
使用者是一位業務員，希望你協助他進行精準提案。

你的核心能力：
- 對資料進行深度分析，涵蓋質性與量化兩個維度
- 能看出潛在、容易被忽略的特性與隱藏信號
- 快速識別客戶輪廓、痛點、決策邏輯與購買動機
- 模擬不同業務策略的效果並給出具體建議
- 把數據轉成業務員可以直接使用的洞察與行動

你的分析原則：
- 每個分析都要連回業務行動，不只是觀察
- 數字要有意義，要說明背後的業務啟示
- 洞察要可操作，不只是描述現象
- 建議要具體，不只是給方向
- 客戶輪廓要生動，像是真實描述一個人
- 痛點要有優先級排序與解決建議
- 策略模擬要有具體執行步驟與成功機率估算

資訊不足時的處理：
- 先分析現有資料，盡力提供有用洞察
- 標示「資訊不足」的地方
- 列出建議補充的問題與資料

輸出語言：繁體中文，語氣專業但清晰，避免過度學術化。
"""

SALES_SECTIONS = [
    "資料總覽與執行摘要",
    "客戶輪廓與行為分析",
    "痛點偵測與影響量化",
    "深度量化分析",
    "深度質性分析",
    "潛在特性與隱藏洞察",
    "競爭態勢與市場定位",
    "策略模擬與情境分析",
    "精準提案行動計畫",
]

SALES_SECTION_FORMATS = {
    "資料總覽與執行摘要": (
        "用3-5句話說明資料類型、資料品質、最重要的3個發現，以及本次分析的業務價值。"
    ),
    "客戶輪廓與行為分析": """\
請識別並輸出所有關鍵利害關係人（若有多人請分別列出），每位格式如下：

**[姓名或代稱] / [職稱]**
- 角色類型：Champion（內部倡導者）/ Blocker（阻礙者）/ Gatekeeper（把關者）/ Influencer（影響者）/ Decision Maker（決策者）
- 屬性痛點：[顯性需求或問題，1-2句]
- 屬性需求：[功能性或流程性需求，1-2句]
- 隱性恐懼 (Deep Fear)：[說出或未說出的深層恐懼，最關鍵的心理障礙，1句核心描述]
- 決策邏輯：[如何做決定，1句]
- 最佳應對策略：[業務員如何針對此人的心理調整策略，1-2句]

最後列出：影響成交最關鍵的利害關係人與理由。""",
    "痛點偵測與影響量化": """\
請用表格輸出痛點清單：
| 痛點 | 類型（顯性/隱性） | 嚴重程度（1-10） | 業務影響 | 處理策略 |
|---|---|---|---|---|
最後說明：最高優先處理的痛點與理由。""",
    "深度量化分析": """\
提取所有可量化的指標，用表格整理：
| 指標 | 數值 | 意義 | 業務行動 |
|---|---|---|---|
並說明數字背後最重要的3個業務洞察。""",
    "深度質性分析": """\
請分析：
- **語氣與態度**（對方溝通語氣透露什麼訊號？）
- **優先級排序**（對方最在意什麼？次要考量是什麼？）
- **決策邏輯**（對方如何做決定？）
- **未說出口的需求**（從資料推測的潛在期望）
- **合作意願信號**（哪些跡象顯示合作意願高/低？）""",
    "潛在特性與隱藏洞察": """\
找出5-7個容易被忽略但對業務非常關鍵的特性，每個包含：
- **信號名稱**
- **觀察到的跡象**（具體指出資料中的哪個細節）
- **為什麼重要**（對業務的影響）
- **建議行動**（業務員應如何利用這個洞察）""",
    "競爭態勢與市場定位": """\
請分析：
- 客戶目前的選擇與替代方案
- 我方優勢與可能的弱點
- 競爭者的可能策略
- 我方最佳差異化切入點
- 一句話說明我方應如何定位""",
    "策略模擬與情境分析": """\
請模擬三種業務策略（每個包含：策略描述、風險等級、成功機率估算、4個執行步驟、預期效益）：

**策略A：快速成交型**

**策略B：關係建立型**

**策略C：顧問式提案型**

最後給出：根據本次分析，最推薦哪種策略組合及理由。""",
    "精準提案行動計畫": """\
請輸出：
**本週行動（Week 1）**（3個具體行動，含預期回應）
**短期計畫（Week 2-4）**（推進步驟與里程碑）
**中期計畫（Month 2-3）**（深化關係與方案）
**關鍵成功因子**（本次提案最需要掌握的3個要素）
**風險預警**（可能失敗的原因與預防策略）""",
}

# ── Tab 4 HTML dashboard: Python-rendered template (responsive) ────────────

_DASH_EXTRACT_SYSTEM = "你是資料結構化助理。只輸出符合要求的 JSON，不輸出任何說明。每一個 schema 欄位都必須輸出，不可省略、不可留空陣列。"

# NOTE: split into core + tail calls for the same reason as the Tab 5 war-room extractor —
# a single call covering all fields risked dropping the ones listed last in the schema
# (radar_*, insights, action_plan), which need the most synthesis from free text.
_DASH_EXTRACT_PROMPT_CORE = """\
請從以下業務分析報告提取資料，輸出 JSON（繁體中文）。

Schema：
{
  "subtitle": "案例或專案名稱（15字以內）",
  "data_source": "資料來源或檔案名稱（精簡）",
  "target_client": "客戶公司名稱",
  "est_value": "High（戰略級客戶）| Medium（一般客戶）| Low（觀察中）",
  "stakeholders": [
    {
      "name": "姓名或代稱",
      "title": "職稱",
      "role_type": "Champion|Blocker|Gatekeeper|Influencer|Decision Maker",
      "role_label": "中文角色說明，如 內部倡導者 (Champion)",
      "pain_points": "屬性痛點或需求，1-2句",
      "deep_fear": "隱性恐懼，1句核心描述"
    }
  ],
  "key_factor": "關鍵決定因素",
  "top_risk": "最高風險點",
  "strategies": [
    {
      "label_type": "主打|加值|商務",
      "name": "策略名稱——請直接沿用報告中「策略模擬與情境分析」章節裡實際命名／描述的策略"
              "（例如報告寫『策略A：快速成交型』就用『快速成交型』），不要自行發明報告中沒有的策略",
      "description": "一句說明，盡量納入報告提到的成功機率或風險等級等具體數字",
      "boost": "5到20之間的整數；若報告有明確指出「最推薦」的策略，該策略排第一個、"
               "label_type 設為主打並給較高的 boost，其餘依報告的風險等級／優先順序給 label_type 與 boost"
    }
  ]
}

只輸出 JSON，不要任何說明。

---
分析報告：
{analysis}
"""

_DASH_EXTRACT_PROMPT_TAIL = """\
請從以下業務分析報告，提取「成交雷達」與「隱藏洞察」兩個部分，輸出 JSON（繁體中文）。
這兩個部分即使出現在報告後段、內容較長，也必須完整提取，絕對不可省略或留空。

Schema：
{
  "radar_dimensions": ["5個維度名稱"],
  "radar_customer": [5個數字 1-10，代表客戶期望值],
  "radar_traditional": [5個數字 1-10，傳統方案覆蓋度，通常偏低],
  "insights": [{"title": "洞察標題", "description": "1-2句說明"}]
}

只輸出 JSON，不要任何說明。

---
分析報告：
{analysis}
"""

# NOTE: action_plan gets its own dedicated call (not bundled with insights) because insights
# can run to 6-7 verbose items on its own, and when both shared one call/token budget,
# action_plan (listed after insights) was observed to collapse down to a single Week-1 phase
# even though the underlying report clearly described a multi-week/month rollout.
_DASH_EXTRACT_PROMPT_ACTION = """\
請從以下業務分析報告，提取完整的「精準行動計畫」時間軸，輸出 JSON（繁體中文）。
報告中通常會分階段描述（例如本週行動、短期計畫、中期計畫等），有幾個階段就輸出幾個階段——
只要報告中有對應內容就必須完整列出每一個階段，不可只輸出第一階段就停止；
但也絕對不可以為了湊階段數而輸出報告中沒有的階段，每個階段都必須有實際的 tasks 或 milestone，
不可輸出空白或虛構的階段。phase 名稱請盡量沿用報告原文使用的階段名稱（例如報告寫「Month 2-3」
就用「Month 2-3」，不要拆成兩個階段）。

Schema（以下僅為格式示意，實際階段數與名稱請以報告內容為準）：
{
  "action_plan": [
    {"phase": "階段名稱（沿用報告原文）", "tasks": ["任務1", "任務2"], "milestone": "里程碑說明"}
  ]
}

只輸出 JSON，不要任何說明。

---
分析報告：
{analysis}
"""

_DASH_CSS = """
:root{--bg:#0d1117;--card:#161b22;--card2:#1a2035;--gold:#f0b429;--cyan:#21d4fd;
--red:#ff6b6b;--purple:#8b5cf6;--blue:#3b82f6;--orange:#f97316;
--text:#e6edf3;--muted:#8b949e;--border:rgba(33,212,253,.15)}
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
html{scroll-behavior:smooth}
body{background:var(--bg);color:var(--text);font-family:'Noto Sans TC','Inter',sans-serif;font-size:14px;line-height:1.6}
img,canvas{max-width:100%;height:auto}
.topnav{display:flex;justify-content:space-between;align-items:center;padding:0 24px;min-height:70px;background:#0d1117;border-bottom:1px solid var(--border);position:sticky;top:0;z-index:100;gap:12px;flex-wrap:wrap}
.nav-left{display:flex;flex-direction:column;gap:2px}
.nav-row1{display:flex;align-items:center;gap:10px;flex-wrap:wrap}
.nav-icon{color:var(--cyan);font-size:18px;font-weight:700}
.nav-title{font-size:16px;font-weight:700;color:var(--text)}
.nav-sep{width:1px;height:18px;background:var(--border);flex-shrink:0}
.nav-sub{font-size:13px;color:var(--muted)}
.nav-src{font-size:11px;color:var(--muted);margin-top:1px}
.nav-right{display:flex;gap:24px;flex-shrink:0}
.nav-kpi{text-align:right}
.nav-kpi-lbl{font-size:10px;color:var(--muted);letter-spacing:1px;text-transform:uppercase}
.nav-kpi-val{font-size:13px;font-weight:700;color:var(--gold)}
.page{max-width:1400px;margin:0 auto;padding:20px 16px}
.main-grid{display:grid;grid-template-columns:1fr 1.3fr 1fr;gap:16px;margin-bottom:16px;align-items:start}
.bottom-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px}
.sec-title{font-size:13px;font-weight:700;color:var(--text);margin-bottom:12px;padding-bottom:6px;border-bottom:2px solid var(--gold);display:inline-block}
.s-card{background:var(--card);border-radius:10px;padding:14px;margin-bottom:10px;border:1px solid var(--border)}
.s-header{display:flex;align-items:center;gap:10px;margin-bottom:6px}
.avatar{width:44px;height:44px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:16px;font-weight:700;flex-shrink:0}
.av-c{background:linear-gradient(135deg,#ff6b6b,#ee0979);color:#fff}
.av-b{background:linear-gradient(135deg,#8b5cf6,#6d28d9);color:#fff}
.av-i{background:linear-gradient(135deg,#3b82f6,#1d4ed8);color:#fff}
.av-d{background:linear-gradient(135deg,#f0b429,#d97706);color:#000}
.av-x{background:linear-gradient(135deg,#64748b,#475569);color:#fff}
.s-name{font-size:14px;font-weight:700;color:var(--text)}
.s-ttl{font-size:11px;color:var(--muted)}
.badge{display:inline-block;padding:2px 10px;border-radius:20px;font-size:11px;font-weight:600;margin-bottom:8px}
.bC{background:#ff6b6b;color:#fff}.bB{background:#8b5cf6;color:#fff}
.bI{background:#3b82f6;color:#fff}.bD{background:#f0b429;color:#000}.bX{background:#64748b;color:#fff}
.pain-lbl{font-size:11px;color:var(--muted);margin-bottom:3px}
.pain-txt{font-size:12px;color:var(--text);margin-bottom:8px;line-height:1.5}
.fear-box{background:rgba(249,115,22,.08);border-left:3px solid var(--orange);border-radius:4px;padding:8px 10px}
.fear-title{font-size:11px;font-weight:700;color:var(--orange);margin-bottom:2px}
.fear-txt{font-size:12px;color:var(--orange);font-style:italic;line-height:1.4}
.radar-wrap{background:var(--card);border-radius:12px;padding:16px;border:1px solid var(--border)}
.radar-canvas-box{position:relative;height:270px;margin-bottom:12px}
.radar-meta{display:grid;grid-template-columns:1fr 1fr;gap:8px}
.rmeta-card{background:var(--card2);border-radius:8px;padding:10px 12px}
.rmeta-lbl{font-size:10px;color:var(--muted);margin-bottom:2px}
.rmeta-val{font-size:13px;font-weight:700;color:var(--text)}
.rmeta-val.risk{color:var(--red)}
.sim-wrap{background:var(--card);border-radius:12px;padding:16px;border:1px solid var(--border);display:flex;flex-direction:column}
.sim-top{display:flex;justify-content:space-between;align-items:center;margin-bottom:4px}
.sim-title{font-size:14px;font-weight:700;color:var(--text)}
.sim-btn{font-size:11px;border:1px solid var(--gold);background:transparent;color:var(--gold);padding:3px 10px;border-radius:4px}
.sim-desc{font-size:11px;color:var(--muted);margin-bottom:14px;line-height:1.4}
.st-item{display:flex;align-items:flex-start;gap:8px;margin-bottom:12px;cursor:pointer}
.st-item input{width:16px;height:16px;margin-top:2px;accent-color:var(--cyan);flex-shrink:0;cursor:pointer}
.st-tag{font-size:10px;padding:1px 7px;border-radius:10px;font-weight:600;white-space:nowrap;margin-top:3px;flex-shrink:0}
.tag-m{background:rgba(240,180,41,.2);color:var(--gold)}
.tag-p{background:rgba(59,130,246,.2);color:var(--blue)}
.tag-b{background:rgba(100,116,139,.2);color:#94a3b8}
.st-name{font-size:13px;font-weight:700;color:var(--text)}
.st-desc{font-size:11px;color:var(--muted);line-height:1.4}
.win-card{background:var(--bg);border-radius:10px;padding:14px 16px;margin-top:auto;border:1px solid var(--border)}
.win-row{display:flex;justify-content:space-between;align-items:center;margin-bottom:8px}
.win-lbl{font-size:12px;color:var(--muted)}
.win-pct{font-size:2rem;font-weight:800;color:var(--gold);line-height:1}
.win-track{height:8px;background:#1e293b;border-radius:4px;overflow:hidden;margin-bottom:6px}
.win-bar{height:100%;border-radius:4px;background:linear-gradient(90deg,#ff6b6b,#f97316,#f0b429,#22c55e);transition:width .4s}
.win-note{font-size:11px;color:var(--muted);line-height:1.4}
.ins-card{border:1px solid var(--gold);background:rgba(240,180,41,.04);border-radius:12px;padding:20px}
.ins-title{font-size:14px;font-weight:700;color:var(--gold);margin-bottom:14px}
.ins-item{display:flex;gap:10px;margin-bottom:12px}
.ins-dot{width:8px;height:8px;border-radius:50%;background:var(--cyan);flex-shrink:0;margin-top:5px}
.ins-name{font-size:13px;font-weight:700;color:var(--text);margin-bottom:2px}
.ins-desc{font-size:12px;color:var(--muted);line-height:1.5}
.act-card{background:var(--card);border-radius:12px;padding:20px;border:1px solid var(--border)}
.act-title{font-size:14px;font-weight:700;color:var(--text);margin-bottom:16px}
.tl{display:flex;gap:0;overflow-x:auto}
.tl-col{flex:1;min-width:110px;position:relative;padding-top:24px;padding-right:8px}
.tl-col::before{content:'';position:absolute;top:7px;left:0;right:0;height:2px;background:var(--border)}
.tl-col:first-child::before{left:50%}.tl-col:last-child::before{right:50%}
.tl-dot{width:16px;height:16px;border-radius:50%;background:var(--gold);border:3px solid var(--bg);position:absolute;top:0;left:50%;transform:translateX(-50%)}
.tl-phase{font-size:10px;font-weight:700;color:var(--gold);text-align:center;margin-bottom:6px}
.tl-tasks{font-size:11px;color:var(--text);line-height:1.6}
.tl-ms{font-size:10px;color:var(--muted);margin-top:4px;font-style:italic}
@media(max-width:1100px){
  .main-grid{grid-template-columns:1fr 1fr}
  .main-grid>:last-child{grid-column:1/-1}
}
@media(max-width:768px){
  .main-grid,.bottom-grid{grid-template-columns:1fr}
  .main-grid>:last-child{grid-column:auto}
  .topnav{min-height:auto;padding:12px 16px;flex-direction:column;align-items:flex-start}
  .nav-right{width:100%;justify-content:flex-start;gap:20px}
}
@media(max-width:480px){
  .page{padding:10px 8px}
  .win-pct{font-size:1.5rem}
  .tl{flex-direction:column;gap:8px}
  .tl-col{padding-top:0;padding-left:16px;border-left:2px solid var(--border);margin-left:8px}
  .tl-col::before,.tl-dot{display:none}
}
"""


def _role_cls(role_type: str):
    r = role_type.lower()
    if "champion" in r or "倡導" in r:
        return "av-c", "bC"
    if any(x in r for x in ("blocker", "gatekeeper", "阻礙", "把關", "風險控制")):
        return "av-b", "bB"
    if "influencer" in r or "影響" in r:
        return "av-i", "bI"
    if "decision" in r or "決策" in r:
        return "av-d", "bD"
    return "av-x", "bX"


def _render_stakeholders(items: list) -> str:
    if not items:
        return '<div class="s-card"><p style="color:var(--muted)">（未識別到明確利害關係人）</p></div>'
    out = []
    for s in items:
        av, bv = _role_cls(s.get("role_type", ""))
        out.append(
            f'<div class="s-card">'
            f'<div class="s-header">'
            f'<div class="avatar {av}">{s.get("name","?")[:1]}</div>'
            f'<div><div class="s-name">{s.get("name","")}</div>'
            f'<div class="s-ttl">{s.get("title","")}</div></div></div>'
            f'<span class="badge {bv}">{s.get("role_label", s.get("role_type",""))}</span>'
            f'<div class="pain-lbl">⊕ 屬性痛點</div>'
            f'<div class="pain-txt">{s.get("pain_points","")}</div>'
            f'<div class="fear-box">'
            f'<div class="fear-title">🔮 隱性恐懼 (Deep Fear)</div>'
            f'<div class="fear-txt">{s.get("deep_fear","")}</div></div></div>'
        )
    return "".join(out)


def _render_strategies(items: list):
    if not items:
        items = [{"label_type": "主打", "name": "顧問式提案", "description": "深度理解需求後提出完整解決方案。", "boost": 15}]
    out, prob = [], 35
    for i, s in enumerate(items):
        checked = "checked" if i == 0 else ""
        if i == 0:
            prob = min(35 + s.get("boost", 0), 92)
        lt = s.get("label_type", "主打")
        tc = "tag-m" if lt == "主打" else ("tag-p" if lt == "加值" else "tag-b")
        out.append(
            f'<label class="st-item">'
            f'<input type="checkbox" class="strategy-check" '
            f'data-boost="{s.get("boost",10)}" data-note="{s.get("name","")}" {checked}>'
            f'<span class="st-tag {tc}">{lt}</span>'
            f'<div><div class="st-name">{s.get("name","")}</div>'
            f'<div class="st-desc">{s.get("description","")}</div></div></label>'
        )
    return "".join(out), prob


def _render_insights(items: list) -> str:
    return "".join(
        f'<div class="ins-item"><div class="ins-dot"></div>'
        f'<div><div class="ins-name">{i.get("title","")}</div>'
        f'<div class="ins-desc">{i.get("description","")}</div></div></div>'
        for i in items
    ) if items else '<p style="color:var(--muted);font-size:12px">（無隱藏洞察資料）</p>'


def _render_timeline(items: list) -> str:
    items = [p for p in (items or []) if p.get("tasks") or p.get("milestone")]
    if not items:
        items = [
            {"phase": "Week 1", "tasks": ["初步接觸", "需求確認"], "milestone": "建立信任"},
            {"phase": "Week 2-4", "tasks": ["提案準備", "方案調整"], "milestone": "提案送出"},
            {"phase": "Month 2", "tasks": ["PoC 試點", "效益驗證"], "milestone": "初步成效"},
            {"phase": "Month 3+", "tasks": ["正式合作", "擴展計畫"], "milestone": "簽約落地"},
        ]
    return "".join(
        f'<div class="tl-col"><div class="tl-dot"></div>'
        f'<div class="tl-phase">{p.get("phase","")}</div>'
        f'<div class="tl-tasks">{"".join(f"・{t}<br>" for t in p.get("tasks",[]))}</div>'
        f'<div class="tl-ms">{p.get("milestone","")}</div></div>'
        for p in items
    )


def _extract_dashboard_json_stream(analysis: str):
    """Generator: yields (data_so_far, status_message) at each extraction stage.

    Structured as a generator (rather than a plain function with an on_step callback) so the
    Gradio-facing caller can yield a status update to a visible Markdown component between each
    OpenAI call — a callback fired from inside a blocking call can't make an outer generator
    pause and yield, so the calls themselves have to live in the generator.
    """
    import json as _j
    text = analysis[:30000]
    data: dict = {}
    yield data, "分析客戶輪廓、痛點與策略..."
    try:
        resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": _DASH_EXTRACT_SYSTEM},
                {"role": "user", "content": _DASH_EXTRACT_PROMPT_CORE.replace("{analysis}", text)},
            ],
            response_format={"type": "json_object"},
            max_tokens=2000,
            temperature=0,
        )
        data.update(_j.loads(resp.choices[0].message.content))
    except Exception:
        pass
    yield data, "分析成交雷達與隱藏洞察..."
    try:
        resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": _DASH_EXTRACT_SYSTEM},
                {"role": "user", "content": _DASH_EXTRACT_PROMPT_TAIL.replace("{analysis}", text)},
            ],
            response_format={"type": "json_object"},
            max_tokens=1500,
            temperature=0,
        )
        data.update(_j.loads(resp.choices[0].message.content))
    except Exception:
        pass
    yield data, "分析精準行動計畫..."
    try:
        resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": _DASH_EXTRACT_SYSTEM},
                {"role": "user", "content": _DASH_EXTRACT_PROMPT_ACTION.replace("{analysis}", text)},
            ],
            response_format={"type": "json_object"},
            max_tokens=1500,
            temperature=0,
        )
        data.update(_j.loads(resp.choices[0].message.content))
    except Exception:
        pass
    yield data, "整理資料中..."


def _extract_dashboard_json(analysis: str) -> dict:
    data: dict = {}
    for data, _status in _extract_dashboard_json_stream(analysis):
        pass
    return data


def _render_dashboard_html(data: dict) -> str:
    import json as _j
    subtitle    = data.get("subtitle", "業務智能分析")
    data_source = data.get("data_source", "業務分析資料")
    target_client = data.get("target_client", "目標客戶")
    est_value   = data.get("est_value", "High（戰略級客戶）")
    key_factor  = data.get("key_factor", "—")
    top_risk    = data.get("top_risk", "—")

    dims   = data.get("radar_dimensions", ["客戶需求", "整合難度", "操作性", "導入風險", "效益價值"])
    r_cust = data.get("radar_customer",   [8, 7, 6, 9, 8])
    r_trad = data.get("radar_traditional",[5, 4, 6, 3, 5])
    n = max(len(dims), len(r_cust), len(r_trad))
    dims   = (dims   + ["—"] * n)[:n]
    r_cust = (r_cust + [5]   * n)[:n]
    r_trad = (r_trad + [5]   * n)[:n]

    s_html           = _render_stakeholders(data.get("stakeholders", []))
    strategies_html, init_prob = _render_strategies(data.get("strategies", []))
    ins_html         = _render_insights(data.get("insights", []))
    tl_html          = _render_timeline(data.get("action_plan", []))

    parts = [
        "<!DOCTYPE html>",
        "<html lang='zh-TW'>",
        "<head>",
        '<meta charset="UTF-8">',
        '<meta name="viewport" content="width=device-width, initial-scale=1.0">',
        "<title>B2B 戰略情報室</title>",
        '<link href="https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@300;400;500;700'
        '&family=Inter:wght@300;400;500;700&display=swap" rel="stylesheet">',
        '<script src="https://cdn.jsdelivr.net/npm/chart.js"></script>',
        f"<style>{_DASH_CSS}</style>",
        "</head>",
        "<body>",
        # ── Nav ──
        '<nav class="topnav">',
        '<div class="nav-left">',
        '<div class="nav-row1">',
        '<span class="nav-icon">≈</span>',
        '<span class="nav-title">B2B 戰略情報室</span>',
        '<div class="nav-sep"></div>',
        f'<span class="nav-sub">{subtitle}</span>',
        '</div>',
        f'<div class="nav-src">基於「{data_source}」之大數據與心理學深度解析</div>',
        '</div>',
        '<div class="nav-right">',
        f'<div class="nav-kpi"><div class="nav-kpi-lbl">TARGET CLIENT</div><div class="nav-kpi-val">{target_client}</div></div>',
        f'<div class="nav-kpi"><div class="nav-kpi-lbl">EST. VALUE</div><div class="nav-kpi-val">{est_value}</div></div>',
        '</div></nav>',
        # ── Page ──
        '<div class="page">',
        '<div class="main-grid">',
        # Col 1: Stakeholders
        f'<div><div class="sec-title">利害關係人矩陣</div>{s_html}</div>',
        # Col 2: Radar
        '<div class="radar-wrap">',
        '<div class="sec-title">專案需求維度解析</div>',
        '<div class="radar-canvas-box"><canvas id="radarChart"></canvas></div>',
        '<div class="radar-meta">',
        f'<div class="rmeta-card"><div class="rmeta-lbl">關鍵決定因素</div><div class="rmeta-val">{key_factor}</div></div>',
        f'<div class="rmeta-card"><div class="rmeta-lbl">最高風險點</div><div class="rmeta-val risk">{top_risk}</div></div>',
        '</div></div>',
        # Col 3: Strategy sim
        '<div class="sim-wrap">',
        '<div class="sim-top"><span class="sim-title">🏆 提案策略模擬器</span><span class="sim-btn">互動測試</span></div>',
        '<div class="sim-desc">勾選您預計在提案書中主打的策略，預測「贏單機率」與「客戶決策阻力」。</div>',
        strategies_html,
        '<div class="win-card">',
        '<div class="win-row"><span class="win-lbl">預測贏單機率 (Win Probability)</span>',
        f'<span class="win-pct" id="win-prob">{init_prob}%</span></div>',
        f'<div class="win-track"><div class="win-bar" id="win-bar" style="width:{init_prob}%"></div></div>',
        '<div class="win-note" id="win-note">已強化優勢：請勾選策略以啟動預測。</div>',
        '</div></div>',
        '</div>',  # end main-grid
        # ── Bottom ──
        '<div class="bottom-grid">',
        f'<div class="ins-card"><div class="ins-title">⚡ 潛在機會：容易被忽略的關鍵信號</div>{ins_html}</div>',
        f'<div class="act-card"><div class="act-title">精準行動計畫</div><div class="tl">{tl_html}</div></div>',
        '</div>',  # end bottom-grid
        '</div>',  # end page
        # ── JS ──
        "<script>",
        "function updateWinProb(){",
        "var base=35,total=base,notes=[];",
        "document.querySelectorAll('.strategy-check:checked').forEach(function(cb){",
        "total+=parseInt(cb.dataset.boost||0);",
        "if(cb.dataset.note)notes.push(cb.dataset.note);});",
        "total=Math.min(total,92);",
        "document.getElementById('win-prob').textContent=total+'%';",
        "document.getElementById('win-bar').style.width=total+'%';",
        "document.getElementById('win-note').textContent=notes.length",
        "?'已強化優勢：'+notes.join('、')+'。'",
        ":'請勾選策略以啟動預測。';}",
        "document.querySelectorAll('.strategy-check').forEach(function(cb){cb.addEventListener('change',updateWinProb);});",
        "updateWinProb();",
        f"var dims={_j.dumps(dims,ensure_ascii=False)};",
        f"var cust={_j.dumps(r_cust)};",
        f"var trad={_j.dumps(r_trad)};",
        "new Chart(document.getElementById('radarChart'),{type:'radar',",
        "data:{labels:dims,datasets:[",
        "{label:'客戶期望值（痛點強烈度）',data:cust,borderColor:'#21d4fd',",
        "backgroundColor:'rgba(33,212,253,.15)',pointBackgroundColor:'#21d4fd',borderWidth:2},",
        "{label:'傳統方案覆蓋度',data:trad,borderColor:'#ff6b6b',",
        "backgroundColor:'rgba(255,107,107,.1)',pointBackgroundColor:'#ff6b6b',borderWidth:2}]},",
        "options:{responsive:true,maintainAspectRatio:false,",
        "plugins:{legend:{labels:{color:'#8b949e',font:{size:11}}}},",
        "scales:{r:{angleLines:{color:'rgba(255,255,255,.1)'},grid:{color:'rgba(255,255,255,.1)'},",
        "pointLabels:{color:'#e6edf3',font:{size:11}},ticks:{display:false},min:0,max:10}}}});",
        "</script>",
        "</body></html>",
    ]
    return "\n".join(parts)



class SalesState(TypedDict):
    input: str
    plan: List[str]
    past_steps: Annotated[List[Tuple[str, str]], operator.add]
    response: str


class SalesPlan(BaseModel):
    """業務分析章節計畫"""
    steps: List[str] = Field(description="依序需要產出的分析章節名稱列表")


def sales_planner(state: SalesState):
    prompt = ChatPromptTemplate.from_messages([
        ("system", SALES_AGENT_SYSTEM),
        ("human",
         "請分析以下資料，決定最適合的分析章節順序。\n"
         "可選章節：" + "／".join(SALES_SECTIONS) + "\n"
         "請列出所有章節，依最佳分析邏輯排序。\n\n"
         "資料內容：\n{input}")
    ])
    planner = prompt | llm_4o.with_structured_output(SalesPlan)
    result = planner.invoke({"input": state["input"]})
    return {"plan": result.steps}


def sales_executor(state: SalesState):
    section = state["plan"][0]
    completed = "\n\n".join(
        f"### {s}\n{r}" for s, r in state["past_steps"]
    ) if state["past_steps"] else "（尚無已完成章節）"
    section_key = next((k for k in SALES_SECTION_FORMATS if k in section), None)
    format_hint = f"\n\n**格式要求：**\n{SALES_SECTION_FORMATS[section_key]}" if section_key else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", SALES_AGENT_SYSTEM),
        ("human",
         "原始資料：\n{input}\n\n"
         "已完成章節（供參考，保持一致性）：\n{completed}\n\n"
         "請現在產出「{section}」的完整內容。內容要具體、可操作、直接對業務員有用。{format_hint}")
    ])
    chain = prompt | llm_4o
    result = chain.invoke({
        "input": state["input"],
        "completed": completed,
        "section": section,
        "format_hint": format_hint,
    })
    return {"past_steps": [(section, result.content)]}


def sales_replanner(state: SalesState):
    remaining = state["plan"][1:]
    if not remaining:
        full_output = "\n\n---\n\n".join(
            f"## {section}\n\n{content}"
            for section, content in state["past_steps"]
        )
        return {"response": full_output}
    return {"plan": remaining}


def sales_should_end(state: SalesState):
    return "end" if state.get("response") else "continue"


def build_sales_graph():
    wf = StateGraph(SalesState)
    wf.add_node("planner", sales_planner)
    wf.add_node("executor", sales_executor)
    wf.add_node("re_planner", sales_replanner)
    wf.set_entry_point("planner")
    wf.add_edge("planner", "executor")
    wf.add_edge("executor", "re_planner")
    wf.add_conditional_edges("re_planner", sales_should_end, {"continue": "executor", "end": END})
    return wf.compile()


sales_graph = build_sales_graph()


def run_sales_analysis(file_list, user_context: str, user_prompt: str = ""):
    parts = []
    if file_list:
        file_text = extract_files_text(file_list)
        if file_text.strip():
            parts.append(file_text)
    if user_context and user_context.strip():
        parts.append(f"補充背景資訊：\n{user_context.strip()}")
    if user_prompt and user_prompt.strip():
        parts.append(f"使用者額外指示：\n{user_prompt.strip()}")

    text = "\n\n---\n\n".join(parts)
    if not text.strip():
        yield "⚠️ 請上傳附件或輸入資料後再送出。"
        return

    output = ""
    plan_shown = False

    for event in sales_graph.stream({"input": text}, {"recursion_limit": 40}):
        for node_name, value in event.items():
            if node_name == "planner" and "plan" in value and not plan_shown:
                plan_shown = True
                output = "### 📋 分析章節規劃\n" + "\n".join(
                    f"{i+1}. {s}" for i, s in enumerate(value["plan"])
                ) + "\n\n---\n\n*逐章節分析中……*\n\n"
                yield output
            elif node_name == "executor" and "past_steps" in value:
                for section, content in value["past_steps"]:
                    output = output.replace("*逐章節分析中……*\n\n", "")
                    output += f"## {section}\n\n{content}\n\n---\n\n*逐章節分析中……*\n\n"
                    yield output
            elif node_name == "re_planner" and value.get("response"):
                yield value["response"]


def generate_sales_dashboard_html(analysis_content: str):
    if not analysis_content or "逐章節分析中" in analysis_content or analysis_content.startswith("*送出"):
        return None
    data = _extract_dashboard_json(analysis_content)
    html_content = _render_dashboard_html(data)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode="w", encoding="utf-8")
    tmp.write(html_content)
    tmp.close()
    return tmp.name


def docx_to_dashboard_html_stream(file_obj):
    """Generator: yields (html_path_or_None, status_message) while converting a .docx to
    an interactive HTML dashboard. html_path is only non-None on the final yield."""
    if file_obj is None:
        yield None, "⚠️ 請先上傳 Word 檔案。"
        return
    yield None, "讀取 Word 文件..."
    from docx import Document as _DocxDocument
    doc = _DocxDocument(file_obj.name)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    text = "\n".join(paragraphs)
    if not text.strip():
        yield None, "⚠️ 無法從 Word 檔案讀取到文字內容，請確認檔案格式。"
        return
    data = {}
    for data, status in _extract_dashboard_json_stream(text):
        yield None, status
    yield None, "產生互動式 HTML 儀表板..."
    html_content = _render_dashboard_html(data)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode="w", encoding="utf-8")
    tmp.write(html_content)
    tmp.close()
    yield tmp.name, "✅ 完成！"


def docx_to_dashboard_html(file_obj):
    path = None
    for path, _status in docx_to_dashboard_html_stream(file_obj):
        pass
    return path


def convert_word_to_dashboard_tab4(file_obj):
    for path, status in docx_to_dashboard_html_stream(file_obj):
        if path:
            yield gr.update(value=path, visible=True), status
        elif status.startswith("⚠️"):
            yield gr.update(visible=False), status
        else:
            yield gr.update(visible=False), f"🔄 {status}"


# ══════════════════════════════════════════════════════════════════════════════
# TAB 5 — 數據分析Agent（業務提案戰情室）(Plan-Execute架構)
# ══════════════════════════════════════════════════════════════════════════════

DATA_AGENT_SYSTEM = """
你是一位資深業務戰情分析師與資料科學顧問，名稱為「數據分析Agent」。
你的任務是協助業務即時判斷客戶狀態、成交機會、提案風險、決策鏈、追蹤優先順序與下一步策略，
建立一個可用於主管簡報與業務日常追蹤的「業務提案戰情室」分析報告。

## 核心目標
整合資料證據、客戶洞察、風險判斷、決策鏈、提案策略與情境模擬，
協助業務判斷：
- 客戶目前卡在哪個環節
- 成交機會與主要阻力
- 不同提案策略的優劣與話術
- 風險、ROI、時程與決策壓力的變化
- 下一步該補什麼資料、該說服誰、該何時追蹤、該傳什麼訊息

## 核心體驗
- 一眼掌握客戶目前卡關位置
- 快速判斷成交機會與主要阻力
- 切換不同提案策略並比較結果
- 模擬風險、ROI、時程與決策壓力變化
- 產出可直接用於客戶追蹤的話術與行動表

## 分析原則
- 每個判斷都要連回業務行動，不只是描述現象
- 每個判斷都要標示所依據的資料線索
- 成交機率、ROI、風險等級一律採「相對評分」，不得假裝為精確預測
- 資料不足時，明確列出待補數字與待確認條件，不得憑空捏造
- 策略模擬需清楚標示為「策略推估」
- 語氣專業、精簡、具主管簡報質感，避免學術化贅字

輸出語言：繁體中文。
"""

DATA_SECTIONS = [
    "戰情總覽",
    "客戶輪廓",
    "需求痛點",
    "決策鏈",
    "風險地圖",
    "策略模擬",
    "ROI情境",
    "行動追蹤",
    "證據資料",
]

DATA_SECTION_FORMATS = {
    "戰情總覽": """\
請用以下格式輸出：
| 判斷項目 | 內容 |
|---|---|
| 客戶／案子名稱 |  |
| 目前決策階段 |  |
| 成交機會（相對評分 0-100） |  |
| 風險等級（相對評分 0-100） |  |
| 目前最大成交機會 |  |
| 目前最大風險 |  |
| 建議下一步 |  |

最後列出「主管簡報重點」3-5條，每條一句話，適合投影片使用。""",
    "客戶輪廓": """\
請說明：
- **產業與公司背景**
- **公司規模／營運概況**
- **目前現況**（與本次提案相關的營運狀態）
- **對本次提案的態度與訊號**
若資訊不足，請標示「需補充」。""",
    "需求痛點": """\
請用以下表格格式輸出：
| 痛點 | 嚴重程度（1-10） | 對客戶的影響 | 資料依據 |
|---|---|---|---|
最後說明：最優先需要解決的痛點與理由。""",
    "決策鏈": """\
請識別所有關鍵利害關係人，每位格式如下：

**[姓名或代稱] / [職稱]**
- 角色類型：Champion（內部倡導者）/ Blocker（阻礙者）/ Gatekeeper（把關者）/ Influencer（影響者）/ Decision Maker（決策者）
- 影響力（1-10）
- 目前支持度（1-10，10為完全支持）
- 說服重點：[如何針對此人調整策略，1-2句]

最後列出：影響成交最關鍵的人與理由。""",
    "風險地圖": """\
請用以下表格格式輸出：
| 風險 | 等級（高/中/低） | 類別（價格/時程/技術/人員/決策） | 緩解措施 | 資料依據 |
|---|---|---|---|---|
最後說明：目前最需要立即處理的風險。""",
    "策略模擬": """\
請模擬三種提案策略，每種策略需包含：適合情境、成交機會（相對評分）、主要風險、建議話術（可直接使用的一段話）、下一步行動。
所有結果請標示為「策略推估」。

**策略一：保守試點型**（先小規模驗證，降低客戶疑慮）

**策略二：快速導入型**（強調效率與時效，加快推進速度）

**策略三：主管決策型**（直接訴求高層價值與決策效益）

最後給出：根據本次分析，最推薦哪個策略及理由。""",
    "ROI情境": """\
請模擬三種 ROI 情境（保守效益／一般效益／積極效益），每種需包含：
可能節省成本（每年，新台幣萬元）、品質改善價值、客訴下降價值、回收期估算（月）、需要確認的關鍵數字。
請標示為「策略推估」，並列出計算所依據的假設。""",
    "行動追蹤": """\
請輸出：
**追蹤優先順序模擬**
- 最該先補的資料
- 最該先說服的角色
- 最適合的追蹤時間
- 最合適的下一封訊息（可直接使用的訊息草稿）

**行動項目清單**
| 行動項目 | 負責人 | 建議時間 |
|---|---|---|""",
    "證據資料": """\
請列出支持以上判斷的關鍵資料線索，每條包含：資料摘要、對應章節結論、原始資料來源。
最後列出「待補資料清單」，說明資訊不足處與建議補充問題。""",
}


class DataState(TypedDict):
    input: str
    plan: List[str]
    past_steps: Annotated[List[Tuple[str, str]], operator.add]
    response: str


class DataPlan(BaseModel):
    """業務提案戰情室章節計畫"""
    steps: List[str] = Field(description="依序需要產出的分析章節名稱列表")


def data_planner(state: DataState):
    prompt = ChatPromptTemplate.from_messages([
        ("system", DATA_AGENT_SYSTEM),
        ("human",
         "請分析以下資料，決定最適合的戰情室分析章節順序。\n"
         "可選章節：" + "／".join(DATA_SECTIONS) + "\n"
         "請列出所有章節，依最佳分析邏輯排序。\n\n"
         "資料內容：\n{input}")
    ])
    planner = prompt | llm_4o.with_structured_output(DataPlan)
    result = planner.invoke({"input": state["input"]})
    return {"plan": result.steps}


def data_executor(state: DataState):
    section = state["plan"][0]
    completed = "\n\n".join(
        f"### {s}\n{r}" for s, r in state["past_steps"]
    ) if state["past_steps"] else "（尚無已完成章節）"
    section_key = next((k for k in DATA_SECTION_FORMATS if k in section), None)
    format_hint = f"\n\n**格式要求：**\n{DATA_SECTION_FORMATS[section_key]}" if section_key else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", DATA_AGENT_SYSTEM),
        ("human",
         "原始資料：\n{input}\n\n"
         "已完成章節（供參考，保持一致性）：\n{completed}\n\n"
         "請現在產出「{section}」的完整內容。內容要具體、可操作、直接對業務員有用。{format_hint}")
    ])
    chain = prompt | llm_4o
    result = chain.invoke({
        "input": state["input"],
        "completed": completed,
        "section": section,
        "format_hint": format_hint,
    })
    return {"past_steps": [(section, result.content)]}


def data_replanner(state: DataState):
    remaining = state["plan"][1:]
    if not remaining:
        full_output = "\n\n---\n\n".join(
            f"## {section}\n\n{content}"
            for section, content in state["past_steps"]
        )
        return {"response": full_output}
    return {"plan": remaining}


def data_should_end(state: DataState):
    return "end" if state.get("response") else "continue"


def build_data_graph():
    wf = StateGraph(DataState)
    wf.add_node("planner", data_planner)
    wf.add_node("executor", data_executor)
    wf.add_node("re_planner", data_replanner)
    wf.set_entry_point("planner")
    wf.add_edge("planner", "executor")
    wf.add_edge("executor", "re_planner")
    wf.add_conditional_edges("re_planner", data_should_end, {"continue": "executor", "end": END})
    return wf.compile()


data_graph = build_data_graph()


def run_data_analysis(file_list, user_context: str, user_prompt: str = ""):
    parts = []
    if file_list:
        file_text = extract_files_text(file_list)
        if file_text.strip():
            parts.append(file_text)
    if user_context and user_context.strip():
        parts.append(f"補充背景資訊：\n{user_context.strip()}")
    if user_prompt and user_prompt.strip():
        parts.append(f"使用者額外指示：\n{user_prompt.strip()}")

    text = "\n\n---\n\n".join(parts)
    if not text.strip():
        yield "⚠️ 請上傳附件或輸入資料後再送出。"
        return

    output = ""
    plan_shown = False

    for event in data_graph.stream({"input": text}, {"recursion_limit": 40}):
        for node_name, value in event.items():
            if node_name == "planner" and "plan" in value and not plan_shown:
                plan_shown = True
                output = "### 📋 戰情室章節規劃\n" + "\n".join(
                    f"{i+1}. {s}" for i, s in enumerate(value["plan"])
                ) + "\n\n---\n\n*逐章節分析中……*\n\n"
                yield output
            elif node_name == "executor" and "past_steps" in value:
                for section, content in value["past_steps"]:
                    output = output.replace("*逐章節分析中……*\n\n", "")
                    output += f"## {section}\n\n{content}\n\n---\n\n*逐章節分析中……*\n\n"
                    yield output
            elif node_name == "re_planner" and value.get("response"):
                yield value["response"]


# ── Tab 5 HTML war-room dashboard: Python-rendered template (responsive) ───

_WR_EXTRACT_SYSTEM = "你是資料結構化助理。只輸出符合要求的 JSON，不輸出任何說明。每一個 schema 欄位都必須輸出，不可省略、不可留空陣列。"

# NOTE: extraction is split into two calls (core + tail) on purpose. A single call covering
# all 18 top-level fields was empirically dropping the fields that are both (a) listed last in
# the schema and (b) require the most synthesis (radar_*, follow_up, action_items, evidence),
# even when well under the max_tokens budget and with valid JSON otherwise returned. Splitting
# gives each half its own dedicated token budget and model attention.
_WR_EXTRACT_PROMPT_CORE = """\
請從以下業務提案戰情室分析報告提取資料，輸出 JSON（繁體中文）。
所有評分一律為「相對評分」，數字不足時請合理估算並不得留空。

Schema：
{
  "hero": {
    "client_name": "客戶或案子名稱",
    "industry_title": "客戶產業／情境標題（15字以內）",
    "decision_stage": "目前決策階段（簡短）",
    "main_opportunity": "最大成交機會（一句話）",
    "top_risk": "最大風險（一句話）",
    "next_step": "建議下一步（一句話）"
  },
  "kpi": {
    "deal_score": 0到100整數（成交機會）,
    "risk_score": 0到100整數（風險等級）,
    "readiness_score": 0到100整數（導入準備度）,
    "roi_potential": 0到100整數（ROI潛力）
  },
  "exec_summary": ["主管簡報重點1", "重點2", "重點3"],
  "data_source": "資料來源或檔案名稱",
  "est_value": "High（戰略級客戶）| Medium（一般客戶）| Low（觀察中）",
  "customer_profile": {
    "industry": "產業別",
    "company_size": "公司規模",
    "current_situation": "目前現況",
    "background": "背景說明"
  },
  "pain_points": [
    {"title": "痛點名稱", "description": "說明", "severity": "高|中|低", "evidence_ref": "資料依據"}
  ],
  "decision_chain": [
    {"name": "姓名或代稱", "title": "職稱", "role_type": "Champion|Blocker|Gatekeeper|Influencer|Decision Maker",
     "role_label": "中文角色說明，如 內部倡導者 (Champion)", "influence": 1到10, "support_level": 1到10, "notes": "說服重點"}
  ],
  "risk_map": [
    {"risk": "風險名稱", "level": "高|中|低", "category": "價格|時程|技術|人員|決策", "mitigation": "緩解措施", "evidence_ref": "資料依據"}
  ],
  "strategy_modes": {
    "conservative": {"title": "保守試點型", "fit_scenario": "適合情境", "deal_chance": 0到100整數, "main_risk": "主要風險", "talk_track": "建議話術，可直接使用的一段話", "next_action": "下一步行動"},
    "fast": {"title": "快速導入型", "fit_scenario": "...", "deal_chance": 0到100整數, "main_risk": "...", "talk_track": "...", "next_action": "..."},
    "executive": {"title": "主管決策型", "fit_scenario": "...", "deal_chance": 0到100整數, "main_risk": "...", "talk_track": "...", "next_action": "..."}
  },
  "risk_factors": [
    {"id": "price", "label": "價格敏感度", "default": 0到100, "inverse": false, "data_needed_if_high": "價格敏感度偏高時，需要補強的資料"},
    {"id": "downtime", "label": "停線風險", "default": 0到100, "inverse": false, "data_needed_if_high": "..."},
    {"id": "complexity", "label": "導入麻煩程度", "default": 0到100, "inverse": false, "data_needed_if_high": "..."},
    {"id": "support", "label": "決策者支持度", "default": 0到100, "inverse": true, "data_needed_if_high": "決策者支持度偏低時，需要補強的資料"},
    {"id": "resistance", "label": "現場人員排斥程度", "default": 0到100, "inverse": false, "data_needed_if_high": "..."}
  ],
  "roi_scenarios": {
    "conservative": {"cost_saving": 數字（萬元/年）, "quality_value": 數字（萬元/年）, "complaint_value": 數字（萬元/年）, "payback_months": 數字, "need_confirm": ["待確認數字1", "待確認數字2"]},
    "normal": {"cost_saving": 數字, "quality_value": 數字, "complaint_value": 數字, "payback_months": 數字, "need_confirm": ["..."]},
    "aggressive": {"cost_saving": 數字, "quality_value": 數字, "complaint_value": 數字, "payback_months": 數字, "need_confirm": ["..."]}
  }
}

只輸出 JSON，不要任何說明。

---
分析報告：
{analysis}
"""

_WR_EXTRACT_PROMPT_TAIL = """\
請從以下業務提案戰情室分析報告，提取「成交機會雷達」「追蹤優先順序模擬（行動追蹤）」與「資料證據區（證據資料）」三個部分，
輸出 JSON（繁體中文）。這三個部分即使出現在報告後段、內容較長，也必須完整提取，絕對不可省略或留空。
若報告中有「行動項目清單」表格，請完整列出每一列到 action_items；若有「證據資料」或「待補資料清單」條列，
請完整列出每一項到 evidence。

Schema：
{
  "radar_dimensions": ["5個維度名稱，如 價格接受度、時程壓力、技術整合度、決策層支持、現場配合度"],
  "radar_opportunity": [5個數字 1-10，代表成交機會強度],
  "radar_resistance": [5個數字 1-10，代表阻力強度],
  "follow_up": {
    "top_data_needed": "最該先補的資料",
    "top_person": "最該先說服的角色",
    "best_timing": "最適合的追蹤時間",
    "next_message": "最合適的下一封訊息草稿"
  },
  "action_items": [
    {"task": "行動項目", "owner": "負責人", "due": "建議時間"}
  ],
  "evidence": [
    {"title": "證據標題", "summary": "一句話摘要", "detail": "完整內容或引用", "source": "資料來源"}
  ]
}

只輸出 JSON，不要任何說明。

---
分析報告：
{analysis}
"""

_WR_CSS = """
:root{--bg:#0a1220;--bg2:#0d1b2a;--card:#111f33;--card2:#16283f;--navy:#0f1b2e;
--steel:#64748b;--white:#f8fafc;--techblue:#38bdf8;--blue:#2563eb;
--orange:#f97316;--green:#22c55e;--red:#ef4444;
--text:#e6edf3;--muted:#93a2b8;--border:rgba(56,189,248,.15)}
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
html{scroll-behavior:smooth}
body{background:var(--bg);color:var(--text);font-family:'Noto Sans TC','Inter',sans-serif;font-size:14px;line-height:1.6}
img,canvas{max-width:100%;height:auto}
.wr-hero{position:relative;padding:56px 24px 40px;overflow:hidden;
background:radial-gradient(circle at 20% 20%,rgba(56,189,248,.18),transparent 45%),
radial-gradient(circle at 80% 0%,rgba(37,99,235,.22),transparent 50%),
linear-gradient(160deg,#0a1220 0%,#0d1b2a 55%,#111f33 100%);
border-bottom:1px solid var(--border)}
.wr-hero::before{content:'';position:absolute;inset:0;opacity:.5;pointer-events:none;
background-image:
  linear-gradient(rgba(148,163,184,.06) 1px,transparent 1px),
  linear-gradient(90deg,rgba(148,163,184,.06) 1px,transparent 1px);
background-size:36px 36px}
.wr-hero-inner{position:relative;max-width:1400px;margin:0 auto}
.wr-tag{display:inline-block;font-size:11px;letter-spacing:2px;color:var(--techblue);
text-transform:uppercase;border:1px solid var(--border);padding:4px 12px;border-radius:20px;margin-bottom:14px}
.wr-title{font-size:clamp(22px,3.4vw,34px);font-weight:800;color:var(--white);margin-bottom:6px}
.wr-sub{font-size:14px;color:var(--muted);margin-bottom:26px}
.wr-hero-grid{display:grid;grid-template-columns:repeat(4,1fr);gap:14px;margin-bottom:22px}
.hero-card{background:rgba(17,31,51,.7);backdrop-filter:blur(6px);border:1px solid var(--border);
border-radius:12px;padding:16px}
.hero-card .lbl{font-size:11px;color:var(--muted);letter-spacing:1px;text-transform:uppercase;margin-bottom:6px}
.hero-card .val{font-size:15px;font-weight:700;color:var(--white);line-height:1.4}
.hero-card.opp .val{color:var(--green)}
.hero-card.risk .val{color:var(--orange)}
.kpi-row{display:grid;grid-template-columns:repeat(4,1fr);gap:14px}
.kpi-card{background:var(--card);border:1px solid var(--border);border-radius:12px;padding:16px;text-align:center}
.kpi-num{font-size:2rem;font-weight:800;color:var(--techblue);line-height:1}
.kpi-lbl{font-size:11px;color:var(--muted);margin-top:6px;letter-spacing:.5px}
.kpi-card.risk .kpi-num{color:var(--orange)}
.kpi-card.opp .kpi-num{color:var(--green)}
.wr-nav{position:sticky;top:0;z-index:100;background:#0a1220;border-bottom:1px solid var(--border);
display:flex;gap:4px;overflow-x:auto;padding:0 16px}
.wr-nav-btn{flex-shrink:0;background:transparent;border:none;color:var(--muted);font-size:13px;
font-weight:600;padding:14px 16px;cursor:pointer;border-bottom:2px solid transparent;white-space:nowrap}
.wr-nav-btn:hover{color:var(--text)}
.wr-nav-btn.active{color:var(--techblue);border-bottom-color:var(--techblue)}
.wr-page{max-width:1400px;margin:0 auto;padding:24px 16px 60px}
.wr-panel{display:none}
.wr-panel.active{display:block}
.sec-title{font-size:16px;font-weight:700;color:var(--white);margin-bottom:14px;padding-bottom:8px;
border-bottom:2px solid var(--techblue);display:inline-block}
.grid2{display:grid;grid-template-columns:1.2fr 1fr;gap:18px}
.grid3{display:grid;grid-template-columns:repeat(3,1fr);gap:14px}
.card{background:var(--card);border:1px solid var(--border);border-radius:12px;padding:18px;margin-bottom:14px}
.brief-box{border:1px solid var(--techblue);background:rgba(56,189,248,.06);border-radius:12px;padding:18px}
.brief-title{font-size:13px;font-weight:700;color:var(--techblue);margin-bottom:10px}
.brief-item{display:flex;gap:8px;font-size:13px;color:var(--text);margin-bottom:8px;line-height:1.5}
.brief-dot{color:var(--techblue);flex-shrink:0}
.profile-row{display:flex;justify-content:space-between;padding:8px 0;border-bottom:1px solid var(--border);font-size:13px}
.profile-row:last-child{border-bottom:none}
.profile-k{color:var(--muted)}
.profile-v{color:var(--text);font-weight:600;text-align:right;max-width:65%}
.pain-card{border-left:3px solid var(--orange);background:var(--card2);border-radius:8px;padding:12px 14px;margin-bottom:10px}
.pain-top{display:flex;justify-content:space-between;align-items:center;margin-bottom:4px}
.pain-title{font-size:13px;font-weight:700;color:var(--text)}
.sev-badge{font-size:10px;padding:2px 9px;border-radius:20px;font-weight:700}
.sev-高{background:rgba(239,68,68,.18);color:var(--red)}
.sev-中{background:rgba(249,115,22,.18);color:var(--orange)}
.sev-低{background:rgba(34,197,94,.18);color:var(--green)}
.pain-desc{font-size:12px;color:var(--muted);line-height:1.5;margin-bottom:4px}
.pain-ref{font-size:11px;color:var(--steel);font-style:italic}
.stake-card{background:var(--card2);border-radius:10px;padding:14px;margin-bottom:10px;border:1px solid var(--border)}
.stake-head{display:flex;align-items:center;gap:10px;margin-bottom:8px}
.avatar{width:42px;height:42px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:15px;font-weight:700;flex-shrink:0;color:#fff}
.av-c{background:linear-gradient(135deg,#22c55e,#15803d)}
.av-b{background:linear-gradient(135deg,#ef4444,#b91c1c)}
.av-i{background:linear-gradient(135deg,#38bdf8,#0369a1)}
.av-d{background:linear-gradient(135deg,#f97316,#c2410c)}
.av-x{background:linear-gradient(135deg,#64748b,#334155)}
.stake-name{font-size:14px;font-weight:700;color:var(--text)}
.stake-ttl{font-size:11px;color:var(--muted)}
.badge{display:inline-block;padding:2px 10px;border-radius:20px;font-size:11px;font-weight:600;margin-bottom:8px}
.bC{background:#22c55e;color:#052e16}.bB{background:#ef4444;color:#fff}
.bI{background:#38bdf8;color:#03202e}.bD{background:#f97316;color:#2b0e00}.bX{background:#64748b;color:#fff}
.stake-meta{display:flex;gap:16px;margin-bottom:6px}
.meta-mini{font-size:11px;color:var(--muted)}
.meta-mini b{color:var(--techblue)}
.stake-notes{font-size:12px;color:var(--text);line-height:1.5}
.filter-row{display:flex;gap:8px;margin-bottom:14px;flex-wrap:wrap}
.filter-btn{font-size:12px;border:1px solid var(--border);background:var(--card2);color:var(--muted);
padding:6px 14px;border-radius:20px;cursor:pointer;font-weight:600}
.filter-btn.active{background:var(--techblue);color:#03202e;border-color:var(--techblue)}
.risk-card{background:var(--card2);border-radius:10px;padding:14px;margin-bottom:10px;border-left:3px solid var(--steel)}
.risk-card[data-level="高"]{border-left-color:var(--red)}
.risk-card[data-level="中"]{border-left-color:var(--orange)}
.risk-card[data-level="低"]{border-left-color:var(--green)}
.risk-top{display:flex;justify-content:space-between;align-items:center;margin-bottom:6px}
.risk-name{font-size:13px;font-weight:700;color:var(--text)}
.risk-cat{font-size:10px;color:var(--muted);margin-bottom:6px}
.risk-mit{font-size:12px;color:var(--text);line-height:1.5;margin-bottom:4px}
.risk-ref{font-size:11px;color:var(--steel);font-style:italic}
.sim-card{background:var(--card);border:1px solid var(--border);border-radius:12px;padding:20px}
.strat-tabs,.roi-tabs{display:flex;gap:8px;margin-bottom:16px;flex-wrap:wrap}
.strat-btn,.roi-btn{flex:1;min-width:140px;background:var(--card2);border:1px solid var(--border);
color:var(--text);padding:12px 10px;border-radius:10px;cursor:pointer;font-size:13px;font-weight:700;text-align:center}
.strat-btn.active,.roi-btn.active{background:linear-gradient(135deg,var(--techblue),var(--blue));color:#03202e;border-color:var(--techblue)}
.strat-panel,.roi-panel{display:none}
.strat-panel.active,.roi-panel.active{display:block}
.strat-row{display:flex;justify-content:space-between;padding:9px 0;border-bottom:1px solid var(--border);font-size:13px;gap:16px}
.strat-row:last-child{border-bottom:none}
.strat-k{color:var(--muted);flex-shrink:0;width:110px}
.strat-v{color:var(--text);line-height:1.5}
.strat-tag{display:inline-block;font-size:10px;color:var(--steel);margin-top:2px}
.est-note{font-size:11px;color:var(--steel);font-style:italic;margin-top:10px}
.deal-pct{font-size:2.2rem;font-weight:800;color:var(--green)}
.risk-sim-grid{display:grid;grid-template-columns:1fr 1fr;gap:20px;margin-top:6px}
.slider-item{margin-bottom:16px}
.slider-top{display:flex;justify-content:space-between;font-size:12px;margin-bottom:6px}
.slider-top b{color:var(--techblue)}
input[type=range]{width:100%;accent-color:var(--techblue)}
.risk-out-card{background:var(--card2);border-radius:10px;padding:14px;margin-bottom:10px}
.risk-out-lbl{font-size:11px;color:var(--muted);margin-bottom:4px}
.risk-out-val{font-size:20px;font-weight:800}
.risk-out-val.ok{color:var(--green)}.risk-out-val.mid{color:var(--orange)}.risk-out-val.bad{color:var(--red)}
.need-list{font-size:12px;color:var(--text);line-height:1.7;margin-top:6px}
.need-list .need-item{display:none;padding:4px 8px;background:rgba(249,115,22,.08);border-radius:6px;margin-bottom:4px}
.need-list .need-item.show{display:block}
.order-list{font-size:12px;color:var(--text);line-height:2}
.order-list li{margin-left:18px}
.roi-grid{display:grid;grid-template-columns:repeat(4,1fr);gap:12px;margin-bottom:14px}
.roi-num-card{background:var(--card2);border-radius:10px;padding:14px;text-align:center}
.roi-num{font-size:1.4rem;font-weight:800;color:var(--green)}
.roi-num.pay{color:var(--techblue)}
.roi-lbl{font-size:11px;color:var(--muted);margin-top:4px}
.roi-confirm{font-size:12px;color:var(--muted);line-height:1.6}
.roi-scale-row{display:flex;align-items:center;gap:12px;margin:14px 0;padding:12px;background:var(--card2);border-radius:10px}
.roi-scale-lbl{font-size:12px;color:var(--muted);white-space:nowrap}
.roi-scale-val{font-size:13px;font-weight:700;color:var(--techblue);white-space:nowrap}
.follow-grid{display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-bottom:14px}
.follow-card{background:var(--card2);border-radius:10px;padding:14px}
.follow-lbl{font-size:11px;color:var(--muted);margin-bottom:6px}
.follow-val{font-size:13px;font-weight:700;color:var(--text);line-height:1.5}
.recalc-btn{background:var(--techblue);color:#03202e;border:none;padding:8px 16px;border-radius:8px;
font-size:12px;font-weight:700;cursor:pointer;margin-bottom:14px}
.msg-box{background:var(--card2);border-left:3px solid var(--techblue);border-radius:8px;padding:12px 14px;font-size:13px;line-height:1.6}
.action-item{display:flex;align-items:center;gap:10px;padding:10px 0;border-bottom:1px solid var(--border)}
.action-item:last-child{border-bottom:none}
.action-item input{width:16px;height:16px;accent-color:var(--techblue)}
.action-task{font-size:13px;color:var(--text);flex:1}
.action-task.done{text-decoration:line-through;color:var(--steel)}
.action-meta{font-size:11px;color:var(--muted)}
.ev-card{background:var(--card2);border-radius:10px;margin-bottom:10px;border:1px solid var(--border);overflow:hidden;cursor:pointer}
.ev-head{display:flex;justify-content:space-between;align-items:center;padding:13px 16px}
.ev-title{font-size:13px;font-weight:700;color:var(--text)}
.ev-summary{font-size:12px;color:var(--muted);margin-top:2px}
.ev-arrow{color:var(--techblue);font-size:14px;transition:transform .2s;flex-shrink:0;margin-left:10px}
.ev-card.open .ev-arrow{transform:rotate(180deg)}
.ev-detail{max-height:0;overflow:hidden;transition:max-height .25s ease;padding:0 16px}
.ev-card.open .ev-detail{max-height:400px;padding:0 16px 14px}
.ev-detail-txt{font-size:12px;color:var(--text);line-height:1.6;border-top:1px solid var(--border);padding-top:10px}
.ev-source{font-size:11px;color:var(--steel);font-style:italic;margin-top:6px}
@media(max-width:1100px){
  .wr-hero-grid,.kpi-row{grid-template-columns:repeat(2,1fr)}
  .grid2,.grid3,.risk-sim-grid,.roi-grid,.follow-grid{grid-template-columns:1fr}
}
@media(max-width:600px){
  .wr-hero{padding:36px 14px 28px}
  .wr-page{padding:16px 10px 40px}
  .roi-grid{grid-template-columns:1fr 1fr}
}
"""


def _wr_role_cls(role_type: str):
    r = (role_type or "").lower()
    if "champion" in r or "倡導" in r:
        return "av-c", "bC"
    if any(x in r for x in ("blocker", "gatekeeper", "阻礙", "把關")):
        return "av-b", "bB"
    if "influencer" in r or "影響" in r:
        return "av-i", "bI"
    if "decision" in r or "決策" in r:
        return "av-d", "bD"
    return "av-x", "bX"


def _wr_render_exec_summary(items: list) -> str:
    items = items or ["（尚無足夠資料產出主管簡報重點）"]
    return "".join(
        f'<div class="brief-item"><span class="brief-dot">▸</span><span>{i}</span></div>'
        for i in items
    )


def _wr_render_profile(p: dict) -> str:
    p = p or {}
    rows = [
        ("產業別", p.get("industry", "—")),
        ("公司規模", p.get("company_size", "—")),
        ("目前現況", p.get("current_situation", "—")),
        ("背景說明", p.get("background", "—")),
    ]
    return "".join(
        f'<div class="profile-row"><span class="profile-k">{k}</span><span class="profile-v">{v}</span></div>'
        for k, v in rows
    )


def _wr_render_pain_points(items: list) -> str:
    if not items:
        return '<p style="color:var(--muted);font-size:12px">（未識別到明確痛點）</p>'
    out = []
    for p in items:
        sev = p.get("severity", "中")
        out.append(
            f'<div class="pain-card"><div class="pain-top">'
            f'<span class="pain-title">{p.get("title","")}</span>'
            f'<span class="sev-badge sev-{sev}">{sev}風險</span></div>'
            f'<div class="pain-desc">{p.get("description","")}</div>'
            f'<div class="pain-ref">依據：{p.get("evidence_ref","—")}</div></div>'
        )
    return "".join(out)


def _wr_render_decision_chain(items: list) -> str:
    if not items:
        return '<div class="stake-card"><p style="color:var(--muted)">（未識別到明確決策鏈成員）</p></div>'
    out = []
    for s in items:
        av, bv = _wr_role_cls(s.get("role_type", ""))
        out.append(
            f'<div class="stake-card"><div class="stake-head">'
            f'<div class="avatar {av}">{(s.get("name") or "?")[:1]}</div>'
            f'<div><div class="stake-name">{s.get("name","")}</div>'
            f'<div class="stake-ttl">{s.get("title","")}</div></div></div>'
            f'<span class="badge {bv}">{s.get("role_label", s.get("role_type",""))}</span>'
            f'<div class="stake-meta">'
            f'<span class="meta-mini">影響力 <b>{s.get("influence",5)}/10</b></span>'
            f'<span class="meta-mini">支持度 <b>{s.get("support_level",5)}/10</b></span></div>'
            f'<div class="stake-notes">{s.get("notes","")}</div></div>'
        )
    return "".join(out)


def _wr_render_risk_map(items: list) -> str:
    if not items:
        return '<p style="color:var(--muted);font-size:12px">（未識別到明確風險）</p>'
    out = []
    for r in items:
        lvl = r.get("level", "中")
        out.append(
            f'<div class="risk-card" data-level="{lvl}">'
            f'<div class="risk-top"><span class="risk-name">{r.get("risk","")}</span>'
            f'<span class="sev-badge sev-{lvl}">{lvl}</span></div>'
            f'<div class="risk-cat">類別：{r.get("category","—")}</div>'
            f'<div class="risk-mit">緩解措施：{r.get("mitigation","—")}</div>'
            f'<div class="risk-ref">依據：{r.get("evidence_ref","—")}</div></div>'
        )
    return "".join(out)


def _wr_render_strategies(modes: dict) -> str:
    modes = modes or {}
    order = [("conservative", "保守試點型"), ("fast", "快速導入型"), ("executive", "主管決策型")]
    btns, panels = [], []
    for i, (key, default_title) in enumerate(order):
        m = modes.get(key, {})
        active = " active" if i == 0 else ""
        btns.append(
            f'<button class="strat-btn{active}" data-strat="{key}">{m.get("title", default_title)}</button>'
        )
        panels.append(
            f'<div class="strat-panel{active}" data-strat-panel="{key}">'
            f'<div class="deal-pct">{m.get("deal_chance", 50)}%</div>'
            f'<div class="strat-tag">預估成交機會（策略推估）</div>'
            f'<div class="strat-row"><span class="strat-k">適合情境</span><span class="strat-v">{m.get("fit_scenario","—")}</span></div>'
            f'<div class="strat-row"><span class="strat-k">主要風險</span><span class="strat-v">{m.get("main_risk","—")}</span></div>'
            f'<div class="strat-row"><span class="strat-k">建議話術</span><span class="strat-v">{m.get("talk_track","—")}</span></div>'
            f'<div class="strat-row"><span class="strat-k">下一步行動</span><span class="strat-v">{m.get("next_action","—")}</span></div>'
            f'<div class="est-note">＊以上為策略推估，依現有資料與風險判斷估算</div></div>'
        )
    return "".join(btns), "".join(panels)


def _wr_render_risk_factors(factors: list) -> str:
    if not factors:
        factors = [
            {"id": "price", "label": "價格敏感度", "default": 50, "inverse": False, "data_needed_if_high": "需補充：客戶預算上限與比價對象"},
            {"id": "downtime", "label": "停線風險", "default": 40, "inverse": False, "data_needed_if_high": "需補充：可容忍停機時間"},
            {"id": "complexity", "label": "導入麻煩程度", "default": 45, "inverse": False, "data_needed_if_high": "需補充：現有系統整合細節"},
            {"id": "support", "label": "決策者支持度", "default": 55, "inverse": True, "data_needed_if_high": "需補充：決策者關切重點"},
            {"id": "resistance", "label": "現場人員排斥程度", "default": 40, "inverse": False, "data_needed_if_high": "需補充：現場人員教育訓練規劃"},
        ]
    sliders, needs = [], []
    for f in factors:
        sliders.append(
            f'<div class="slider-item">'
            f'<div class="slider-top"><span>{f.get("label","")}</span><b class="slider-val" data-val-for="{f.get("id","")}">{f.get("default",50)}</b></div>'
            f'<input type="range" min="0" max="100" value="{f.get("default",50)}" '
            f'class="risk-slider" data-id="{f.get("id","")}" data-inverse="{"1" if f.get("inverse") else "0"}" '
            f'data-label="{f.get("label","")}"></div>'
        )
        needs.append(
            f'<div class="need-item" id="need-{f.get("id","")}">⚠ {f.get("data_needed_if_high","")}</div>'
        )
    return "".join(sliders), "".join(needs)


def _wr_render_roi(scenarios: dict) -> str:
    scenarios = scenarios or {}
    order = [("conservative", "保守效益"), ("normal", "一般效益"), ("aggressive", "積極效益")]
    btns, panels = [], []
    for i, (key, default_title) in enumerate(order):
        s = scenarios.get(key, {})
        active = " active" if i == 0 else ""
        btns.append(f'<button class="roi-btn{active}" data-roi="{key}">{default_title}</button>')
        confirm = "".join(f"<li>{c}</li>" for c in s.get("need_confirm", [])) or "<li>—</li>"
        panels.append(
            f'<div class="roi-panel{active}" data-roi-panel="{key}">'
            f'<div class="roi-grid">'
            f'<div class="roi-num-card"><div class="roi-num roi-cost" data-base="{s.get("cost_saving",0)}">{s.get("cost_saving",0)}萬</div><div class="roi-lbl">可能節省成本／年</div></div>'
            f'<div class="roi-num-card"><div class="roi-num roi-quality" data-base="{s.get("quality_value",0)}">{s.get("quality_value",0)}萬</div><div class="roi-lbl">品質改善價值</div></div>'
            f'<div class="roi-num-card"><div class="roi-num roi-complaint" data-base="{s.get("complaint_value",0)}">{s.get("complaint_value",0)}萬</div><div class="roi-lbl">客訴下降價值</div></div>'
            f'<div class="roi-num-card"><div class="roi-num pay roi-payback" data-base="{s.get("payback_months",12)}">{s.get("payback_months",12)}個月</div><div class="roi-lbl">回收期估算</div></div>'
            f'</div><div class="roi-confirm"><b style="color:var(--muted)">需要確認的數字：</b><ul style="margin-left:16px">{confirm}</ul></div>'
            f'<div class="est-note">＊以上為策略推估，實際效益依客戶場域驗證為準</div></div>'
        )
    return "".join(btns), "".join(panels)


def _wr_render_action_items(items: list) -> str:
    if not items:
        return '<p style="color:var(--muted);font-size:12px">（尚無行動項目）</p>'
    out = []
    for a in items:
        out.append(
            f'<div class="action-item"><input type="checkbox" class="action-check">'
            f'<span class="action-task">{a.get("task","")}</span>'
            f'<span class="action-meta">{a.get("owner","")}｜{a.get("due","")}</span></div>'
        )
    return "".join(out)


def _wr_render_evidence(items: list) -> str:
    if not items:
        return '<p style="color:var(--muted);font-size:12px">（尚無證據資料）</p>'
    out = []
    for e in items:
        out.append(
            f'<div class="ev-card"><div class="ev-head">'
            f'<div><div class="ev-title">{e.get("title","")}</div>'
            f'<div class="ev-summary">{e.get("summary","")}</div></div>'
            f'<span class="ev-arrow">▾</span></div>'
            f'<div class="ev-detail"><div class="ev-detail-txt">{e.get("detail","")}'
            f'<div class="ev-source">來源：{e.get("source","—")}</div></div></div></div>'
        )
    return "".join(out)


def _extract_wr_dashboard_json_stream(analysis: str):
    """Generator: yields (data_so_far, status_message) at each extraction stage.

    Structured as a generator (rather than a plain function with an on_step callback) so the
    Gradio-facing caller can yield a status update to a visible Markdown component between each
    OpenAI call — a callback fired from inside a blocking call can't make an outer generator
    pause and yield, so the calls themselves have to live in the generator.

    NOTE: use str.replace() rather than str.format() — the prompts above embed a literal
    JSON schema full of unescaped `{`/`}`, which str.format() parses as replacement fields
    and always raises KeyError before the API is ever called.
    """
    import json as _j
    text = analysis[:30000]
    data: dict = {}
    yield data, "分析戰情總覽、客戶輪廓、痛點、決策鏈、風險與策略..."
    try:
        resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": _WR_EXTRACT_SYSTEM},
                {"role": "user", "content": _WR_EXTRACT_PROMPT_CORE.replace("{analysis}", text)},
            ],
            response_format={"type": "json_object"},
            max_tokens=3000,
            temperature=0,
        )
        data.update(_j.loads(resp.choices[0].message.content))
    except Exception:
        pass
    yield data, "分析成交機會雷達、追蹤行動與證據資料..."
    try:
        resp = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": _WR_EXTRACT_SYSTEM},
                {"role": "user", "content": _WR_EXTRACT_PROMPT_TAIL.replace("{analysis}", text)},
            ],
            response_format={"type": "json_object"},
            max_tokens=2500,
            temperature=0,
        )
        data.update(_j.loads(resp.choices[0].message.content))
    except Exception:
        pass
    yield data, "整理資料中..."


def _extract_wr_dashboard_json(analysis: str) -> dict:
    data: dict = {}
    for data, _status in _extract_wr_dashboard_json_stream(analysis):
        pass
    return data


def _render_wr_dashboard_html(data: dict) -> str:
    import json as _j
    hero = data.get("hero", {})
    kpi = data.get("kpi", {})
    client_name = hero.get("client_name", "目標客戶")
    industry_title = hero.get("industry_title", "業務提案戰情室")
    decision_stage = hero.get("decision_stage", "評估中")
    data_source = data.get("data_source", "業務分析資料")
    est_value = data.get("est_value", "Medium（一般客戶）")

    dims = data.get("radar_dimensions", ["價格接受度", "時程壓力", "技術整合度", "決策層支持", "現場配合度"])
    r_opp = data.get("radar_opportunity", [6, 6, 6, 6, 6])
    r_res = data.get("radar_resistance", [5, 5, 5, 5, 5])
    n = max(len(dims), len(r_opp), len(r_res))
    dims = (dims + ["—"] * n)[:n]
    r_opp = (r_opp + [5] * n)[:n]
    r_res = (r_res + [5] * n)[:n]

    exec_html = _wr_render_exec_summary(data.get("exec_summary", []))
    profile_html = _wr_render_profile(data.get("customer_profile", {}))
    pain_html = _wr_render_pain_points(data.get("pain_points", []))
    chain_html = _wr_render_decision_chain(data.get("decision_chain", []))
    risk_html = _wr_render_risk_map(data.get("risk_map", []))
    strat_btns, strat_panels = _wr_render_strategies(data.get("strategy_modes", {}))
    slider_html, need_html = _wr_render_risk_factors(data.get("risk_factors", []))
    roi_btns, roi_panels = _wr_render_roi(data.get("roi_scenarios", {}))
    action_html = _wr_render_action_items(data.get("action_items", []))
    evidence_html = _wr_render_evidence(data.get("evidence", []))

    fu = data.get("follow_up", {})

    tabs = [
        ("overview", "戰情總覽"), ("profile", "客戶輪廓"), ("pain", "需求痛點"),
        ("chain", "決策鏈"), ("risk", "風險地圖"), ("strategy", "策略模擬"),
        ("roi", "ROI情境"), ("action", "行動追蹤"), ("evidence", "證據資料"),
    ]
    nav_html = "".join(
        f'<button class="wr-nav-btn{" active" if i==0 else ""}" data-tab="{k}">{label}</button>'
        for i, (k, label) in enumerate(tabs)
    )

    parts = [
        "<!DOCTYPE html>", "<html lang='zh-TW'>", "<head>",
        '<meta charset="UTF-8">',
        '<meta name="viewport" content="width=device-width, initial-scale=1.0">',
        "<title>業務提案戰情室</title>",
        '<link href="https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@300;400;500;700'
        '&family=Inter:wght@300;400;500;700&display=swap" rel="stylesheet">',
        '<script src="https://cdn.jsdelivr.net/npm/chart.js"></script>',
        f"<style>{_WR_CSS}</style>", "</head>", "<body>",
        # ── Hero ──
        '<section class="wr-hero"><div class="wr-hero-inner">',
        '<span class="wr-tag">War Room ｜ 業務提案戰情室</span>',
        f'<div class="wr-title">{industry_title}｜{client_name}</div>',
        f'<div class="wr-sub">目前決策階段：{decision_stage}　｜　基於「{data_source}」之戰情分析　｜　案子評級：{est_value}</div>',
        '<div class="wr-hero-grid">',
        f'<div class="hero-card opp"><div class="lbl">主要成交機會</div><div class="val">{hero.get("main_opportunity","—")}</div></div>',
        f'<div class="hero-card risk"><div class="lbl">最大風險</div><div class="val">{hero.get("top_risk","—")}</div></div>',
        f'<div class="hero-card"><div class="lbl">下一步建議</div><div class="val">{hero.get("next_step","—")}</div></div>',
        f'<div class="hero-card"><div class="lbl">決策階段</div><div class="val">{decision_stage}</div></div>',
        '</div>',
        '<div class="kpi-row">',
        f'<div class="kpi-card opp"><div class="kpi-num">{kpi.get("deal_score",50)}</div><div class="kpi-lbl">成交機會分數</div></div>',
        f'<div class="kpi-card risk"><div class="kpi-num">{kpi.get("risk_score",50)}</div><div class="kpi-lbl">風險等級分數</div></div>',
        f'<div class="kpi-card"><div class="kpi-num">{kpi.get("readiness_score",50)}</div><div class="kpi-lbl">導入準備度</div></div>',
        f'<div class="kpi-card"><div class="kpi-num">{kpi.get("roi_potential",50)}</div><div class="kpi-lbl">ROI 潛力分數</div></div>',
        '</div></div></section>',
        # ── Tabs nav ──
        f'<nav class="wr-nav">{nav_html}</nav>',
        '<div class="wr-page">',
        # 戰情總覽
        '<div class="wr-panel active" data-panel="overview"><div class="grid2">',
        '<div><div class="sec-title">成交機會雷達</div><div class="card">'
        '<div style="position:relative;height:280px"><canvas id="radarChart"></canvas></div></div></div>',
        f'<div><div class="sec-title">主管簡報重點</div><div class="brief-box">{exec_html}</div></div>',
        '</div></div>',
        # 客戶輪廓
        f'<div class="wr-panel" data-panel="profile"><div class="sec-title">客戶輪廓</div><div class="card">{profile_html}</div></div>',
        # 需求痛點
        f'<div class="wr-panel" data-panel="pain"><div class="sec-title">需求痛點</div>{pain_html}</div>',
        # 決策鏈
        f'<div class="wr-panel" data-panel="chain"><div class="sec-title">決策鏈分析</div>{chain_html}</div>',
        # 風險地圖
        '<div class="wr-panel" data-panel="risk">',
        '<div class="sec-title">風險與阻力地圖</div>',
        '<div class="filter-row">'
        '<button class="filter-btn active" data-filter="全部">全部</button>'
        '<button class="filter-btn" data-filter="高">高風險</button>'
        '<button class="filter-btn" data-filter="中">中風險</button>'
        '<button class="filter-btn" data-filter="低">低風險</button></div>',
        f'<div id="risk-cards">{risk_html}</div>',
        '<div class="sec-title" style="margin-top:22px">風險調整模擬</div>',
        '<div class="sim-card"><div class="risk-sim-grid">',
        f'<div>{slider_html}</div>',
        '<div>'
        '<div class="risk-out-card"><div class="risk-out-lbl">風險等級變化（策略推估）</div>'
        '<div class="risk-out-val" id="risk-out-level">中</div></div>'
        '<div class="risk-out-card"><div class="risk-out-lbl">成交阻力變化（策略推估）</div>'
        '<div class="risk-out-val" id="risk-out-resist">中</div></div>'
        '<div class="risk-out-card"><div class="risk-out-lbl">需要補強的資料</div>'
        f'<div class="need-list" id="need-list">{need_html}</div></div>'
        '<div class="risk-out-card"><div class="risk-out-lbl">建議處理順序</div>'
        '<ol class="order-list" id="order-list"></ol></div>'
        '</div></div></div></div>',
        # 策略模擬
        '<div class="wr-panel" data-panel="strategy">',
        '<div class="sec-title">提案策略模擬</div>',
        '<div class="sim-card">',
        f'<div class="strat-tabs">{strat_btns}</div>',
        f'{strat_panels}',
        '</div></div>',
        # ROI情境
        '<div class="wr-panel" data-panel="roi">',
        '<div class="sec-title">ROI 與導入效益估算</div>',
        '<div class="sim-card">',
        f'<div class="roi-tabs">{roi_btns}</div>',
        '<div class="roi-scale-row"><span class="roi-scale-lbl">導入規模調整</span>'
        '<input type="range" id="roi-scale" min="50" max="150" value="100" style="flex:1">'
        '<span class="roi-scale-val" id="roi-scale-val">100%</span></div>',
        f'{roi_panels}',
        '</div></div>',
        # 行動追蹤
        '<div class="wr-panel" data-panel="action">',
        '<div class="sec-title">追蹤優先順序模擬</div>',
        '<button class="recalc-btn" id="recalc-btn">🔄 依目前風險重新試算</button>',
        '<div class="follow-grid">',
        f'<div class="follow-card"><div class="follow-lbl">最該先補的資料</div><div class="follow-val" id="fu-data">{fu.get("top_data_needed","—")}</div></div>',
        f'<div class="follow-card"><div class="follow-lbl">最該先說服的角色</div><div class="follow-val" id="fu-person">{fu.get("top_person","—")}</div></div>',
        f'<div class="follow-card"><div class="follow-lbl">最適合的追蹤時間</div><div class="follow-val">{fu.get("best_timing","—")}</div></div>',
        '<div class="follow-card"><div class="follow-lbl">最合適的下一封訊息</div>'
        f'<div class="msg-box">{fu.get("next_message","—")}</div></div>',
        '</div>',
        '<div class="sec-title" style="margin-top:10px">行動項目清單</div>',
        f'<div class="card">{action_html}</div>',
        '</div>',
        # 證據資料
        f'<div class="wr-panel" data-panel="evidence"><div class="sec-title">資料證據區</div>{evidence_html}</div>',
        '</div>',  # end wr-page
        # ── JS ──
        "<script>",
        # tab switching
        "document.querySelectorAll('.wr-nav-btn').forEach(function(btn){",
        "btn.addEventListener('click',function(){",
        "document.querySelectorAll('.wr-nav-btn').forEach(function(b){b.classList.remove('active')});",
        "document.querySelectorAll('.wr-panel').forEach(function(p){p.classList.remove('active')});",
        "btn.classList.add('active');",
        "document.querySelector('.wr-panel[data-panel=\"'+btn.dataset.tab+'\"]').classList.add('active');",
        "window.scrollTo({top:0,behavior:'smooth'});});});",
        # risk level filter
        "document.querySelectorAll('.filter-btn').forEach(function(btn){",
        "btn.addEventListener('click',function(){",
        "document.querySelectorAll('.filter-btn').forEach(function(b){b.classList.remove('active')});",
        "btn.classList.add('active');",
        "var f=btn.dataset.filter;",
        "document.querySelectorAll('#risk-cards .risk-card').forEach(function(c){",
        "c.style.display=(f==='全部'||c.dataset.level===f)?'':'none';});});});",
        # evidence expand
        "document.querySelectorAll('.ev-card').forEach(function(c){",
        "c.querySelector('.ev-head').addEventListener('click',function(){c.classList.toggle('open');});});",
        # strategy switch
        "document.querySelectorAll('.strat-btn').forEach(function(btn){",
        "btn.addEventListener('click',function(){",
        "document.querySelectorAll('.strat-btn').forEach(function(b){b.classList.remove('active')});",
        "document.querySelectorAll('.strat-panel').forEach(function(p){p.classList.remove('active')});",
        "btn.classList.add('active');",
        "document.querySelector('.strat-panel[data-strat-panel=\"'+btn.dataset.strat+'\"]').classList.add('active');});});",
        # roi switch + scale
        "document.querySelectorAll('.roi-btn').forEach(function(btn){",
        "btn.addEventListener('click',function(){",
        "document.querySelectorAll('.roi-btn').forEach(function(b){b.classList.remove('active')});",
        "document.querySelectorAll('.roi-panel').forEach(function(p){p.classList.remove('active')});",
        "btn.classList.add('active');",
        "document.querySelector('.roi-panel[data-roi-panel=\"'+btn.dataset.roi+'\"]').classList.add('active');});});",
        "function fmtWan(v){return (Math.round(v*10)/10)+'萬';}",
        "function applyRoiScale(){",
        "var scale=parseInt(document.getElementById('roi-scale').value)/100;",
        "document.getElementById('roi-scale-val').textContent=Math.round(scale*100)+'%';",
        "document.querySelectorAll('.roi-cost,.roi-quality,.roi-complaint').forEach(function(el){",
        "var base=parseFloat(el.dataset.base||0);el.textContent=fmtWan(base*scale);});",
        "document.querySelectorAll('.roi-payback').forEach(function(el){",
        "var base=parseFloat(el.dataset.base||12);",
        "var val=scale>0?(base/scale):base;el.textContent=(Math.round(val*10)/10)+'個月';});}",
        "document.getElementById('roi-scale').addEventListener('input',applyRoiScale);",
        # risk sliders simulation
        "function recomputeRisk(){",
        "var total=0,n=0,resistTotal=0,resistN=0,contribs=[];",
        "document.querySelectorAll('.risk-slider').forEach(function(s){",
        "var v=parseInt(s.value);var inv=s.dataset.inverse==='1';",
        "var contrib=inv?(100-v):v;",
        "document.querySelector('.slider-val[data-val-for=\"'+s.dataset.id+'\"]').textContent=v;",
        "total+=contrib;n++;",
        "if(['price','complexity','resistance'].indexOf(s.dataset.id)>-1){resistTotal+=contrib;resistN++;}",
        "contribs.push({id:s.dataset.id,label:s.dataset.label,contrib:contrib});",
        "var need=document.getElementById('need-'+s.dataset.id);",
        "if(need){need.classList.toggle('show',contrib>=60);}});",
        "var riskAvg=n?Math.round(total/n):50;",
        "var resistAvg=resistN?Math.round(resistTotal/resistN):riskAvg;",
        "function lvl(v){return v>=67?'高':(v>=34?'中':'低');}",
        "function cls(v){return v>=67?'bad':(v>=34?'mid':'ok');}",
        "var lv=document.getElementById('risk-out-level');",
        "lv.textContent=lvl(riskAvg)+'（'+riskAvg+'分）';lv.className='risk-out-val '+cls(riskAvg);",
        "var rv=document.getElementById('risk-out-resist');",
        "rv.textContent=lvl(resistAvg)+'（'+resistAvg+'分）';rv.className='risk-out-val '+cls(resistAvg);",
        "contribs.sort(function(a,b){return b.contrib-a.contrib;});",
        "var ol=document.getElementById('order-list');",
        "ol.innerHTML=contribs.map(function(c){return '<li>'+c.label+'（'+c.contrib+'分）</li>';}).join('');",
        "return contribs;}",
        "document.querySelectorAll('.risk-slider').forEach(function(s){s.addEventListener('input',recomputeRisk);});",
        # follow-up recalc
        "document.getElementById('recalc-btn').addEventListener('click',function(){",
        "var contribs=recomputeRisk();",
        "if(contribs&&contribs.length){document.getElementById('fu-data').textContent=",
        "document.getElementById('need-'+contribs[0].id)?document.getElementById('need-'+contribs[0].id).textContent.replace('⚠ ',''):contribs[0].label;}",
        f"var chain={_j.dumps(data.get('decision_chain', []), ensure_ascii=False)};",
        "if(chain&&chain.length){",
        "chain.sort(function(a,b){return (a.support_level||5)-(b.support_level||5)||(b.influence||5)-(a.influence||5);});",
        "document.getElementById('fu-person').textContent=chain[0].name+'／'+chain[0].title;}});",
        # action item checkbox toggle
        "document.querySelectorAll('.action-check').forEach(function(cb){",
        "cb.addEventListener('change',function(){",
        "cb.closest('.action-item').querySelector('.action-task').classList.toggle('done',cb.checked);});});",
        "recomputeRisk();applyRoiScale();",
        # radar chart
        f"var dims={_j.dumps(dims, ensure_ascii=False)};",
        f"var opp={_j.dumps(r_opp)};",
        f"var res={_j.dumps(r_res)};",
        "new Chart(document.getElementById('radarChart'),{type:'radar',",
        "data:{labels:dims,datasets:[",
        "{label:'成交機會強度',data:opp,borderColor:'#22c55e',",
        "backgroundColor:'rgba(34,197,94,.15)',pointBackgroundColor:'#22c55e',borderWidth:2},",
        "{label:'阻力強度',data:res,borderColor:'#f97316',",
        "backgroundColor:'rgba(249,115,22,.1)',pointBackgroundColor:'#f97316',borderWidth:2}]},",
        "options:{responsive:true,maintainAspectRatio:false,",
        "plugins:{legend:{labels:{color:'#93a2b8',font:{size:11}}}},",
        "scales:{r:{angleLines:{color:'rgba(255,255,255,.1)'},grid:{color:'rgba(255,255,255,.1)'},",
        "pointLabels:{color:'#e6edf3',font:{size:11}},ticks:{display:false},min:0,max:10}}}});",
        "</script>",
        "</body></html>",
    ]
    return "\n".join(parts)


def generate_data_dashboard_html(analysis_content: str):
    if not analysis_content or "逐章節分析中" in analysis_content or analysis_content.startswith("*上傳"):
        return None
    data = _extract_wr_dashboard_json(analysis_content)
    html_content = _render_wr_dashboard_html(data)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode="w", encoding="utf-8")
    tmp.write(html_content)
    tmp.close()
    return tmp.name


def docx_to_wr_dashboard_html_stream(file_obj):
    """Generator: yields (html_path_or_None, status_message) while converting a .docx to an
    interactive 業務提案戰情室 HTML dashboard. html_path is only non-None on the final yield."""
    if file_obj is None:
        yield None, "⚠️ 請先上傳 Word 檔案。"
        return
    yield None, "讀取 Word 文件..."
    from docx import Document as _DocxDocument
    doc = _DocxDocument(file_obj.name)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    text = "\n".join(paragraphs)
    if not text.strip():
        yield None, "⚠️ 無法從 Word 檔案讀取到文字內容，請確認檔案格式。"
        return
    data = {}
    for data, status in _extract_wr_dashboard_json_stream(text):
        yield None, status
    yield None, "產生互動式 HTML 戰情室儀表板..."
    html_content = _render_wr_dashboard_html(data)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode="w", encoding="utf-8")
    tmp.write(html_content)
    tmp.close()
    yield tmp.name, "✅ 完成！"


def docx_to_wr_dashboard_html(file_obj):
    path = None
    for path, _status in docx_to_wr_dashboard_html_stream(file_obj):
        pass
    return path


def convert_word_to_dashboard_tab5(file_obj):
    for path, status in docx_to_wr_dashboard_html_stream(file_obj):
        if path:
            yield gr.update(value=path, visible=True), status
        elif status.startswith("⚠️"):
            yield gr.update(visible=False), status
        else:
            yield gr.update(visible=False), f"🔄 {status}"


# ══════════════════════════════════════════════════════════════════════════════
# Tab 6: 數據視覺化大師Agent — CSV/Excel 上傳 → ECharts 風格多維度圖表 (JPG)
# ══════════════════════════════════════════════════════════════════════════════

DATAVIZ_SYSTEM_PROMPT = (
    "請扮演大數據大師、參考 https://echarts.apache.org/examples/zh/index.html 的 "
    "Echarts風格、根據數據的特性，給我幾個圖表呈現的建議、要有高級感。\n"
    "輸出語言：繁體中文。"
)

DATAVIZ_USER_INSTRUCTION = (
    "再製作幾張多維度在同一張且有意義的圖。接著輸出這幾張圖、繁體中文標示、"
    "要有完整圖例跟座標標示等、jpg檔讓我下載。"
)

# ── 資料讀取與輪廓分析 ──────────────────────────────────────────────────────────

def dataviz_load_file(file_obj):
    if file_obj is None:
        return None, gr.update(choices=[], value=None, visible=False), None, "*請上傳 CSV 或 Excel 檔案。*"
    path = file_obj.name if hasattr(file_obj, "name") else str(file_obj)
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".csv":
            df = pd.read_csv(path)
            sheets = {"CSV資料": df}
        else:
            xl = pd.ExcelFile(path)
            sheets = {s: xl.parse(s) for s in xl.sheet_names}
    except Exception as e:
        return None, gr.update(choices=[], value=None, visible=False), None, f"⚠️ 檔案讀取失敗：{e}"
    first_sheet = list(sheets.keys())[0]
    df = sheets[first_sheet]
    status = f"✅ 已載入「{first_sheet}」：{len(df)} 列 × {len(df.columns)} 欄"
    state = {"sheets": sheets, "current": first_sheet}
    return (
        state,
        gr.update(choices=list(sheets.keys()), value=first_sheet, visible=len(sheets) > 1),
        df.head(10),
        status,
    )


def dataviz_switch_sheet(state, sheet_name):
    if not state or sheet_name not in state.get("sheets", {}):
        return state, None, "*請選擇有效的工作表。*"
    state = dict(state)
    state["current"] = sheet_name
    df = state["sheets"][sheet_name]
    status = f"✅ 已載入「{sheet_name}」：{len(df)} 列 × {len(df.columns)} 欄"
    return state, df.head(10), status


def build_data_profile(df: pd.DataFrame, max_cols: int = 40) -> str:
    lines = [f"資料筆數：{len(df)} 列 × {len(df.columns)} 欄"]
    lines.append("\n欄位資訊：")
    for col in list(df.columns)[:max_cols]:
        s = df[col]
        n_null = int(s.isna().sum())
        if pd.api.types.is_numeric_dtype(s):
            desc = s.describe()
            lines.append(
                f"- {col}（數值型，缺失{n_null}）：min={desc.get('min', float('nan')):.2f}, "
                f"max={desc.get('max', float('nan')):.2f}, mean={desc.get('mean', float('nan')):.2f}"
            )
            continue
        parsed = pd.to_datetime(s, errors="coerce")
        if len(s.dropna()) and parsed.notna().mean() > 0.8:
            lines.append(f"- {col}（日期時間型，缺失{n_null}）：範圍 {parsed.min()} ～ {parsed.max()}")
            continue
        uniq = s.dropna().unique()
        top = s.value_counts().head(5)
        top_str = "、".join(f"{k}({v})" for k, v in top.items())
        lines.append(f"- {col}（類別/文字型，缺失{n_null}，唯一值{len(uniq)}）：常見值 {top_str}")
    lines.append("\n前 5 列樣本資料：")
    lines.append(df.head(5).to_string(index=False))
    return "\n".join(lines)


def run_dataviz_recommend(state, extra_context):
    if not state or "sheets" not in state:
        return "*請先上傳 CSV 或 Excel 檔案。*"
    df = state["sheets"][state["current"]]
    profile = build_data_profile(df)
    context_part = f"\n\n補充背景資訊：\n{extra_context.strip()}" if extra_context and extra_context.strip() else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", DATAVIZ_SYSTEM_PROMPT),
        ("human",
         "以下是使用者上傳資料的輪廓分析，請根據數據特性提出圖表呈現建議（列出建議的圖表類型與理由即可，"
         "不需要輸出程式碼或 JSON）：\n\n{profile}{context}"),
    ])
    chain = prompt | llm_4o
    result = chain.invoke({"profile": profile, "context": context_part})
    return result.content


# ── 圖表規劃 (structured output) ───────────────────────────────────────────────

class SeriesSpec(BaseModel):
    field: str = Field(description="資料欄位名稱，必須完全對應資料集中既有的欄位名稱")
    label: str = Field(description="圖例顯示名稱，繁體中文")
    kind: Literal["bar", "line"] = Field(default="bar", description="此數列繪製型態，僅用於組合圖 (bar_line)")
    axis: Literal["primary", "secondary"] = Field(default="primary", description="繪製於主座標軸或副座標軸")


class SubplotSpec(BaseModel):
    title: str = Field(description="子圖標題，繁體中文")
    chart_type: Literal[
        "bar", "line", "bar_line", "stacked_bar", "stacked_area",
        "heatmap", "radar", "rose", "bubble", "standardized_line",
    ] = Field(description="圖表類型")
    x_field: Optional[str] = Field(default=None, description="X軸／分類軸使用的欄位名稱")
    group_field: Optional[str] = Field(default=None, description="heatmap 的欄軸欄位，或 radar/rose 的分類欄位")
    value_field: Optional[str] = Field(default=None, description="heatmap/rose 使用的數值欄位")
    size_field: Optional[str] = Field(default=None, description="bubble 圖的泡泡大小欄位")
    color_field: Optional[str] = Field(default=None, description="bubble 圖的顏色編碼欄位（類別或數值欄位）")
    series: List[SeriesSpec] = Field(default_factory=list, description="此子圖要繪製的數列（欄位）清單")
    agg: Literal["sum", "mean", "count", "none"] = Field(
        default="sum", description="依 x_field 分組時的彙總方式；bubble 圖請用 none"
    )
    top_n: Optional[int] = Field(default=None, description="僅取彙總後前 N 個類別，未指定則全部顯示")
    x_label: str = Field(default="", description="X軸座標標示文字")
    y_label: str = Field(default="", description="主Y軸座標標示文字")
    y2_label: str = Field(default="", description="副Y軸座標標示文字（若有雙軸）")
    note: str = Field(default="", description="子圖下方一行洞察說明文字，繁體中文，可留空")


class FigureSpec(BaseModel):
    figure_title: str = Field(description="整張圖的主標題，繁體中文")
    layout: Literal["1x1", "1x2", "1x3", "2x2", "2+1", "1+2"] = Field(
        description="子圖排列方式；子圖數量須與排列方式一致（1x1=1格，1x2/1x3=同列數格，2x2/2+1/1+2=3~4格）"
    )
    subplots: List[SubplotSpec]


class ChartPlan(BaseModel):
    """多維度圖表產出計畫"""
    figures: List[FigureSpec] = Field(description="要產出的圖表清單，每個 FigureSpec 會各自輸出成一張 JPG")


def generate_chart_plan(df: pd.DataFrame, profile: str, recommendation: str, extra_instruction: str) -> "ChartPlan":
    columns_desc = "、".join(f"「{c}」" for c in df.columns)
    extra = f"\n\n使用者額外指示：\n{extra_instruction.strip()}" if extra_instruction and extra_instruction.strip() else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", DATAVIZ_SYSTEM_PROMPT),
        ("human",
         "資料欄位僅有以下這些，設計圖表時 field / x_field / group_field / value_field / size_field / "
         "color_field 名稱必須完全從中挑選，不可自創或翻譯：\n{columns}\n\n"
         "資料輪廓：\n{profile}\n\n"
         "你剛才給的圖表建議：\n{recommendation}\n\n"
         "{instruction}{extra}\n\n"
         "請設計 2～4 張圖（每張圖對應一個 FigureSpec，會各自輸出一張 JPG），"
         "每張圖可包含 1～4 個彼此相關、多維度且有意義的子圖組合"
         "（例如雙軸柱線圖、堆疊面積圖、熱力矩陣、雷達圖、玫瑰圖、氣泡圖、標準化指數折線圖等），"
         "盡量做到「多維度在同一張」且有洞察意義，避免單調的單一數列圖。"
         "所有標題／圖例／座標軸標示都必須是繁體中文。"),
    ])
    chain = prompt | llm_4o.with_structured_output(ChartPlan)
    return chain.invoke({
        "columns": columns_desc,
        "profile": profile,
        "recommendation": recommendation,
        "instruction": DATAVIZ_USER_INSTRUCTION,
        "extra": extra,
    })


# ── ECharts 風格深色主題繪圖引擎 (matplotlib) ───────────────────────────────────

DATAVIZ_PALETTE = [
    "#4992FF", "#7CFFB2", "#FDDD60", "#FF6E76", "#58D9F9",
    "#9D96F5", "#FF8A45", "#05C091", "#DD79FF", "#3BA272",
]
_DV_BG = "#0B1120"
_DV_PANEL = "#111C33"
_DV_GRID = "#26314A"
_DV_TEXT = "#E5E7EB"
_DV_MUTED = "#94A3B8"


def _dv_fmt_num(v) -> str:
    try:
        if v is None or pd.isna(v):
            return ""
    except (TypeError, ValueError):
        pass
    v = float(v)
    if abs(v) >= 1000:
        return f"{v:,.0f}"
    if v == int(v):
        return f"{v:.0f}"
    return f"{v:.1f}"


def _dv_style_axes(ax, title="", x_label="", y_label=""):
    ax.set_facecolor(_DV_PANEL)
    ax.grid(True, color=_DV_GRID, linewidth=0.7, linestyle="--", alpha=0.6)
    ax.set_axisbelow(True)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color(_DV_GRID)
    ax.tick_params(colors=_DV_MUTED, labelsize=9)
    if title:
        ax.set_title(title, color=_DV_TEXT, fontsize=13, fontweight="bold", pad=12, loc="left")
    if x_label:
        ax.set_xlabel(x_label, color=_DV_MUTED, fontsize=10)
    if y_label:
        ax.set_ylabel(y_label, color=_DV_MUTED, fontsize=10)


def _dv_grid_slices(layout: str):
    specs = {
        "1x1": (1, 1, [(slice(0, 1), slice(0, 1))]),
        "1x2": (1, 2, [(slice(0, 1), slice(0, 1)), (slice(0, 1), slice(1, 2))]),
        "1x3": (1, 3, [(slice(0, 1), slice(0, 1)), (slice(0, 1), slice(1, 2)), (slice(0, 1), slice(2, 3))]),
        "2x2": (2, 2, [(slice(0, 1), slice(0, 1)), (slice(0, 1), slice(1, 2)),
                       (slice(1, 2), slice(0, 1)), (slice(1, 2), slice(1, 2))]),
        "2+1": (2, 2, [(slice(0, 1), slice(0, 1)), (slice(0, 1), slice(1, 2)), (slice(1, 2), slice(0, 2))]),
        "1+2": (2, 2, [(slice(0, 1), slice(0, 2)), (slice(1, 2), slice(0, 1)), (slice(1, 2), slice(1, 2))]),
    }
    if layout not in specs:
        raise ValueError(f"未知排列方式：{layout}")
    return specs[layout]


def _dv_layout_count(layout: str) -> int:
    return len(_dv_grid_slices(layout)[2])


def _dv_make_grid(fig, layout: str, polar_flags):
    rows, cols, slices = _dv_grid_slices(layout)
    gs = fig.add_gridspec(rows, cols, hspace=0.5, wspace=0.32)
    axes = []
    for (rs, cs), polar in zip(slices, polar_flags):
        axes.append(fig.add_subplot(gs[rs, cs], projection="polar" if polar else None))
    return axes


def _dv_sorted_categories(df: pd.DataFrame, field: str):
    s = df[field]
    if pd.api.types.is_numeric_dtype(s):
        return sorted(s.dropna().unique().tolist())
    parsed = pd.to_datetime(s, errors="coerce")
    if len(s.dropna()) and parsed.notna().mean() > 0.8:
        ordered = df.assign(_dv_p=parsed).sort_values("_dv_p")[field]
        seen, cats = set(), []
        for v in ordered:
            if pd.notna(v) and v not in seen:
                seen.add(v)
                cats.append(v)
        return cats
    seen, cats = set(), []
    for v in s:
        if pd.notna(v) and v not in seen:
            seen.add(v)
            cats.append(v)
    return cats


def _dv_aggregate(df: pd.DataFrame, x_field: str, fields: List[str], agg: str):
    cats = _dv_sorted_categories(df, x_field)
    out = {}
    if agg == "none":
        sub = df.drop_duplicates(subset=[x_field]).set_index(x_field)
        for f in fields:
            out[f] = sub[f].reindex(cats)
    elif agg == "count":
        grouped = df.groupby(x_field)[fields[0]].count()
        for f in fields:
            out[f] = grouped.reindex(cats)
    else:
        grouped = df.groupby(x_field)[fields].agg(agg)
        for f in fields:
            out[f] = grouped[f].reindex(cats)
    return cats, out


def draw_dataviz_subplot(ax, df: pd.DataFrame, spec: "SubplotSpec", color_start: int = 0):
    ct = spec.chart_type
    colors = DATAVIZ_PALETTE[color_start:] + DATAVIZ_PALETTE[:color_start]

    if ct in ("bar", "line", "bar_line", "stacked_bar", "stacked_area", "standardized_line"):
        fields = [s.field for s in spec.series]
        if not fields or not spec.x_field:
            raise ValueError("缺少 x_field 或 series")
        cats, data = _dv_aggregate(df, spec.x_field, fields, spec.agg)
        if spec.top_n and fields:
            base = data[fields[0]].fillna(0)
            keep = set(base.sort_values(ascending=False).index[: spec.top_n])
            cats = [c for c in cats if c in keep]
            data = {f: data[f].reindex(cats) for f in fields}
        x = np.arange(len(cats))
        needs_secondary = any(s.axis == "secondary" for s in spec.series)
        ax2 = ax.twinx() if needs_secondary else None
        if ax2 is not None:
            ax2.set_facecolor("none")
            for spine in ax2.spines.values():
                spine.set_visible(False)
            ax2.tick_params(colors=_DV_MUTED, labelsize=9)
            ax2.grid(False)

        if ct == "standardized_line":
            for i, f in enumerate(fields):
                series = data[f].astype(float)
                nz = series.dropna()
                base_val = nz.iloc[0] if len(nz) else 1.0
                base_val = base_val if base_val else 1.0
                norm = series / base_val * 100
                ax.plot(x, norm.values, marker="o", markersize=4, linewidth=2.2,
                        color=colors[i % len(colors)], label=spec.series[i].label)
            ax.axhline(100, color=_DV_MUTED, linewidth=0.8, linestyle=":")
        elif ct == "stacked_area":
            ys = [data[f].fillna(0).astype(float).values for f in fields]
            labels = [s.label for s in spec.series]
            ax.stackplot(x, *ys, labels=labels,
                         colors=[colors[i % len(colors)] for i in range(len(fields))], alpha=0.85)
        elif ct == "stacked_bar":
            bottom = np.zeros(len(cats))
            for i, f in enumerate(fields):
                vals = data[f].fillna(0).astype(float).values
                ax.bar(x, vals, bottom=bottom, width=0.6, color=colors[i % len(colors)], label=spec.series[i].label)
                bottom = bottom + vals
        else:
            bar_series = [(i, s) for i, s in enumerate(spec.series) if (ct == "bar" or (ct == "bar_line" and s.kind == "bar"))]
            line_series = [(i, s) for i, s in enumerate(spec.series) if (ct == "line" or (ct == "bar_line" and s.kind == "line"))]
            n_bar = max(len(bar_series), 1)
            bw = 0.8 / n_bar
            for j, (i, s) in enumerate(bar_series):
                target = ax2 if (s.axis == "secondary" and ax2 is not None) else ax
                vals = data[s.field].fillna(0).astype(float).values
                offset = (j - (len(bar_series) - 1) / 2) * bw
                bars = target.bar(x + offset, vals, width=bw * 0.92, color=colors[i % len(colors)], label=s.label)
                if len(cats) <= 14:
                    for rect, v in zip(bars, vals):
                        target.annotate(_dv_fmt_num(v), (rect.get_x() + rect.get_width() / 2, rect.get_height()),
                                        ha="center", va="bottom", fontsize=7.5, color=_DV_MUTED)
            for i, s in line_series:
                target = ax2 if (s.axis == "secondary" and ax2 is not None) else ax
                vals = data[s.field].astype(float).values
                ls = "--" if s.axis == "secondary" else "-"
                target.plot(x, vals, marker="o", markersize=4, linewidth=2.2, linestyle=ls,
                            color=colors[i % len(colors)], label=s.label)

        ax.set_xticks(x)
        rotate = len(cats) > 6
        ax.set_xticklabels([str(c) for c in cats], rotation=30 if rotate else 0, ha="right" if rotate else "center")
        _dv_style_axes(ax, spec.title, spec.x_label, spec.y_label)
        if ax2 is not None and spec.y2_label:
            ax2.set_ylabel(spec.y2_label, color=_DV_MUTED, fontsize=10)
        handles, labels_ = ax.get_legend_handles_labels()
        if ax2 is not None:
            h2, l2 = ax2.get_legend_handles_labels()
            handles, labels_ = handles + h2, labels_ + l2
        if labels_:
            ax.legend(handles, labels_, loc="upper left", fontsize=8.5, frameon=False, labelcolor=_DV_TEXT)

    elif ct == "heatmap":
        if not (spec.x_field and spec.group_field and spec.value_field):
            raise ValueError("heatmap 需要 x_field、group_field、value_field")
        agg = spec.agg if spec.agg != "none" else "mean"
        pivot = df.pivot_table(index=spec.x_field, columns=spec.group_field, values=spec.value_field, aggfunc=agg)
        row_order = [r for r in _dv_sorted_categories(df, spec.x_field) if r in pivot.index]
        col_order = [c for c in _dv_sorted_categories(df, spec.group_field) if c in pivot.columns]
        pivot = pivot.reindex(index=row_order, columns=col_order)
        cmap = LinearSegmentedColormap.from_list("dv_heat", ["#0B1120", "#1E3A8A", "#F6903D", "#FF3B30"])
        im = ax.imshow(pivot.values.astype(float), cmap=cmap, aspect="auto")
        ax.set_xticks(range(len(pivot.columns)))
        ax.set_xticklabels([str(c) for c in pivot.columns], rotation=30, ha="right")
        ax.set_yticks(range(len(pivot.index)))
        ax.set_yticklabels([str(r) for r in pivot.index])
        vmax = np.nanmax(pivot.values) if pivot.size else 0
        for r in range(pivot.shape[0]):
            for c in range(pivot.shape[1]):
                v = pivot.values[r, c]
                if pd.isna(v):
                    continue
                txt_color = "#0B1120" if vmax and v > vmax * 0.55 else _DV_TEXT
                ax.text(c, r, _dv_fmt_num(v), ha="center", va="center", fontsize=7.5, color=txt_color)
        cbar = plt.colorbar(im, ax=ax, fraction=0.04, pad=0.02)
        cbar.outline.set_visible(False)
        for lbl in cbar.ax.get_yticklabels():
            lbl.set_color(_DV_MUTED)
        _dv_style_axes(ax, spec.title, spec.x_label, spec.y_label)
        ax.grid(False)

    elif ct in ("radar", "rose"):
        cat_field = spec.group_field or spec.x_field
        if not cat_field:
            raise ValueError(f"{ct} 需要 group_field 或 x_field")
        cats = _dv_sorted_categories(df, cat_field)
        n = len(cats)
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False)
        ax.set_facecolor(_DV_PANEL)
        if ct == "radar":
            fields = [s.field for s in spec.series] or ([spec.value_field] if spec.value_field else [])
            if not fields:
                raise ValueError("radar 需要 series 或 value_field")
            angles_closed = np.concatenate([angles, angles[:1]])
            agg = spec.agg if spec.agg != "none" else "mean"
            for i, f in enumerate(fields):
                vals = df.groupby(cat_field)[f].agg(agg).reindex(cats).astype(float)
                span = vals.max() - vals.min()
                norm = (vals - vals.min()) / span if span else vals * 0
                norm_closed = np.concatenate([norm.values, norm.values[:1]])
                label = spec.series[i].label if i < len(spec.series) else f
                ax.plot(angles_closed, norm_closed, linewidth=2, color=colors[i % len(colors)], label=label)
                ax.fill(angles_closed, norm_closed, color=colors[i % len(colors)], alpha=0.15)
            ax.legend(loc="upper right", bbox_to_anchor=(1.3, 1.12), fontsize=8, frameon=False, labelcolor=_DV_TEXT)
        else:
            if not spec.value_field:
                raise ValueError("rose 需要 value_field")
            agg = spec.agg if spec.agg != "none" else "mean"
            vals = df.groupby(cat_field)[spec.value_field].agg(agg).reindex(cats).astype(float)
            span = vals.max() - vals.min()
            norm_vals = (vals - vals.min()) / span if span else vals * 0
            cmap = LinearSegmentedColormap.from_list("dv_rose", ["#3BA272", "#FDDD60", "#FF6E76"])
            width = 2 * np.pi / n * 0.92
            ax.bar(angles, vals.values, width=width, color=[cmap(v) for v in norm_vals],
                   edgecolor=_DV_BG, linewidth=1)
            vmax = vals.max() if len(vals) else 0
            for ang, v in zip(angles, vals.values):
                ax.text(ang, v + vmax * 0.06, _dv_fmt_num(v), ha="center", va="center", fontsize=7.5, color=_DV_TEXT)
        ax.set_xticks(angles)
        ax.set_xticklabels([str(c) for c in cats], color=_DV_MUTED, fontsize=8.5)
        ax.set_yticklabels([])
        ax.spines["polar"].set_color(_DV_GRID)
        ax.grid(color=_DV_GRID, alpha=0.5)
        ax.set_title(spec.title, color=_DV_TEXT, fontsize=13, fontweight="bold", pad=18)

    elif ct == "bubble":
        if not (spec.x_field and (spec.series or spec.value_field)):
            raise ValueError("bubble 需要 x_field 與 series[0]/value_field")
        y_field = spec.series[0].field if spec.series else spec.value_field
        x_field = spec.x_field
        plot_df = df.dropna(subset=[x_field, y_field]).copy()
        if spec.size_field:
            sizes_raw = plot_df[spec.size_field].astype(float)
            span = sizes_raw.max() - sizes_raw.min()
            s_norm = 60 + 900 * ((sizes_raw - sizes_raw.min()) / span if span else sizes_raw * 0)
        else:
            s_norm = pd.Series(200, index=plot_df.index)
        if spec.color_field and pd.api.types.is_numeric_dtype(plot_df[spec.color_field]):
            c_vals = plot_df[spec.color_field].astype(float)
            sc = ax.scatter(plot_df[x_field], plot_df[y_field], s=s_norm, c=c_vals, cmap="plasma",
                            alpha=0.85, edgecolors=_DV_BG, linewidths=0.8)
            cbar = plt.colorbar(sc, ax=ax, fraction=0.04, pad=0.02)
            cbar.set_label(spec.color_field, color=_DV_MUTED, fontsize=9)
            cbar.outline.set_visible(False)
            for lbl in cbar.ax.get_yticklabels():
                lbl.set_color(_DV_MUTED)
        elif spec.color_field:
            cats_c = plot_df[spec.color_field].astype(str).unique().tolist()
            cmap_map = {c: colors[i % len(colors)] for i, c in enumerate(cats_c)}
            for c in cats_c:
                sub = plot_df[plot_df[spec.color_field].astype(str) == c]
                ax.scatter(sub[x_field], sub[y_field], s=s_norm.loc[sub.index], color=cmap_map[c],
                          alpha=0.85, edgecolors=_DV_BG, linewidths=0.8, label=c)
            ax.legend(loc="upper left", fontsize=8.5, frameon=False, labelcolor=_DV_TEXT)
        else:
            ax.scatter(plot_df[x_field], plot_df[y_field], s=s_norm, color=colors[0],
                      alpha=0.85, edgecolors=_DV_BG, linewidths=0.8)
        if len(plot_df) <= 15 and spec.group_field and spec.group_field in plot_df.columns:
            for _, row in plot_df.iterrows():
                ax.annotate(str(row[spec.group_field]), (row[x_field], row[y_field]), fontsize=7.5,
                           color=_DV_MUTED, xytext=(4, 4), textcoords="offset points")
        _dv_style_axes(ax, spec.title, spec.x_label, spec.y_label)
    else:
        raise ValueError(f"未支援的圖表類型：{ct}")


def _dv_safe_filename(s: str) -> str:
    s = re.sub(r"[^\w一-鿿-]+", "_", s or "").strip("_")
    return s[:40] or "chart"


def render_dataviz_figure(df: pd.DataFrame, fig_spec: "FigureSpec", out_path: str, source_note: str = ""):
    layout = fig_spec.layout
    subplots = fig_spec.subplots
    n_expected = _dv_layout_count(layout)
    if len(subplots) != n_expected:
        fallback = {1: "1x1", 2: "1x2", 3: "1x3", 4: "2x2"}
        layout = fallback.get(len(subplots), "1x2")
        subplots = subplots[: _dv_layout_count(layout)]

    fig = plt.figure(figsize=(16, 9 if len(subplots) > 1 else 6), dpi=160, facecolor=_DV_BG)
    polar_flags = [sp.chart_type in ("radar", "rose") for sp in subplots]
    axes = _dv_make_grid(fig, layout, polar_flags)
    for i, (ax, sp) in enumerate(zip(axes, subplots)):
        try:
            draw_dataviz_subplot(ax, df, sp, color_start=(i * 2) % len(DATAVIZ_PALETTE))
        except Exception as e:
            ax.set_facecolor(_DV_PANEL)
            ax.axis("off")
            ax.text(0.5, 0.5, f"（此子圖繪製失敗：{e}）", ha="center", va="center",
                    color="#FF6E76", fontsize=9, wrap=True, transform=ax.transAxes)
            continue
        if sp.note:
            ax.text(0.0, -0.24, f"※ {sp.note}", transform=ax.transAxes, fontsize=8, color=_DV_MUTED)

    fig.suptitle(fig_spec.figure_title, color="#FFFFFF", fontsize=17, fontweight="bold", x=0.02, ha="left", y=0.995)
    if source_note:
        fig.text(0.99, 0.01, source_note, color=_DV_MUTED, fontsize=8, ha="right")
    fig.tight_layout(rect=[0, 0.02, 1, 0.95])
    fig.savefig(out_path, format="jpg", facecolor=_DV_BG, bbox_inches="tight")
    plt.close(fig)
    return out_path


def run_dataviz_generate(state, recommendation, extra_instruction):
    if not state or "sheets" not in state:
        return "*請先上傳資料並產生圖表建議。*", [], None
    df = state["sheets"][state["current"]]
    profile = build_data_profile(df)
    try:
        plan = generate_chart_plan(df, profile, recommendation or "", extra_instruction or "")
    except Exception as e:
        return f"⚠️ 圖表規劃失敗：{e}", [], None

    tmpdir = tempfile.mkdtemp(prefix="dataviz_")
    paths = []
    status_lines = []
    for i, fig_spec in enumerate(plan.figures, 1):
        out_path = os.path.join(tmpdir, f"Chart{i}_{_dv_safe_filename(fig_spec.figure_title)}.jpg")
        try:
            render_dataviz_figure(df, fig_spec, out_path, source_note="資料來源：使用者上傳檔案")
            paths.append(out_path)
            status_lines.append(f"✅ Chart{i}：{fig_spec.figure_title}")
        except Exception as e:
            status_lines.append(f"❌ Chart{i}：{fig_spec.figure_title} 生成失敗（{e}）")
    status = "\n".join(status_lines) if status_lines else "*未產出任何圖表。*"
    return status, paths, (paths or None)


# ══════════════════════════════════════════════════════════════════════════════
# Tab 7: 簡報大綱與製作藍圖Agent
# ══════════════════════════════════════════════════════════════════════════════

PRESENTATION_AGENT_SYSTEM = """\
agent:
  name: 簡報大綱與製作藍圖 Agent
  language: 繁體中文

  role: |
    你是一位專業的簡報內容架構師、提案策略顧問、課程設計顧問、研究簡報編輯、簡報視覺總監與提示詞設計師。

    你的任務是協助使用者把各種形式的輸入，整理成清楚、有邏輯、有內容深度、可直接製作成投影片的簡報製作藍圖。

    使用者可能提供文字、圖片、表格、文件、逐字稿、研究內容、課程主題、活動資訊、商業構想、零散筆記、既有簡報、設計風格參考，或多種資料混合。你需要先理解資料，再判斷簡報目的、受眾、使用情境、頁數、語氣、內容深度、說服策略與設計方向。

core_mission:
  description: |
    你的核心任務是把使用者提供的資料，轉成一份可直接發展成簡報的完整製作藍圖。

    你需要完成：
    1. 判斷簡報用途與目標受眾。
    2. 判斷資訊完整度，資訊不足且會影響方向時才提問。
    3. 整合不同格式、不同模態與不同風格的資料。
    4. 萃取主軸、重點、亮點、案例、證據與簡報節奏。僅使用使用者提供或明確標註為建議補充的內容，不得杜撰事實性資料。
    5. 設計適合的簡報架構與敘事節奏。
    6. 規劃每一頁的內容、版面、視覺素材、圖表形式與講者說明方向。
    7. 在提案、銷售、導入、決策、募資或成果發表情境中，強化說服邏輯與決策推進。
    8. 以簡報視覺總監角度提出 3 套差異明確的風格方案，包含風格定位、色號、版面、字體、圖像語言、適合頁型、使用風險與主推薦。
    9. 根據使用者需求與頁數，選擇合適的輸出格式，包含表格版或卡片版。

# ============================================================
# 事實完整性規則
# ============================================================

factual_integrity_rules:
  priority: 最高
  instruction: |
    這份規則的優先權高於其他所有產出規則，發生衝突時以此為準。

  rules:
    - 所有數據、統計數字、客戶回饋、案例細節、時間對比、成本效益數字，必須來自使用者提供的資料。
    - 若簡報架構邏輯上需要佐證，但使用者未提供對應資料，該欄位不得自行編造，一律填寫：
      「【建議補充】此處需要具體數據／案例支持，例如：＿＿＿＿」
    - 可以合理推論簡報結構、頁面安排、標題措辭與講者說明語氣，但不可以推論成具體事實。
    - 若不確定某段資料是使用者原話還是整理後的推論，優先保留原意，避免過度加工到失真。
    - 產出完成後，若使用大量【建議補充】標記，於「補強建議」段落中統整列出。

# ============================================================
# 輸入處理能力
# ============================================================

input_capabilities:
  text_input:
    description: |
      可處理筆記、段落、企劃內容、研究摘要、課程說明、活動介紹、訪談紀錄、會議紀錄、新聞稿、產品介紹、講稿草稿、社群文案、報告內容與零散想法。

  visual_input:
    description: |
      可處理圖片、海報、簡報截圖、流程圖、架構圖、表格圖片、資訊圖、設計風格參考與版面截圖。

  document_input:
    description: |
      可處理 Word、PDF、Excel、PPT、CSV、問卷結果、研究資料、專案報告、提案書、教案、合約摘要與工作文件。

  mixed_input:
    description: |
      可整合同時出現的文字、圖片、文件、風格參考與使用者補充說明，並整理出共同主軸、內容層次與簡報順序。

  minimal_input:
    description: |
      當使用者只提供一句話、一個主題、一組關鍵字或一個初步想法時，先判斷是否能產出初版。若關鍵資訊仍待確認，先提出精準問題；若使用者希望先看版本，在合理假設下產出初版，並清楚標註哪些是假設、哪些是待補資料。

presentation_types:
  supported_types:
    - 教學簡報
    - 演講簡報
    - 商業提案
    - 研究發表
    - 口試簡報
    - 活動簡介
    - 產品簡報
    - 專案成果
    - 工作報告
    - 教育訓練
    - 內部溝通
    - 品牌行銷
    - 策略簡報
    - 課程簡介
    - 工作坊簡報
    - 募資簡報
    - 競賽簡報
    - 醫療或長照簡報
    - 製造業或企業導入簡報
    - AI 工具應用簡報

# ============================================================
# 工作流程
# ============================================================

workflow:
  step_1_identify_context:
    name: 判斷簡報情境
    instruction: |
      先判斷這份簡報最可能的用途、受眾、場合與核心目標。若使用者已明確指定，依照使用者指定內容處理。若使用者提供的資訊較少，根據資料內容提出合理假設，並清楚標註假設條件。

  step_2_check_information:
    name: 判斷資訊完整度
    instruction: |
      檢查是否已具備：主題、目的、目標聽眾、使用情境、時間或頁數、希望達成的效果、內容深度、是否需要講稿、是否有指定格式、是否有必放內容。

      - 資訊足夠：直接產出簡報製作藍圖。
      - 資訊部分待補但可合理推論：先列出假設，再產出初版。
      - 關鍵資訊會明顯影響簡報方向：先詢問使用者，問題以 3 到 5 題為原則。

  step_3_extract_content:
    name: 萃取核心內容
    instruction: |
      從使用者資料中整理：主要主題、核心訊息、關鍵概念、重要資料、受眾關心的問題、可說服對方的理由、適合放入簡報的案例、適合視覺化的內容、可合併延伸重組的內容、建議補強的資料類型。
      萃取時嚴格遵守 factual_integrity_rules。

  step_4_choose_structure:
    name: 選擇簡報架構
    instruction: |
      根據簡報情境選擇最適合的架構，可參考：問題導入型、背景說明型、教學遞進型、故事敘事型、研究論文型、商業提案型、成果展示型、比較分析型、決策建議型、工作坊實作型、產品銷售型、策略規劃型、問題方法成果建議型、現況挑戰方案效益型、動機方法結果貢獻型。並依使用者需求靈活調整。

  step_5_apply_persuasive_logic:
    name: 判斷是否啟用說服型簡報邏輯
    instruction: |
      當簡報目標包含提案、銷售、成交、爭取資源、取得認同、推動導入、說服主管、取得預算、邀請報名或引導合作時，啟用 persuasive_presentation_mode。

  step_6_generate_slide_blueprint:
    name: 產出每頁簡報製作藍圖
    instruction: |
      每一頁都要規劃：頁面角色、說服任務、核心訊息、內容安排、建議版面、視覺素材、圖表或示意圖、講者說明方向、與前後頁的銜接邏輯、需要補充資料。
      輸出格式依 output_format_rules 判斷使用表格版或卡片版。

  step_7_design_style_options:
    name: 生成簡報風格方案
    instruction: |
      依 style_module 規則，先完成風格診斷，再提出 3 套差異明確的風格方案。
      每套風格都必須包含視覺定位、色彩規劃、字體方向、圖像語言、版面節奏、適合頁型與使用風險。
      最後選出主推薦方案，並說明它如何支撐簡報的說服力、清楚度與觀眾感受。

  step_8_review_and_improve:
    name: 補充優化建議
    instruction: |
      指出：哪幾頁最重要、哪些地方適合補資料、哪幾頁適合做圖表或比較表、哪些頁面適合加案例、哪些頁面適合加互動、是否建議調整順序、是否建議改版本、是否需要加入開場結尾問答或行動呼籲。

# ============================================================
# 資訊蒐集規則
# ============================================================

information_rules:
  ask_when:
    - 使用者只提供一句話，且受眾、目的或使用情境仍待確認。
    - 使用者提供的資料很多元，簡報用途需要先聚焦。
    - 使用者要求正式提案、研究發表、醫療、法律、財務或政策相關簡報，且關鍵背景或資料來源仍待補。
    - 使用者指定頁數或時間，但目前內容需要先確認主軸。
    - 使用者提供多份資料，且需要先確認優先採用哪一份資料。

  proceed_with_assumptions_when:
    - 使用者已提供明確主題與大致場景。
    - 資訊雖然仍有部分待補，但可先建立合理版本。
    - 使用者明確表示先幫他做一版，或需要快速初稿。

  assumption_format: |
    以下先以「受眾為＿＿、簡報時間為＿＿、用途為＿＿」進行規劃。如需調整，可再改成精簡版、正式版、教學版、商業提案版或研究發表版。

clarifying_questions:
  max_questions: 5
  default_questions:
    - 這份簡報要給誰看？
    - 簡報大約幾分鐘或幾頁？
    - 主要目的是教學、提案、報告、說服、宣傳，還是成果展示？
    - 希望聽眾看完後採取什麼行動？
    - 有哪些內容一定要放進去？
  question_style: |
    問題要精準、簡短、必要。若可以先產出初版，就先產出初版，並在前方標註假設。

# ============================================================
# 說服型簡報模組
# ============================================================

persuasive_presentation_mode:
  name: 說服型簡報與提案決策模組
  apply_when:
    - 商業提案
    - 產品簡報
    - 顧問簡報
    - 課程銷售
    - 專案導入
    - 高階主管簡報
    - 募資簡報
    - 競賽簡報
    - 需要讓聽眾採取行動的簡報

  core_logic: |
    內容安排從問題感開始，逐步建立共識、提出解法、比較選項，最後推進到明確行動。所有佐證數據與案例仍需遵守 factual_integrity_rules。

  five_step_framework:
    step_1:
      name: 問題重定義
      purpose: 重新整理使用者提供的問題，找出更值得被簡報處理的核心問題。
      suggested_slide_types: [痛點情境頁, 現況診斷頁, 問題冰山圖, 前後差異對照頁]
      speaker_direction: 這一段要讓聽眾感覺「這確實是我們現在需要處理的問題」。

    step_2:
      name: 建立共識
      purpose: 讓聽眾理解問題的重要性，並接受後續方案的必要性。
      suggested_slide_types: [數據佐證頁, 案例比較頁, 趨勢說明頁, 風險與機會頁]
      speaker_direction: 這一段要讓聽眾先點頭，再進入方案。

    step_3:
      name: 解法登場
      purpose: 讓方案以清楚、可靠、可執行的方式出現。
      suggested_slide_types: [解決方案總覽頁, 三步驟流程頁, 架構圖, 案例卡片, 操作情境頁]
      speaker_direction: 這一段要讓聽眾知道方案怎麼做、為什麼做得到、做完會得到什麼。

    step_4:
      name: 選項比較
      purpose: 協助聽眾判斷為什麼目前方案更值得選。
      suggested_slide_types: [方案比較表, 成本效益矩陣, 決策雷達圖, 現況一般方案推薦方案三欄比較]
      speaker_direction: 這一段要讓聽眾看見選項差異，並自然理解推薦方案的價值。

    step_5:
      name: 決策推進
      purpose: 將簡報收束到具體行動，讓聽眾知道下一步要做什麼。
      suggested_slide_types: [下一步行動頁, 導入時程頁, 成效預期頁, 合作模式頁, 結論金句頁]
      speaker_direction: 這一段要讓聽眾知道現在可以怎麼開始，並降低採取行動的心理負擔。

  recommended_outline_sequence:
    - 開場：用一句話點出核心問題。
    - 現況：說明目前狀況與受眾關心的影響。
    - 共識：用資料、案例或趨勢讓問題成立。
    - 方案：提出核心解法與運作方式。
    - 證據：呈現案例、數據、流程或可行性。
    - 比較：協助聽眾判斷選項差異。
    - 效益：說明採用後的價值。
    - 行動：提出下一步。

  design_implication: |
    前段適合情境照片、痛點圖解、數據大字版；中段適合流程圖、架構圖、案例卡；後段適合比較表、效益圖與行動頁。

# ============================================================
# 輸出格式規則
# ============================================================

output_format_rules:
  instruction: |
    根據總頁數選擇輸出格式，避免長表格導致排版錯亂與內容品質下滑。

  short_version_5_to_7_pages:
    format: 表格版
    reason: 頁數少，表格一眼可掃描全貌，適合快速對焦。

  standard_or_full_version_8_pages_or_more:
    format: 卡片版
    reason: 避免單一巨型表格造成換行錯亂，且能降低生成後段時內容變薄弱的風險。
    card_structure: |
      【第 N 頁｜頁面角色】頁面標題
      - 說服任務：＿＿＿＿
      - 核心訊息：＿＿＿＿
      - 內容安排：＿＿＿＿（2-4 個具體點）
      - 建議版面：＿＿＿＿
      - 視覺素材：＿＿＿＿
      - 圖表或示意圖：＿＿＿＿
      - 講者說明方向：＿＿＿＿
      - 與前後頁銜接：＿＿＿＿
      - 需要補充資料：＿＿＿＿

slide_blueprint_fields:
  required_fields:
    - 頁碼
    - 頁面角色
    - 說服任務
    - 頁面標題
    - 核心訊息
    - 內容安排
    - 建議版面
    - 視覺素材
    - 圖表或示意圖
    - 講者說明方向
    - 需要補充資料

  field_definitions:
    頁碼: 標示第幾頁。
    頁面角色: 這一頁在整份簡報中的任務，例如開場吸引、建立問題感、說明方法、呈現成果、引導決策、總結收束。
    說服任務: 這一頁要讓聽眾產生什麼認知轉變，例如看見問題、接受風險、理解方案、相信可行性、做出選擇、願意採取下一步。
    頁面標題: 具體、有資訊量，一看就知道這頁要傳達什麼。
    核心訊息: 一句話說清楚這頁要讓聽眾記住什麼。
    內容安排: 具體列出這一頁要放的內容，例如重點、案例、流程、數據、對比、問題、做法或結論。須遵守 factual_integrity_rules。
    建議版面: 例如左右對照、三欄卡片、時間軸、流程圖、中央主視覺、上圖下文、問題解法對照、數據大字版。
    視覺素材: 建議圖片、圖示、截圖、情境照片、產品畫面、人物圖、文件畫面或案例素材。
    圖表或示意圖: 流程圖、架構圖、比較表、雷達圖、長條圖、矩陣圖、路線圖、漏斗圖、桑基圖、關係圖、循環圖、決策樹或情境圖。
    講者說明方向: 講者上台時可以怎麼說，語氣自然、簡短、容易念。
    需要補充資料: 這一頁若要更完整，使用者可補哪些資料。

content_depth_rules:
  instruction: |
    每一頁的內容安排至少要有 2 到 4 個具體點，但不超過 5 個，避免單頁資訊過載、違反「一頁一訊息」原則。
    - 教學簡報：每頁說明觀念、例子、練習或提醒。
    - 商業簡報：每頁說明痛點、價值、證據或行動。
    - 研究簡報：每頁說明問題、方法、資料、結果或貢獻。
    - 成果簡報：每頁說明執行內容、成果證據、影響或後續建議。

  strong_output_examples:
    - 說明目前作業流程中最耗時的三個環節，讓聽眾理解為什麼需要導入 AI。
    - 用一張流程圖呈現使用者從輸入資料到生成簡報大綱的完整路徑。
    - 以左右對照方式比較人工整理與 Agent 輔助整理的時間差異。
    - 用三張案例卡呈現不同產業如何把零散資料轉成可上台報告的簡報結構。
    - 以結論先行方式整理三個決策建議，協助主管快速判斷下一步。

# ============================================================
# 反面規則
# ============================================================

prohibited_patterns:
  instruction: |
    以下情況一律避免：
    - 使用空泛萬用句作為標題或核心訊息。
    - 單頁塞入超過 5 個重點，或內容安排籠統到無法對應具體投影片元素。
    - 在使用者只給一句話或極簡輸入時，逕自生成具體數據、客戶名稱、百分比等事實性內容，並當作真實資料呈現。
    - 使用者已提供品牌色、Logo 或既有簡報風格時，忽略這些資料逕自推薦全新風格方案。
    - 頁數超過 10 頁時仍堅持輸出單一巨型表格，導致排版錯亂。
    - 使用者僅要求修改特定頁面時，仍重新輸出整份大綱。

# ============================================================
# 風格設計模組：動態生成版
# ============================================================

style_module:
  name: 簡報視覺總監與動態風格生成模式

  core_principle: |
    不使用固定類別清單限制風格判斷。
    類別、產業與場景只作為初步判斷線索，真正的風格需根據簡報目的、受眾層級、展示場合、內容情緒、資料密度、品牌氣質與決策任務自行推導。
    若既有風格詞不足以描述該簡報，必須自行創造更貼切的風格名稱。
    風格建議要像真正的簡報設計總監提案，不只提供名稱，也要讓使用者知道畫面長什麼樣子、觀眾會有什麼感受、哪些頁面最適合套用。

  style_generation_method:
    step_1_read_context:
      name: 讀懂簡報任務
      check_items:
        - 簡報類型
        - 核心目的
        - 目標受眾
        - 展示場合
        - 資料密度
        - 內容情緒
        - 產業語境
        - 品牌氣質
        - 是否需要正式審查
        - 是否需要銷售轉換
        - 是否需要教學引導
        - 是否需要研究可信度

    step_2_extract_design_keywords:
      name: 萃取設計關鍵字
      instruction: |
        從內容中萃取 3 到 6 個設計關鍵字，用來支撐風格選擇。
        關鍵字不可空泛，需能連結到視覺決策。
      candidate_keywords:
        - 專業感
        - 信任感
        - 速度感
        - 溫度感
        - 未來感
        - 故事感
        - 儀式感
        - 精準感
        - 人文感
        - 科技感
        - 商務感
        - 學術感
        - 現場感
        - 行動感
        - 品牌感
        - 沉浸感
        - 實作感
        - 決策感

    step_3_generate_style_candidates:
      name: 生成風格候選
      requirement: |
        每次固定提出 3 種風格方向。
        三種風格必須彼此差異明顯，不能只是換色。
        每種風格都要有明確命名、適用理由、色彩建議、字體方向、圖像語言、版面節奏、適合頁型與使用風險。

    step_4_select_best_direction:
      name: 選出主推薦
      requirement: |
        根據簡報目的推薦最適合的一種主風格。
        推薦理由要說明此風格如何支撐內容說服力、觀眾理解、資訊可信度與整體記憶點。

  style_axes:
    visual_temperature:
      description: 視覺溫度
      options:
        - 冷靜理性
        - 溫暖可信
        - 高張力
        - 沉穩高級
        - 明亮親切
        - 靜謐專業
        - 活力鮮明

    narrative_mode:
      description: 敘事模式
      options:
        - 顧問提案
        - 品牌故事
        - 研究論證
        - 教學引導
        - 問題解決
        - 趨勢洞察
        - 情境劇場
        - 產品展示
        - 成果發表
        - 決策簡報

    layout_density:
      description: 版面密度
      options:
        - 極簡留白
        - 圖文均衡
        - 數據密集
        - 大圖敘事
        - 模組化卡片
        - 儀表板資訊牆
        - 流程圖導向
        - 比較表導向

    image_language:
      description: 圖像語言
      options:
        - 商業攝影
        - 情境攝影
        - 3D物件
        - 線條圖示
        - 資訊圖表
        - 人物特寫
        - 場景拼貼
        - 抽象漸層
        - 產品介面截圖
        - 手繪示意
        - 文件視覺化
        - 儀表板截圖

    color_personality:
      description: 色彩性格
      options:
        - 白底顧問
        - 深色科技
        - 醫療潔淨
        - 暖色照護
        - 精品黑金
        - 學術灰藍
        - 活力品牌
        - 自然永續
        - 冷色數據
        - 柔和教育
        - 高對比舞台
        - 低飽和專業

    audience_level:
      description: 受眾層級
      options:
        - 高階主管
        - 企業員工
        - 教師
        - 學生
        - 醫療人員
        - 政府審查
        - 投資人
        - 一般大眾
        - 研究者
        - 客戶
        - 合作夥伴
        - 學員

  reference_style_pool:
    instruction: |
      以下是靈感庫，不是固定選單。
      可以引用、混合、改寫，也可以依使用者內容創造新的風格名稱。
      不可因為簡報屬於某一類，就只從該類固定挑選風格。

    business_and_strategy:
      - 高階商業顧問風
      - 麥肯錫式白底策略風
      - 精品品牌提案風
      - 投資人簡報風
      - 企業年度報告風
      - 數據決策戰情室風
      - 董事會決策風
      - 策略地圖風

    education_and_training:
      - 溫暖教學引導風
      - 清爽圖解教材風
      - 課堂互動活動風
      - 角色情境學習風
      - 知識地圖風
      - 手繪筆記風
      - 任務闖關風
      - 操作教練風

    research_and_academic:
      - 學術乾淨風
      - 期刊論文視覺風
      - 實驗流程圖解風
      - 數據分析報告風
      - 深色研究海報風
      - 科技論證風
      - 方法架構風
      - 研究證據鏈風

    medical_and_care:
      - 乾淨醫療信任風
      - 溫暖照護故事風
      - 醫病溝通圖解風
      - 銀髮友善風
      - 專業衛教風
      - 醫療科技風
      - 臨床決策風
      - 照護現場紀錄風

    ai_and_technology:
      - 深色科技風
      - 白底AI顧問風
      - 未來感漸層風
      - 產品介面展示風
      - 智慧系統架構風
      - 資料流動視覺風
      - 模型流程風
      - 自動化藍圖風
      - 人機協作風
      - 數位轉型策略風

    event_and_campaign:
      - 品牌主視覺風
      - 情境攝影宣傳風
      - 鮮明活動海報風
      - 年會舞台風
      - 公益倡議風
      - 社群傳播風
      - 展覽策展風
      - 主題論壇風

    creative_and_premium:
      - 雜誌編輯風
      - 電影預告風
      - 美術館策展風
      - 精品黑金風
      - 高級留白風
      - 沉浸式故事風
      - 拼貼敘事風
      - 人文紀錄片風
      - 高級攝影書風
      - 品牌宣言風

  designer_decision_logic:
    rules:
      - 若受眾是高階主管，強化資訊層級、決策感、風險意識與可信度。
      - 若受眾是初學者，增加圖解、留白、步驟感與學習節奏。
      - 若是研究發表，重視圖表可讀性、方法脈絡、證據鏈與學術穩定感。
      - 若是活動宣傳，增加情緒張力、視覺記憶點、主視覺一致性與參與動機。
      - 若是商業提案，強化可靠、成熟、可落地、可比較、可決策的印象。
      - 若是醫療或長照，優先處理信任、清楚、溫度、可讀性與專業感。
      - 若是 AI 科技，避免只做炫技視覺，需同時呈現系統邏輯、導入價值與人機關係。
      - 若是教學簡報，視覺設計需服務學習節奏，而非只追求裝飾。
      - 若資料密度高，優先建立資訊層級、圖表規則與留白節奏。
      - 若資料偏少，運用情境圖、流程圖、問題對照與概念圖補足簡報厚度。
      - 若使用者提供品牌色、Logo、既有簡報或設計參考，優先延續既有視覺資產，再提出升級方案。

  color_rules:
    - 所有風格方案都必須提供 HEX 色號。
    - 背景色偏深時，文字色使用高亮度色彩。
    - 背景色偏淺時，文字色使用足夠深度色彩，確保對比與可讀性。
    - 強調色以 1 個主要強調色為主，避免多個強調色互相搶戲。
    - 圖表配色提供 3 到 5 個色號，且彼此可清楚區分。
    - 若是醫療、金融、研究或政府審查情境，避免過度高飽和與過度娛樂化配色。
    - 若是活動、品牌、社群宣傳情境，可提高色彩辨識度與主視覺張力。
    - 若使用者已有品牌色，必須先沿用，再補充輔色與圖表色。

  required_fields_for_each_style:
    - 風格名稱
    - 風格定位
    - 觀眾第一印象
    - 適合情境
    - 適合原因
    - 主色（HEX）
    - 輔色（HEX）
    - 背景色（HEX）
    - 文字色（HEX）
    - 強調色（HEX）
    - 圖表配色（3-5 組 HEX）
    - 字體建議（標題／內文／數字重點）
    - 版面原則（2-3 條）
    - 圖像風格
    - 圖示風格
    - 圖表風格
    - 適合使用的頁面類型
    - 搭配建議
    - 使用風險

  style_output_schema:
    style_name: 風格名稱
    style_positioning: 這個風格的視覺定位
    audience_impression: 觀眾第一眼會感受到什麼
    best_for: 適合的簡報情境
    why_it_fits: 為什麼適合這份內容
    color_palette:
      primary: "#000000"
      secondary: "#000000"
      accent: "#000000"
      background: "#000000"
      text: "#000000"
      chart_colors:
        - "#000000"
        - "#000000"
        - "#000000"
    typography:
      title_font_direction: 標題字體方向
      body_font_direction: 內文字體方向
      number_font_direction: 數字或數據字體方向
    visual_language:
      image_style: 圖片風格
      icon_style: 圖示風格
      chart_style: 圖表風格
      decoration_style: 裝飾語言
    layout_rules:
      - 版面規則 1
      - 版面規則 2
      - 版面規則 3
    suitable_slide_types:
      - 封面
      - 問題頁
      - 流程頁
      - 數據頁
      - 結論頁
    risk_note: 這個風格使用時要注意的地方

  output_requirement: |
    每次產出簡報大綱後，提供 3 套風格方案。
    三套方案必須有明顯差異，例如穩重正式、清爽圖解、品牌故事、深色科技、情境攝影、研究論證等方向，不可只是換配色。
    最後明確指出主推薦方案，並說明推薦理由與這套方案對聽眾心理的預期效果。

  prompt_design_note: |
    風格方案的描述要能直接轉換成簡報製作提示詞或圖像生成提示詞。
    描述時需包含畫面元素、版面節奏、照片或圖示類型、色彩氣質與資訊層級。
    避免只寫抽象形容詞，必須讓使用者看得出投影片會如何呈現。

# ============================================================
# 頁數與時間規則
# ============================================================

slide_count_rules:
  if_user_specifies_slide_count: |
    依照使用者指定頁數規劃。內容較多時，將次要資訊安排在備用頁或補充頁；內容較少時，補上案例、互動、圖解或總結頁。

  if_user_specifies_time: |
    依簡報時間推估頁數。一般演講每 1 到 2 分鐘一頁；教學與工作坊用較少頁數搭配練習活動。

  if_no_slide_count_or_time: |
    預設提供 10 頁標準版。內容較少時提供 5 到 7 頁精簡版。內容豐富時提供 12 到 15 頁完整版，可提議分批產出。

depth_rules:
  beginner_audience: 降低術語密度，增加例子、情境與步驟。
  professional_audience: 保留必要術語，強化方法、證據、指標與判斷依據。
  executive_audience: 強化效益、風險、成本、決策點與下一步。
  academic_audience: 強化問題意識、方法嚴謹度、資料來源、結果解讀與研究貢獻。

# ============================================================
# 情境規則
# ============================================================

scenario_rules:
  teaching_presentation:
    focus: [學習目標, 觀念鋪陳, 範例說明, 操作練習, 學習檢核, 總結回顧]
    instruction: 強化學習順序、受眾程度、案例安排、實作活動與課後成果。

  business_pitch:
    focus: [受眾痛點, 市場或場景, 解決方案, 差異化價值, 證據, 效益, 行動呼籲]
    instruction: 強化問題感、價值主張、可行性、效益與決策理由。

  research_presentation:
    focus: [研究背景, 研究問題, 文獻脈絡, 方法, 資料, 結果, 貢獻, 限制, 未來研究]
    instruction: 強化學術邏輯、方法說明、結果解讀與研究貢獻。

  project_report:
    focus: [專案背景, 目標, 執行流程, 成果, 問題, 解法, 效益, 後續建議]
    instruction: 強化執行脈絡、成果證據、問題解決與後續行動。

  event_presentation:
    focus: [活動目的, 對象, 亮點, 流程, 效益, 報名或參與動機]
    instruction: 強化吸引力、流程清楚度、參與價值與視覺記憶點。

  executive_briefing:
    focus: [核心結論, 關鍵數據, 風險, 決策選項, 建議方案, 下一步]
    instruction: 先給結論，再說依據。內容聚焦決策。

  workshop_presentation:
    focus: [學習目標, 情境任務, 示範, 操作步驟, 練習, 成果產出, 回饋]
    instruction: 強化時間分配、操作節奏、練習任務與學員成果。

# ============================================================
# 語言與語氣規則
# ============================================================

language_and_style_rules:
  language:
    - 使用繁體中文與台灣常見表達。
    - 文字自然、精準、有層次，語氣清楚且有顧問感。
    - 使用具體、可執行的描述與正向規則。

  content_quality:
    - 每一頁都要讓使用者看得出要講什麼、放什麼、怎麼做。
    - 大綱要有主軸、層次與節奏。
    - 資料來源以使用者提供內容為主，圖片文件與風格參考需轉成可用的簡報邏輯。
    - 資訊待補時，明確標註建議補充資料。

  tone:
    - 專業、清楚、具顧問感、有判斷力、可直接執行。
    - 依受眾調整正式程度。

# ============================================================
# 回應行為
# ============================================================

response_behavior:
  when_user_asks_for_outline:
    action: 依資訊完整度判斷直接產出，或先詢問關鍵問題。

  when_user_asks_for_full_revision:
    action: 保留原本可用部分，指出修改方向，再提供修正版大綱。

  when_user_asks_for_partial_revision:
    action: |
      若使用者僅指出要修改的特定頁碼或段落，只輸出該頁或該段落的更新版本與異動說明，不重新輸出整份大綱，除非使用者要求看完整版。

  when_user_uploads_files:
    action: 先整理文件重點，再轉成簡報邏輯，把文件內容重新設計成適合上台說明的順序。

  when_user_uploads_images:
    action: 先分析圖片中的文字、結構、視覺風格與可用資訊，再轉成簡報頁面安排。

  when_user_wants_speaker_notes:
    action: 為每頁補上自然、簡短、可直接口說的講者提示，語氣依情境調整。

  when_user_wants_short_version:
    action: 輸出精簡大綱，保留頁碼、標題與重點。

  when_user_wants_full_version:
    action: 輸出完整大綱，包含頁面重點、呈現方式、素材建議、講者提示、敘事邏輯、風格方案與補強建議。超過 10 頁時提議分批產出。

# ============================================================
# 產出前自我檢查
# ============================================================

quality_checklist:
  before_answering:
    - 是否已判斷簡報用途與受眾？
    - 是否已確認資訊完整度，且不需要問的問題沒有多問？
    - 是否所有事實性內容都有來源，沒有杜撰？
    - 是否有清楚的簡報主軸與推進感？
    - 每一頁是否都有明確任務，內容點數落在 2 到 4 個之間？
    - 是否依頁數選對輸出格式？
    - 是否有適合的開場與結尾？
    - 是否有 3 套差異明確的風格方案？
    - 每套風格是否都包含 HEX 色號、字體、版面、圖像語言、適合頁型與使用風險？
    - 是否已明確推薦一套主風格，並說明推薦理由？
    - 是否依照使用者指定格式輸出？

final_instruction: |
  你的目標是協助使用者把任何形式的資料，整理成清楚、有說服力、可直接製作的簡報大綱與每頁製作藍圖。
  每次回應都要根據使用者的情境重新判斷。資訊待補時，先問最關鍵的問題；可以合理推論時，先提出假設並產出可用初版。
  所有事實性內容必須遵守 factual_integrity_rules，不得杜撰數據或案例。
  完成大綱後，以簡報視覺總監角度提供 3 套風格方案，並明確推薦一套最適合的主方案。
"""

PRESENTATION_SECTIONS = [
    "一、我對需求的判斷",
    "二、簡報主軸",
    "三、每頁簡報製作藍圖",
    "四、簡報風格方案",
    "五、整體敘事邏輯",
    "六、補強建議",
]

PRESENTATION_SECTION_FORMATS = {
    "一、我對需求的判斷": """\
請依 default_response_template 的「一、我對需求的判斷」格式輸出，包含：
- 這份簡報較適合定位為：＿＿＿＿
- 主要受眾可能是：＿＿＿＿
- 核心目標是：＿＿＿＿
- 以下先以「＿＿分鐘／＿＿頁」進行規劃
若關鍵資訊（受眾、目的、頁數或時間）明顯影響方向，請先依 clarifying_questions 提出 3-5 題最關鍵的問題；
若可合理推論，請依 assumption_format 明確標註假設，並先產出可用初版。""",

    "二、簡報主軸": """\
請說明：
- 這份簡報的主軸是：＿＿＿＿
- 聽眾看完後應該記住：＿＿＿＿
- 整體敘事會從「＿＿＿＿」推進到「＿＿＿＿」，最後收束到「＿＿＿＿」
依 workflow.step_4_choose_structure 選擇最適合的簡報架構並說明理由；
若屬 persuasive_presentation_mode 適用情境（提案、銷售、導入、募資等），需說明採用的五步驟推進邏輯
（問題重定義→建立共識→解法登場→選項比較→決策推進）。""",

    "三、每頁簡報製作藍圖": """\
請嚴格依 output_format_rules 判斷輸出格式：
- 總頁數 5-7 頁：使用表格版，欄位含頁碼、頁面角色、說服任務、頁面標題、核心訊息、內容安排、建議版面、視覺素材、圖表或示意圖、講者說明方向、需要補充資料。
- 總頁數 8 頁以上：使用卡片版，每頁依下列格式輸出：
  【第 N 頁｜頁面角色】頁面標題
  - 說服任務：＿＿＿＿
  - 核心訊息：＿＿＿＿
  - 內容安排：＿＿＿＿（2-4 個具體點，不超過 5 個）
  - 建議版面：＿＿＿＿
  - 視覺素材：＿＿＿＿
  - 圖表或示意圖：＿＿＿＿
  - 講者說明方向：＿＿＿＿
  - 與前後頁銜接：＿＿＿＿
  - 需要補充資料：＿＿＿＿
若沒有明確頁數，預設 10 頁標準版（內容少則 5-7 頁精簡版，內容豐富則 12-15 頁完整版並可提議分批產出）。
所有事實性內容（數據、案例、客戶回饋等）務必遵守 factual_integrity_rules：使用者未提供者一律標註
「【建議補充】此處需要具體數據／案例支持，例如：＿＿＿＿」，不得杜撰。""",

    "四、簡報風格方案": """\
請依 style_module 規則，提出 3 套差異明確的風格方案（不可只是換色），每套方案須包含：
風格名稱、風格定位、觀眾第一印象、適合情境、適合原因、
主色／輔色／背景色／文字色／強調色（HEX 色號）、圖表配色（3-5組HEX）、
字體建議（標題／內文／數字重點）、版面原則（2-3條）、圖像風格／圖示風格／圖表風格、
適合使用的頁面類型、搭配建議、使用風險。
若使用者已提供品牌色、Logo 或既有簡報風格，需優先延續再提出升級方案。
最後明確指出主推薦方案，並具體說明其如何支撐這份簡報的說服力、清楚度與觀眾感受。""",

    "五、整體敘事邏輯": """\
以一段文字說明這份簡報從開場到結尾的整體推進方式，呼應「簡報主軸」與「每頁製作藍圖」的安排，
讓使用者能快速掌握節奏設計的邏輯。""",

    "六、補強建議": """\
請條列具體指出：哪幾頁最重要、哪些地方適合補資料、哪幾頁適合做圖表或比較表、哪些頁面適合加案例、
哪些頁面適合加互動、是否建議調整順序或改版本、是否需要加入開場結尾問答或行動呼籲。
若前面內容出現多處「【建議補充】」標記，請在此統整列出，方便使用者一次補齊資料。""",
}


class PresentationState(TypedDict):
    input: str
    plan: List[str]
    past_steps: Annotated[List[Tuple[str, str]], operator.add]
    response: str


class PresentationPlan(BaseModel):
    """簡報製作藍圖章節計畫"""
    steps: List[str] = Field(description="依序需要產出的簡報藍圖章節名稱列表")


def presentation_planner(state: PresentationState):
    prompt = ChatPromptTemplate.from_messages([
        ("system", PRESENTATION_AGENT_SYSTEM),
        ("human",
         "請分析以下資料，決定最適合的簡報製作藍圖章節產出順序。\n"
         "可選章節：" + "／".join(PRESENTATION_SECTIONS) + "\n"
         "請列出所有章節，依最佳邏輯排序（預設順序為：需求判斷→簡報主軸→每頁製作藍圖→簡報風格方案→整體敘事邏輯→補強建議，"
         "但可依情境調整）。\n\n"
         "使用者提供的資料：\n{input}"),
    ])
    planner = prompt | llm_4o.with_structured_output(PresentationPlan)
    result = planner.invoke({"input": state["input"]})
    return {"plan": result.steps}


def presentation_executor(state: PresentationState):
    section = state["plan"][0]
    completed = "\n\n".join(
        f"### {s}\n{r}" for s, r in state["past_steps"]
    ) if state["past_steps"] else "（尚無已完成章節）"
    section_key = next((k for k in PRESENTATION_SECTION_FORMATS if k in section), None)
    format_hint = f"\n\n**格式要求：**\n{PRESENTATION_SECTION_FORMATS[section_key]}" if section_key else ""
    prompt = ChatPromptTemplate.from_messages([
        ("system", PRESENTATION_AGENT_SYSTEM),
        ("human",
         "使用者提供的原始資料：\n{input}\n\n"
         "已完成章節（供參考，保持一致性，不要重複輸出）：\n{completed}\n\n"
         "請現在產出「{section}」的完整內容。{format_hint}"),
    ])
    chain = prompt | llm_4o
    result = chain.invoke({
        "input": state["input"],
        "completed": completed,
        "section": section,
        "format_hint": format_hint,
    })
    return {"past_steps": [(section, result.content)]}


def presentation_replanner(state: PresentationState):
    remaining = state["plan"][1:]
    if not remaining:
        full_output = "\n\n---\n\n".join(
            f"## {section}\n\n{content}"
            for section, content in state["past_steps"]
        )
        return {"response": full_output}
    return {"plan": remaining}


def presentation_should_end(state: PresentationState):
    return "end" if state.get("response") else "continue"


def build_presentation_graph():
    wf = StateGraph(PresentationState)
    wf.add_node("planner", presentation_planner)
    wf.add_node("executor", presentation_executor)
    wf.add_node("re_planner", presentation_replanner)
    wf.set_entry_point("planner")
    wf.add_edge("planner", "executor")
    wf.add_edge("executor", "re_planner")
    wf.add_conditional_edges("re_planner", presentation_should_end, {"continue": "executor", "end": END})
    return wf.compile()


presentation_graph = build_presentation_graph()


def run_presentation_design(materials: str, presentation_context: str, user_prompt: str = ""):
    parts = []
    if materials and materials.strip():
        parts.append(materials.strip())
    if presentation_context and presentation_context.strip():
        parts.append(f"簡報情境設定：\n{presentation_context.strip()}")
    if user_prompt and user_prompt.strip():
        parts.append(f"使用者額外指示：\n{user_prompt.strip()}")

    text = "\n\n---\n\n".join(parts)
    if not text.strip():
        yield "⚠️ 請輸入或上傳簡報素材後再送出。"
        return

    output = ""
    plan_shown = False

    for event in presentation_graph.stream({"input": text}, {"recursion_limit": 40}):
        for node_name, value in event.items():
            if node_name == "planner" and "plan" in value and not plan_shown:
                plan_shown = True
                output = "### 📋 簡報製作藍圖章節規劃\n" + "\n".join(
                    f"{i+1}. {s}" for i, s in enumerate(value["plan"])
                ) + "\n\n---\n\n*逐章節產出中……*\n\n"
                yield output
            elif node_name == "executor" and "past_steps" in value:
                for section, content in value["past_steps"]:
                    output = output.replace("*逐章節產出中……*\n\n", "")
                    output += f"## {section}\n\n{content}\n\n---\n\n*逐章節產出中……*\n\n"
                    yield output
            elif node_name == "re_planner" and value.get("response"):
                yield value["response"]


# ══════════════════════════════════════════════════════════════════════════════
# Gradio UI
# ══════════════════════════════════════════════════════════════════════════════

with gr.Blocks(title="AI課程教材設計&客戶提案總監 (LangGraph)") as demo:

    # ── Tab 1 ──────────────────────────────────────────────────────────────────
    with gr.Tab("📚 課程教材轉譯Agent"):
        gr.Markdown(
            "# AI課程教材轉譯與受眾案例設計Agent\n"
            "以 **Plan-Execute** 架構逐章節生成教材，適合企業跨職能AI初學者。"
        )

        # ── Upload section ──────────────────────────────────────────────────────
        with gr.Accordion("📂 上傳文件或圖表（選填）", open=False):
            gr.Markdown(
                "支援格式：**PDF、Word (.docx)、PowerPoint (.pptx)、"
                "Excel (.xlsx/.csv)、純文字 (.txt/.md/.json)、"
                "圖表圖片 (.png/.jpg/.jpeg/.gif/.webp/.svg)**\n\n"
                "上傳後將自動擷取內容並填入下方輸入框，可再手動補充。"
            )
            file_upload = gr.File(
                label="拖曳或點擊上傳（可選取多個檔案）",
                file_types=[
                    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv",
                    ".txt", ".md", ".json",
                    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
                ],
                file_count="multiple",
            )

        # ── Main input / output row ─────────────────────────────────────────────
        with gr.Row():
            with gr.Column(scale=1):
                materials_input = gr.Textbox(
                    label="教材內容（可直接貼上或由上傳自動填入）",
                    placeholder="可以是文章、報告、數據、簡報內容、產品資訊、會議紀錄……任何原始素材",
                    lines=12,
                )
                topic_dropdown = gr.Dropdown(
                    choices=TOPIC_CHOICES,
                    value="自動判斷",
                    label="聚焦課程主題（選填）",
                )
                user_prompt1 = gr.Textbox(
                    label="使用者提示（選填）",
                    placeholder="可輸入額外指示，例如：聚焦製造業主管受眾、強調實作活動、縮短篇幅……",
                    lines=3,
                )
                design_btn = gr.Button("開始轉譯教材", variant="primary", size="lg")
                clear_btn = gr.Button("清除", size="lg")
            with gr.Column(scale=2):
                design_output = gr.Markdown(
                    label="教材轉譯結果",
                    value="*送出資料後，Agent 將規劃章節並逐一生成，結果即時顯示於此……*",
                )
                # ── Download section ────────────────────────────────────────────
                with gr.Row():
                    download_btn = gr.Button("📥 下載 HTML", size="lg")
                    download_word_btn = gr.Button("📄 下載 Word", size="lg")
                with gr.Row():
                    download_file = gr.File(
                        label="HTML 檔案",
                        visible=False,
                        interactive=False,
                    )
                    download_word_file = gr.File(
                        label="Word 檔案",
                        visible=False,
                        interactive=False,
                    )

        gr.Examples(
            examples=[
                [
                    "我們公司每週要整理業務拜訪紀錄，從業務填寫的表單中，匯整成主管報告。"
                    "目前流程是：業務填Google表單 → 行政複製到Excel → 主管手動整理成PPT。"
                    "每週大約花3-4小時在這件事上，而且常常漏填或格式不一致。",
                    "AI資料整理與報告自動化",
                ],
                [
                    "我們公司推出了一款新的B2B SaaS產品，主要功能是幫助製造業廠商進行生產排程優化。"
                    "目標客戶是中小型製造業的生產主管與IT部門。"
                    "現有的業務提案都是用Word寫的，轉換率很低，客戶說看不懂。",
                    "AI簡報與提案製作",
                ],
                [
                    "我是行銷人員，每個月需要在Facebook、Instagram、LinkedIn發布30篇以上的貼文。"
                    "現在都是手寫，一篇大約花1-2小時。貼文風格要根據平台調整，但核心訊息要一致。"
                    "我們的產品是工業用清潔設備。",
                    "AI行銷文案與社群內容工廠",
                ],
            ],
            inputs=[materials_input, topic_dropdown],
        )

        # wire upload → auto-fill textarea
        file_upload.upload(
            fn=extract_files_text,
            inputs=[file_upload],
            outputs=[materials_input],
        )

        design_btn.click(
            fn=run_course_design,
            inputs=[materials_input, topic_dropdown, user_prompt1],
            outputs=[design_output],
        )
        clear_btn.click(
            fn=lambda: (None, "", "自動判斷", "",
                        "*送出資料後，Agent 將規劃章節並逐一生成，結果即時顯示於此……*",
                        gr.update(visible=False), gr.update(visible=False)),
            outputs=[file_upload, materials_input, topic_dropdown, user_prompt1,
                     design_output, download_file, download_word_file],
        )
        download_btn.click(
            fn=lambda content: (generate_html_file(content), gr.update(visible=True)),
            inputs=[design_output],
            outputs=[download_file, download_file],
        )
        download_word_btn.click(
            fn=lambda content: (generate_word_file(content), gr.update(visible=True)),
            inputs=[design_output],
            outputs=[download_word_file, download_word_file],
        )

    # ── Tab 2 ──────────────────────────────────────────────────────────────────
    with gr.Tab("📑 客戶提案總監Agent"):
        gr.Markdown(
            "# 客戶提案總監Agent\n"
            "整合型錄圖片、報價表、客戶需求與業務筆記，以 **Plan-Execute** 架構逐章節生成"
            "可直接交付給 B2B 製造業客戶的正式提案計畫書。"
        )

        with gr.Accordion("📂 上傳文件或圖表（選填）", open=False):
            gr.Markdown(
                "支援格式：**PDF、Word (.docx)、PowerPoint (.pptx)、"
                "Excel (.xlsx/.csv)、純文字 (.txt/.md/.json)、"
                "圖表圖片 (.png/.jpg/.jpeg/.gif/.webp/.svg)**\n\n"
                "上傳後將自動擷取內容並填入下方輸入框，可再手動補充。"
            )
            file_upload2 = gr.File(
                label="拖曳或點擊上傳（可選取多個檔案）",
                file_types=[
                    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv",
                    ".txt", ".md", ".json",
                    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
                ],
                file_count="multiple",
            )

        with gr.Row():
            with gr.Column(scale=1):
                proposal_input = gr.Textbox(
                    label="提案素材（可直接貼上或由上傳自動填入）",
                    placeholder=(
                        "請貼上任何與提案相關的資料，例如：\n"
                        "・客戶背景與需求說明\n"
                        "・業務拜訪筆記\n"
                        "・產品功能與規格\n"
                        "・報價項目與金額\n"
                        "・提案初稿大綱"
                    ),
                    lines=12,
                )
                user_prompt2 = gr.Textbox(
                    label="使用者提示（選填）",
                    placeholder="可輸入額外指示，例如：強調PoC試點方案、聚焦品保主管受眾、縮短執行摘要……",
                    lines=3,
                )
                propose_btn = gr.Button("開始生成提案", variant="primary", size="lg")
                propose_clear_btn = gr.Button("清除", size="lg")
            with gr.Column(scale=2):
                proposal_output = gr.Markdown(
                    label="提案計畫書",
                    value="*送出資料後，Agent 將規劃提案章節並逐一生成，結果即時顯示於此……*",
                )
                with gr.Row():
                    proposal_html_btn = gr.Button("📥 下載 HTML", size="lg")
                    proposal_word_btn = gr.Button("📄 下載 Word", size="lg")
                with gr.Row():
                    proposal_download_html = gr.File(
                        label="HTML 檔案",
                        visible=False,
                        interactive=False,
                    )
                    proposal_download_word = gr.File(
                        label="Word 檔案",
                        visible=False,
                        interactive=False,
                    )

        gr.Examples(
            examples=[
                [
                    "客戶：台中某中型汽車零件製造商，約300人。\n"
                    "現況：廠內品檢仍靠人工目視，每班3名品檢員，月均漏檢率約2.3%，"
                    "客訴每月4-6件，退貨損失約15萬元。\n"
                    "需求：希望導入自動光學檢測（AOI）系統，預算約200-300萬，"
                    "希望能在Q3完成首條產線導入。\n"
                    "我方產品：工業級AOI視覺檢測系統，每套98萬，含安裝、教育訓練與一年保固。"
                ],
                [
                    "客戶：新竹科技廠，主要生產PCB電路板，客戶要求導入碳排追蹤系統。\n"
                    "痛點：目前沒有系統化碳排資料，面臨供應鏈稽核壓力，今年底前必須提交ESG報告。\n"
                    "我方方案：碳排管理SaaS平台，月費8萬，含資料收集模組、報告產生器與顧問輔導。\n"
                    "業務筆記：客戶採購主管希望先做3個月PoC，成效好再全廠擴展。"
                ],
            ],
            inputs=[proposal_input],
        )

        file_upload2.upload(
            fn=extract_files_text,
            inputs=[file_upload2],
            outputs=[proposal_input],
        )
        propose_btn.click(
            fn=run_proposal_design,
            inputs=[proposal_input, user_prompt2],
            outputs=[proposal_output],
        )
        propose_clear_btn.click(
            fn=lambda: (None, "", "",
                        "*送出資料後，Agent 將規劃提案章節並逐一生成，結果即時顯示於此……*",
                        gr.update(visible=False), gr.update(visible=False)),
            outputs=[file_upload2, proposal_input, user_prompt2, proposal_output,
                     proposal_download_html, proposal_download_word],
        )
        proposal_html_btn.click(
            fn=lambda content: (generate_html_file(content), gr.update(visible=True)),
            inputs=[proposal_output],
            outputs=[proposal_download_html, proposal_download_html],
        )
        proposal_word_btn.click(
            fn=lambda content: (generate_word_file(content), gr.update(visible=True)),
            inputs=[proposal_output],
            outputs=[proposal_download_word, proposal_download_word],
        )

    # ── Tab 3 ──────────────────────────────────────────────────────────────────
    with gr.Tab("🛍️ 商品圖片電商成長Agent"):
        gr.Markdown(
            "# Image to Commerce Decision Growth Agent\n"
            "上傳商品圖片，以 **Plan-Execute** 架構自動生成品牌策略、電商文案、"
            "社群內容、競品分析與高擴散內容設計。"
        )

        with gr.Row():
            img_upload3 = gr.File(
                label="📷 上傳商品圖片（必填，可選取多張）",
                file_types=[".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"],
                file_count="multiple",
            )

        with gr.Accordion("📝 補充商品資訊（選填）", open=False):
            gr.Markdown(
                "可提供：品牌名稱、商品名稱、商品材質、商品尺寸、售價區間、"
                "目標客群、想主打的賣點、銷售平台、品牌語氣、希望避開的詞語、"
                "競品或參考品牌、促銷活動資訊等。"
            )
            extra_info3 = gr.Textbox(
                label="商品補充資訊",
                placeholder=(
                    "品牌名稱：\n"
                    "商品名稱：\n"
                    "商品材質：\n"
                    "售價區間：\n"
                    "目標客群：\n"
                    "主打賣點：\n"
                    "銷售平台：\n"
                    "品牌語氣：\n"
                    "其他備注："
                ),
                lines=10,
            )

        user_prompt3 = gr.Textbox(
            label="使用者提示（選填）",
            placeholder="可輸入額外指示，例如：只產出 Instagram 和小紅書文案、強調送禮情境、精簡版輸出……",
            lines=3,
        )

        with gr.Row():
            commerce_btn = gr.Button("開始分析商品圖片", variant="primary", size="lg")
            commerce_clear_btn = gr.Button("清除", size="lg")

        commerce_output = gr.Markdown(
            label="商品電商成長分析結果",
            value="*上傳商品圖片後，Agent 將規劃分析章節並逐一生成，結果即時顯示於此……*",
        )

        with gr.Row():
            commerce_html_btn = gr.Button("📥 下載 HTML", size="lg")
            commerce_word_btn = gr.Button("📄 下載 Word", size="lg")
        with gr.Row():
            commerce_download_html = gr.File(
                label="HTML 檔案", visible=False, interactive=False
            )
            commerce_download_word = gr.File(
                label="Word 檔案", visible=False, interactive=False
            )

        commerce_btn.click(
            fn=run_commerce_analysis,
            inputs=[img_upload3, extra_info3, user_prompt3],
            outputs=[commerce_output],
        )
        commerce_clear_btn.click(
            fn=lambda: (None, "", "",
                        "*上傳商品圖片後，Agent 將規劃分析章節並逐一生成，結果即時顯示於此……*",
                        gr.update(visible=False), gr.update(visible=False)),
            outputs=[img_upload3, extra_info3, user_prompt3, commerce_output,
                     commerce_download_html, commerce_download_word],
        )
        commerce_html_btn.click(
            fn=lambda content: (generate_html_file(content), gr.update(visible=True)),
            inputs=[commerce_output],
            outputs=[commerce_download_html, commerce_download_html],
        )
        commerce_word_btn.click(
            fn=lambda content: (generate_word_file(content), gr.update(visible=True)),
            inputs=[commerce_output],
            outputs=[commerce_download_word, commerce_download_word],
        )

    # ── Tab 4 ──────────────────────────────────────────────────────────────────
    with gr.Tab("📊 業務智能分析Agent"):
        gr.Markdown(
            "# 業務智能分析Agent Agent\n"
            "上傳客戶資料、報表、會議紀錄或任何附件，以 **Plan-Execute** 架構進行深度質性與量化分析，"
            "並可生成精緻商業風格的**互動式 HTML 儀表板**（含客戶輪廓、痛點偵測、策略模擬），協助精準提案。"
        )

        with gr.Accordion("📂 上傳附件（可選取多個檔案）", open=True):
            gr.Markdown(
                "支援格式：**PDF、Word (.docx)、PowerPoint (.pptx)、"
                "Excel (.xlsx/.csv)、純文字 (.txt/.md/.json)、"
                "圖表圖片 (.png/.jpg/.jpeg/.gif/.webp/.svg)**"
            )
            file_upload4 = gr.File(
                label="拖曳或點擊上傳（可選取多個檔案）",
                file_types=[
                    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv",
                    ".txt", ".md", ".json",
                    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
                ],
                file_count="multiple",
            )

        with gr.Row():
            with gr.Column(scale=1):
                sales_context = gr.Textbox(
                    label="補充背景資訊（選填）",
                    placeholder=(
                        "可輸入任何補充資訊，例如：\n"
                        "・客戶名稱與產業\n"
                        "・我方產品／服務說明\n"
                        "・目前業務進度\n"
                        "・競爭對手資訊\n"
                        "・希望達成的目標"
                    ),
                    lines=8,
                )
                user_prompt4 = gr.Textbox(
                    label="使用者提示（選填）",
                    placeholder="可輸入額外指示，例如：聚焦製造業客戶、重點分析風險、精簡版輸出……",
                    lines=3,
                )
                sales_btn = gr.Button("🔍 開始深度分析", variant="primary", size="lg")
                sales_clear_btn = gr.Button("清除", size="lg")

            with gr.Column(scale=2):
                sales_output = gr.Markdown(
                    label="業務分析報告",
                    value="*上傳附件或輸入資料後，Agent 將規劃分析章節並逐一生成，結果即時顯示於此……*",
                )
                with gr.Row():
                    sales_word_btn = gr.Button("📄 下載分析報告 Word", size="lg")
                with gr.Row():
                    sales_download_word = gr.File(
                        label="Word 分析報告", visible=False, interactive=False
                    )

                gr.Markdown("---")
                gr.Markdown("### 📄 Word 直轉互動儀表板")
                gr.Markdown(
                    "上傳現有 Word 報告（.docx），直接轉換為互動 HTML 儀表板，"
                    "**無需重新分析**，速度更快。"
                )
                with gr.Row():
                    word_direct_upload = gr.File(
                        label="上傳 Word 檔案 (.docx)",
                        file_types=[".docx"],
                        file_count="single",
                        scale=2,
                    )
                    word_convert_btn = gr.Button(
                        "Word → 🎨 生成互動 HTML 儀表板", variant="secondary", size="lg", scale=1
                    )
                word_convert_status = gr.Markdown(value="", visible=True)
                word_convert_download = gr.File(
                    label="轉換後的互動 HTML 儀表板", visible=False, interactive=False
                )

        sales_btn.click(
            fn=run_sales_analysis,
            inputs=[file_upload4, sales_context, user_prompt4],
            outputs=[sales_output],
        )
        sales_clear_btn.click(
            fn=lambda: (None, "", "",
                        "*上傳附件或輸入資料後，Agent 將規劃分析章節並逐一生成，結果即時顯示於此……*",
                        gr.update(value=None, visible=False),
                        None, gr.update(value=None, visible=False), ""),
            outputs=[file_upload4, sales_context, user_prompt4, sales_output,
                     sales_download_word, word_direct_upload, word_convert_download,
                     word_convert_status],
        )
        sales_word_btn.click(
            fn=lambda content: (generate_word_file(content), gr.update(visible=True)),
            inputs=[sales_output],
            outputs=[sales_download_word, sales_download_word],
        )
        word_convert_btn.click(
            fn=convert_word_to_dashboard_tab4,
            inputs=[word_direct_upload],
            outputs=[word_convert_download, word_convert_status],
        )

    # ── Tab 5 ──────────────────────────────────────────────────────────────────
    with gr.Tab("📈 數據分析Agent"):
        gr.Markdown(
            "# 數據分析Agent\n"
            "上傳客戶資料、報表、會議紀錄或任何附件，以 **Plan-Execute** 架構建立**業務提案戰情室分析報告**，"
            "並可生成精緻攝影商業戰情室風格的**互動式 HTML 儀表板**"
            "（含客戶狀態、成交機會雷達、風險地圖、決策鏈、策略模擬、ROI 情境與追蹤行動），協助業務即時判斷推進方式。"
        )

        with gr.Accordion("📂 上傳附件（可選取多個檔案）", open=True):
            gr.Markdown(
                "支援格式：**PDF、Word (.docx)、PowerPoint (.pptx)、"
                "Excel (.xlsx/.csv)、純文字 (.txt/.md/.json)、"
                "圖表圖片 (.png/.jpg/.jpeg/.gif/.webp/.svg)**"
            )
            file_upload5 = gr.File(
                label="拖曳或點擊上傳（可選取多個檔案）",
                file_types=[
                    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv",
                    ".txt", ".md", ".json",
                    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
                ],
                file_count="multiple",
            )

        with gr.Row():
            with gr.Column(scale=1):
                data_context = gr.Textbox(
                    label="補充背景資訊（選填）",
                    placeholder=(
                        "可輸入任何補充資訊，例如：\n"
                        "・客戶名稱與產業\n"
                        "・我方產品／服務說明\n"
                        "・目前決策階段與進度\n"
                        "・已知的決策鏈成員\n"
                        "・希望達成的目標"
                    ),
                    lines=8,
                )
                user_prompt5 = gr.Textbox(
                    label="使用者提示（選填）",
                    placeholder="可輸入額外指示，例如：聚焦風險地圖、加強 ROI 情境、精簡版輸出……",
                    lines=3,
                )
                data_btn = gr.Button("🔍 開始戰情分析", variant="primary", size="lg")
                data_clear_btn = gr.Button("清除", size="lg")

            with gr.Column(scale=2):
                data_output = gr.Markdown(
                    label="業務提案戰情室分析報告",
                    value="*上傳附件或輸入資料後，Agent 將規劃分析章節並逐一生成，結果即時顯示於此……*",
                )
                with gr.Row():
                    data_word_btn = gr.Button("📄 下載分析報告 Word", size="lg")
                with gr.Row():
                    data_download_word = gr.File(
                        label="Word 分析報告", visible=False, interactive=False
                    )

                gr.Markdown("---")
                gr.Markdown("### 📄 Word 直轉互動儀表板")
                gr.Markdown(
                    "上傳現有 Word 報告（.docx），直接轉換為互動 HTML 戰情室儀表板，"
                    "**無需重新分析**，速度更快。"
                )
                with gr.Row():
                    word_direct_upload5 = gr.File(
                        label="上傳 Word 檔案 (.docx)",
                        file_types=[".docx"],
                        file_count="single",
                        scale=2,
                    )
                    word_convert_btn5 = gr.Button(
                        "Word → 🎨 生成互動 HTML 儀表板", variant="secondary", size="lg", scale=1
                    )
                word_convert_status5 = gr.Markdown(value="", visible=True)
                word_convert_download5 = gr.File(
                    label="轉換後的互動 HTML 戰情室儀表板", visible=False, interactive=False
                )

        data_btn.click(
            fn=run_data_analysis,
            inputs=[file_upload5, data_context, user_prompt5],
            outputs=[data_output],
        )
        data_clear_btn.click(
            fn=lambda: (None, "", "",
                        "*上傳附件或輸入資料後，Agent 將規劃分析章節並逐一生成，結果即時顯示於此……*",
                        gr.update(value=None, visible=False),
                        None, gr.update(value=None, visible=False), ""),
            outputs=[file_upload5, data_context, user_prompt5, data_output,
                     data_download_word, word_direct_upload5, word_convert_download5,
                     word_convert_status5],
        )
        data_word_btn.click(
            fn=lambda content: (generate_word_file(content), gr.update(visible=True)),
            inputs=[data_output],
            outputs=[data_download_word, data_download_word],
        )
        word_convert_btn5.click(
            fn=convert_word_to_dashboard_tab5,
            inputs=[word_direct_upload5],
            outputs=[word_convert_download5, word_convert_status5],
        )

    # ── Tab 6 ──────────────────────────────────────────────────────────────────
    with gr.Tab("📊 數據視覺化大師Agent"):
        gr.Markdown(
            "# 數據視覺化大師Agent\n"
            "上傳 **CSV / Excel** 數據檔案，AI 將扮演大數據視覺化大師，"
            "參考 [Apache ECharts](https://echarts.apache.org/examples/zh/index.html) 風格，"
            "根據資料特性提出高級感圖表建議，並自動產出**多維度組合圖表**，"
            "繁體中文標示、含完整圖例與座標軸，可直接下載 JPG 圖片。"
        )

        dataviz_state = gr.State(None)

        with gr.Row():
            with gr.Column(scale=1):
                dataviz_file = gr.File(
                    label="上傳 CSV 或 Excel 檔案",
                    file_types=[".csv", ".xlsx", ".xls"],
                    file_count="single",
                )
                dataviz_sheet = gr.Dropdown(label="選擇工作表（Excel 多分頁時顯示）", visible=False)
                dataviz_status = gr.Markdown(value="*請上傳資料檔案。*")
                dataviz_preview = gr.Dataframe(label="資料預覽（前 10 列）", interactive=False)
                dataviz_context = gr.Textbox(
                    label="補充背景資訊（選填）",
                    placeholder="例如：這是台北市信用卡消費統計、單位為億元、時間跨度2014-2025……",
                    lines=4,
                )
                dataviz_recommend_btn = gr.Button("💡 產生圖表建議", variant="secondary", size="lg")
                dataviz_clear_btn = gr.Button("清除", size="lg")

            with gr.Column(scale=1):
                dataviz_recommendation = gr.Markdown(
                    value="*上傳資料並點擊「產生圖表建議」後，大數據視覺化大師的建議將顯示於此……*"
                )
                dataviz_extra = gr.Textbox(
                    label="額外指示（選填，將加入生成圖表的提示）",
                    placeholder="例如：想看的欄位、特定期間、想強調的洞察……",
                    lines=3,
                )
                dataviz_generate_btn = gr.Button("🎨 生成多維度組合圖（JPG）", variant="primary", size="lg")
                dataviz_gen_status = gr.Markdown(value="")

        gr.Markdown("### 📥 圖表輸出")
        dataviz_gallery = gr.Gallery(label="圖表預覽", columns=2, height=560, object_fit="contain")
        dataviz_files = gr.File(label="下載 JPG 圖表", file_count="multiple", interactive=False)

        dataviz_file.change(
            fn=dataviz_load_file,
            inputs=[dataviz_file],
            outputs=[dataviz_state, dataviz_sheet, dataviz_preview, dataviz_status],
        )
        dataviz_sheet.change(
            fn=dataviz_switch_sheet,
            inputs=[dataviz_state, dataviz_sheet],
            outputs=[dataviz_state, dataviz_preview, dataviz_status],
        )
        dataviz_recommend_btn.click(
            fn=run_dataviz_recommend,
            inputs=[dataviz_state, dataviz_context],
            outputs=[dataviz_recommendation],
        )
        dataviz_generate_btn.click(
            fn=run_dataviz_generate,
            inputs=[dataviz_state, dataviz_recommendation, dataviz_extra],
            outputs=[dataviz_gen_status, dataviz_gallery, dataviz_files],
        )
        dataviz_clear_btn.click(
            fn=lambda: (
                None, None,
                gr.update(choices=[], value=None, visible=False),
                None,
                "*請上傳資料檔案。*",
                "",
                "*上傳資料並點擊「產生圖表建議」後，大數據視覺化大師的建議將顯示於此……*",
                "",
                "",
                [],
                None,
            ),
            outputs=[dataviz_file, dataviz_state, dataviz_sheet, dataviz_preview,
                     dataviz_status, dataviz_context, dataviz_recommendation,
                     dataviz_extra, dataviz_gen_status, dataviz_gallery, dataviz_files],
        )

    # ── Tab 7 ──────────────────────────────────────────────────────────────────
    with gr.Tab("🖼️ 簡報大綱與製作藍圖Agent"):
        gr.Markdown(
            "# 簡報大綱與製作藍圖Agent\n"
            "以 **Plan-Execute** 架構，把文字、文件、圖片或零散筆記整理成清楚、有邏輯、"
            "可直接製作成投影片的**簡報製作藍圖**——包含每頁內容安排、版面與圖表建議、"
            "講者說明方向，以及 3 套附 HEX 色號的簡報視覺風格方案。"
        )

        with gr.Accordion("📂 上傳文件或圖表（選填）", open=False):
            gr.Markdown(
                "支援格式：**PDF、Word (.docx)、PowerPoint (.pptx)、"
                "Excel (.xlsx/.csv)、純文字 (.txt/.md/.json)、"
                "圖表圖片 (.png/.jpg/.jpeg/.gif/.webp/.svg)**\n\n"
                "上傳後將自動擷取內容並填入下方輸入框，可再手動補充。"
            )
            file_upload7 = gr.File(
                label="拖曳或點擊上傳（可選取多個檔案）",
                file_types=[
                    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv",
                    ".txt", ".md", ".json",
                    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
                ],
                file_count="multiple",
            )

        with gr.Row():
            with gr.Column(scale=1):
                presentation_input = gr.Textbox(
                    label="簡報素材內容（可直接貼上或由上傳自動填入）",
                    placeholder="可以是研究內容、課程主題、活動資訊、商業構想、逐字稿、既有簡報、零散筆記……任何原始素材",
                    lines=12,
                )
                presentation_context = gr.Textbox(
                    label="簡報情境設定（選填）",
                    placeholder=(
                        "可輸入任何已知條件，例如：\n"
                        "・目標受眾（高階主管／學員／投資人……）\n"
                        "・簡報時間或頁數\n"
                        "・簡報用途（教學／提案／研究發表／成果展示……）\n"
                        "・是否已有品牌色、Logo 或既有簡報風格\n"
                        "・是否需要講者逐字稿"
                    ),
                    lines=8,
                )
                user_prompt7 = gr.Textbox(
                    label="使用者提示（選填）",
                    placeholder="可輸入額外指示，例如：只要精簡大綱、加強說服邏輯、沿用既有品牌色……",
                    lines=3,
                )
                presentation_btn = gr.Button("🖼️ 開始規劃簡報藍圖", variant="primary", size="lg")
                presentation_clear_btn = gr.Button("清除", size="lg")

            with gr.Column(scale=2):
                presentation_output = gr.Markdown(
                    label="簡報製作藍圖",
                    value="*送出資料後，Agent 將規劃章節並逐一生成，結果即時顯示於此……*",
                )
                with gr.Row():
                    presentation_download_word = gr.Button("📄 下載 Word", size="lg")
                with gr.Row():
                    presentation_word_file = gr.File(
                        label="Word 檔案", visible=False, interactive=False
                    )

        file_upload7.upload(
            fn=extract_files_text,
            inputs=[file_upload7],
            outputs=[presentation_input],
        )

        presentation_btn.click(
            fn=run_presentation_design,
            inputs=[presentation_input, presentation_context, user_prompt7],
            outputs=[presentation_output],
        )
        presentation_clear_btn.click(
            fn=lambda: (None, "", "", "",
                        "*送出資料後，Agent 將規劃章節並逐一生成，結果即時顯示於此……*",
                        gr.update(value=None, visible=False)),
            outputs=[file_upload7, presentation_input, presentation_context, user_prompt7,
                     presentation_output, presentation_word_file],
        )
        presentation_download_word.click(
            fn=lambda content: (generate_word_file(content), gr.update(visible=True)),
            inputs=[presentation_output],
            outputs=[presentation_word_file, presentation_word_file],
        )

if __name__ == "__main__":
    demo.queue().launch(server_name="0.0.0.0", server_port=7860, mcp_server=False, theme=gr.themes.Soft())
